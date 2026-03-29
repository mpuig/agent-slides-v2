from __future__ import annotations

import importlib.util
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation

from agent_slides.io import read_template_manifest

ROOT = Path(__file__).resolve().parents[1]
SCRIPT_PATH = ROOT / "scripts" / "export_layout_inventory.py"


def _load_script_module():
    spec = importlib.util.spec_from_file_location("export_layout_inventory", SCRIPT_PATH)
    assert spec is not None
    assert spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _write_manifest(path: Path) -> None:
    path.write_text(
        json.dumps(
            {
                "source": "fixtures/demo-template.pptx",
                "source_hash": "abc123",
                "slide_masters": [
                    {
                        "index": 0,
                        "name": "Slide Master 1",
                        "layouts": [
                            {
                                "slug": "title_only",
                                "usable": True,
                                "placeholders": [
                                    {"idx": 0, "bounds": {"x": 72, "y": 48, "w": 560, "h": 80}},
                                    {"idx": 1, "bounds": {"x": 72, "y": 136, "w": 560, "h": 36}},
                                ],
                                "slot_mapping": {"heading": 0, "subheading": 1},
                            },
                            {
                                "slug": "body_text",
                                "usable": True,
                                "placeholders": [
                                    {"idx": 0, "bounds": {"x": 72, "y": 48, "w": 560, "h": 80}},
                                    {"idx": 1, "bounds": {"x": 72, "y": 156, "w": 560, "h": 220}},
                                ],
                                "slot_mapping": {"heading": 0, "body": 1},
                            },
                            {
                                "slug": "visual_story",
                                "usable": True,
                                "placeholders": [
                                    {"idx": 0, "bounds": {"x": 72, "y": 48, "w": 560, "h": 80}},
                                    {"idx": 1, "bounds": {"x": 360, "y": 156, "w": 288, "h": 260}},
                                ],
                                "slot_mapping": {"heading": 0, "image": 1},
                            },
                            {
                                "slug": "mixed_story",
                                "usable": True,
                                "placeholders": [
                                    {"idx": 0, "bounds": {"x": 72, "y": 48, "w": 560, "h": 80}},
                                    {"idx": 1, "bounds": {"x": 72, "y": 156, "w": 252, "h": 260}},
                                    {"idx": 2, "bounds": {"x": 360, "y": 156, "w": 288, "h": 260}},
                                ],
                                "slot_mapping": {"heading": 0, "body": 1, "image": 2},
                            },
                            {
                                "slug": "three_up",
                                "usable": True,
                                "placeholders": [
                                    {"idx": 0, "bounds": {"x": 72, "y": 48, "w": 560, "h": 80}},
                                    {"idx": 1, "bounds": {"x": 72, "y": 156, "w": 160, "h": 220}},
                                    {"idx": 2, "bounds": {"x": 256, "y": 156, "w": 160, "h": 220}},
                                    {"idx": 3, "bounds": {"x": 440, "y": 156, "w": 160, "h": 220}},
                                ],
                                "slot_mapping": {"heading": 0, "col1": 1, "col2": 2, "col3": 3},
                            },
                            {
                                "slug": "blankish",
                                "usable": True,
                                "placeholders": [],
                                "slot_mapping": {},
                            },
                            {
                                "slug": "d_body_text",
                                "usable": True,
                                "placeholders": [
                                    {"idx": 0, "bounds": {"x": 72, "y": 48, "w": 560, "h": 80}},
                                    {"idx": 1, "bounds": {"x": 72, "y": 156, "w": 560, "h": 220}},
                                ],
                                "slot_mapping": {"heading": 0, "body": 1},
                            },
                            {
                                "slug": "dict_bounds_layout",
                                "usable": True,
                                "slot_mapping": {
                                    "heading": {"bounds": {"x": 60, "y": 40, "width": 520, "height": 72}},
                                    "body": {"bounds": {"left": 60, "top": 140, "width": 520, "height": 240}},
                                },
                            },
                        ],
                    }
                ],
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )


def _run_script(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, str(SCRIPT_PATH), *args],
        check=False,
        capture_output=True,
        text=True,
    )


def test_export_layout_inventory_cli_outputs_expected_inventory(tmp_path: Path) -> None:
    manifest_path = tmp_path / "template.manifest.json"
    policy_path = tmp_path / "exclude-policy.json"
    _write_manifest(manifest_path)
    policy_path.write_text(
        json.dumps({"layouts": {"three_up": "manual review only"}}) + "\n",
        encoding="utf-8",
    )

    result = _run_script(str(manifest_path), "--exclude-policy", str(policy_path))

    assert result.returncode == 0, result.stderr
    payload = json.loads(result.stdout)
    assert payload["source"] == "fixtures/demo-template.pptx"
    assert payload["layout_count"] == 8
    assert payload["testable_count"] == 5

    layouts = {layout["slug"]: layout for layout in payload["layouts"]}
    assert [layout["slug"] for layout in payload["layouts"]] == sorted(layouts)

    assert layouts["title_only"]["slot_structure"] == "heading_only"
    assert layouts["body_text"]["slot_structure"] == "heading_body"
    assert layouts["visual_story"]["slot_structure"] == "heading_image"
    assert layouts["mixed_story"]["slot_structure"] == "heading_body_image"
    assert layouts["three_up"]["slot_structure"] == "multi_slot"
    assert layouts["blankish"]["slot_structure"] == "blank"

    assert layouts["mixed_story"]["requires_image"] is True
    assert layouts["body_text"]["requires_image"] is False
    assert layouts["d_body_text"]["is_disclaimer_duplicate"] is True
    assert layouts["three_up"]["exclude_reason"] == "manual review only"
    assert layouts["three_up"]["testable"] is False
    assert layouts["blankish"]["testable"] is False
    assert layouts["dict_bounds_layout"]["placeholder_bounds"] == {
        "heading": {"x": 60.0, "y": 40.0, "w": 520.0, "h": 72.0},
        "body": {"x": 60.0, "y": 140.0, "w": 520.0, "h": 240.0},
    }
    assert layouts["mixed_story"]["fillable_slots"] == ["heading", "body", "image"]


def test_export_layout_inventory_output_is_deterministic(tmp_path: Path) -> None:
    manifest_path = tmp_path / "template.manifest.json"
    _write_manifest(manifest_path)

    first = _run_script(str(manifest_path))
    second = _run_script(str(manifest_path))

    assert first.returncode == 0, first.stderr
    assert second.returncode == 0, second.stderr
    assert first.stdout == second.stdout


def test_export_layout_inventory_supports_learned_manifests_from_multiple_templates(tmp_path: Path) -> None:
    module = _load_script_module()

    for index in range(4):
        template_path = tmp_path / f"example-{index + 1}.pptx"
        manifest_path = tmp_path / f"example-{index + 1}.manifest.json"
        Presentation().save(template_path)
        read_template_manifest(template_path, manifest_path)

        payload = module.build_inventory(json.loads(manifest_path.read_text(encoding="utf-8")))

        assert payload["layout_count"] > 0
        assert payload["testable_count"] >= 0
        assert all("slug" in layout for layout in payload["layouts"])
        assert all("slot_structure" in layout for layout in payload["layouts"])


def test_export_layout_inventory_rejects_missing_placeholder_references(tmp_path: Path) -> None:
    manifest_path = tmp_path / "template.manifest.json"
    manifest_path.write_text(
        json.dumps(
            {
                "source": "fixtures/demo-template.pptx",
                "slide_masters": [
                    {
                        "layouts": [
                            {
                                "slug": "broken",
                                "usable": True,
                                "placeholders": [],
                                "slot_mapping": {"heading": 99},
                            }
                        ]
                    }
                ],
            }
        )
        + "\n",
        encoding="utf-8",
    )

    result = _run_script(str(manifest_path))

    assert result.returncode == 1
    assert "references missing placeholder idx 99" in result.stderr
