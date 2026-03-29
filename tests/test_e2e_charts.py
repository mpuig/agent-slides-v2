from __future__ import annotations

import json
from pathlib import Path

from click.testing import CliRunner
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from agent_slides.cli import cli
from agent_slides.errors import CHART_DATA_ERROR


def invoke(
    runner: CliRunner, args: list[str], *, input_text: str | None = None
) -> tuple[int, dict[str, object], str]:
    result = runner.invoke(cli, args, input=input_text)
    payload = json.loads(result.stdout) if result.stdout else {}
    return result.exit_code, payload, result.stderr


def chart_data_payload(payload: dict[str, object]) -> str:
    return json.dumps(payload)


def chart_shapes(path: Path, slide_index: int = 0) -> list[object]:
    presentation = Presentation(str(path))
    return [
        shape for shape in presentation.slides[slide_index].shapes if shape.has_chart
    ]


def chart_shape(path: Path, slide_index: int = 0):
    shapes = chart_shapes(path, slide_index)
    assert len(shapes) == 1
    return shapes[0]


def category_data(chart) -> tuple[tuple[str, ...], list[tuple[str, tuple[float, ...]]]]:
    categories = tuple(chart.plots[0].categories)
    series_data = [(series.name, tuple(series.values)) for series in chart.series]
    return categories, series_data


def scatter_points(series) -> tuple[tuple[float, ...], tuple[float, ...]]:
    x_values = tuple(
        float(value) for value in series._element.xpath(".//c:xVal//c:pt/c:v/text()")
    )
    y_values = tuple(
        float(value) for value in series._element.xpath(".//c:yVal//c:pt/c:v/text()")
    )
    return x_values, y_values


def test_bar_chart_full_flow_builds_editable_native_chart(tmp_path: Path) -> None:
    runner = CliRunner()
    deck_path = tmp_path / "deck.json"
    pptx_path = tmp_path / "bar.pptx"
    chart_data = {
        "categories": ["Q1", "Q2", "Q3"],
        "series": [{"name": "Revenue", "values": [1.5, 2.25, 3.75]}],
    }

    exit_code, payload, _ = invoke(runner, ["init", str(deck_path)])
    assert exit_code == 0
    assert payload["ok"] is True

    exit_code, payload, _ = invoke(
        runner, ["slide", "add", str(deck_path), "--layout", "two_col"]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    exit_code, payload, _ = invoke(
        runner,
        [
            "chart",
            "add",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "left",
            "--type",
            "bar",
            "--data",
            chart_data_payload(chart_data),
            "--title",
            "Revenue",
        ],
    )
    assert exit_code == 0
    assert payload["data"]["chart_type"] == "bar"

    exit_code, payload, _ = invoke(
        runner, ["build", str(deck_path), "-o", str(pptx_path)]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    shape = chart_shape(pptx_path)
    assert shape.has_chart is True
    chart = shape.chart
    assert chart.chart_type == XL_CHART_TYPE.BAR_CLUSTERED
    assert category_data(chart) == (
        ("Q1", "Q2", "Q3"),
        [("Revenue", (1.5, 2.25, 3.75))],
    )


def test_multi_series_line_chart_round_trips_both_series(tmp_path: Path) -> None:
    runner = CliRunner()
    deck_path = tmp_path / "deck.json"
    pptx_path = tmp_path / "line.pptx"
    chart_data = {
        "categories": ["Jan", "Feb", "Mar"],
        "series": [
            {"name": "North", "values": [4.0, 5.5, 6.0]},
            {"name": "South", "values": [3.0, 4.25, 5.75]},
        ],
    }

    invoke(runner, ["init", str(deck_path)])
    invoke(runner, ["slide", "add", str(deck_path), "--layout", "two_col"])

    exit_code, payload, _ = invoke(
        runner,
        [
            "chart",
            "add",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "col1",
            "--type",
            "line",
            "--data",
            chart_data_payload(chart_data),
        ],
    )
    assert exit_code == 0
    assert payload["data"]["chart_type"] == "line"

    exit_code, payload, _ = invoke(
        runner, ["build", str(deck_path), "-o", str(pptx_path)]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    chart = chart_shape(pptx_path).chart
    assert chart.chart_type == XL_CHART_TYPE.LINE
    assert category_data(chart) == (
        ("Jan", "Feb", "Mar"),
        [
            ("North", (4.0, 5.5, 6.0)),
            ("South", (3.0, 4.25, 5.75)),
        ],
    )


def test_scatter_chart_round_trips_xy_data(tmp_path: Path) -> None:
    runner = CliRunner()
    deck_path = tmp_path / "deck.json"
    pptx_path = tmp_path / "scatter.pptx"
    chart_data = {
        "scatter_series": [
            {
                "name": "Observations",
                "points": [
                    {"x": 1.0, "y": 2.5},
                    {"x": 2.5, "y": 3.75},
                    {"x": 4.0, "y": 5.25},
                ],
            }
        ],
    }

    invoke(runner, ["init", str(deck_path)])
    invoke(runner, ["slide", "add", str(deck_path), "--layout", "two_col"])

    exit_code, payload, _ = invoke(
        runner,
        [
            "chart",
            "add",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "col1",
            "--type",
            "scatter",
            "--data",
            chart_data_payload(chart_data),
        ],
    )
    assert exit_code == 0
    assert payload["data"]["chart_type"] == "scatter"

    exit_code, payload, _ = invoke(
        runner, ["build", str(deck_path), "-o", str(pptx_path)]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    chart = chart_shape(pptx_path).chart
    assert chart.chart_type == XL_CHART_TYPE.XY_SCATTER
    series = chart.series[0]
    assert series.name == "Observations"
    assert scatter_points(series) == ((1.0, 2.5, 4.0), (2.5, 3.75, 5.25))


def test_chart_update_builds_with_latest_data(tmp_path: Path) -> None:
    runner = CliRunner()
    deck_path = tmp_path / "deck.json"
    pptx_path = tmp_path / "updated.pptx"
    initial_data = {
        "categories": ["Old A", "Old B"],
        "series": [{"name": "Pipeline", "values": [1.0, 2.0]}],
    }
    updated_data = {
        "categories": ["New A", "New B", "New C"],
        "series": [{"name": "Pipeline", "values": [3.0, 4.5, 6.0]}],
    }

    invoke(runner, ["init", str(deck_path)])
    invoke(runner, ["slide", "add", str(deck_path), "--layout", "two_col"])

    exit_code, payload, _ = invoke(
        runner,
        [
            "chart",
            "add",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "col1",
            "--type",
            "column",
            "--data",
            chart_data_payload(initial_data),
            "--title",
            "Pipeline",
        ],
    )
    assert exit_code == 0
    node_id = str(payload["data"]["node_id"])

    exit_code, payload, _ = invoke(
        runner,
        [
            "chart",
            "update",
            str(deck_path),
            "--node",
            node_id,
            "--data",
            chart_data_payload(updated_data),
        ],
    )
    assert exit_code == 0
    assert payload["data"]["node_id"] == node_id

    exit_code, payload, _ = invoke(
        runner, ["build", str(deck_path), "-o", str(pptx_path)]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    chart = chart_shape(pptx_path).chart
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
    assert category_data(chart) == (
        ("New A", "New B", "New C"),
        [("Pipeline", (3.0, 4.5, 6.0))],
    )


def test_batch_builds_mixed_text_and_chart_slide(tmp_path: Path) -> None:
    runner = CliRunner()
    deck_path = tmp_path / "deck.json"
    pptx_path = tmp_path / "batch.pptx"
    batch_input = json.dumps(
        [
            {"command": "slide_add", "args": {"layout": "two_col"}},
            {
                "command": "slot_set",
                "args": {"slide": 0, "slot": "title", "text": "Quarterly results"},
            },
            {
                "command": "chart_add",
                "args": {
                    "slide": 0,
                    "slot": "left",
                    "type": "bar",
                    "data": {
                        "categories": ["Q1", "Q2"],
                        "series": [{"name": "Revenue", "values": [9.0, 11.0]}],
                    },
                },
            },
            {
                "command": "slot_set",
                "args": {"slide": 0, "slot": "right", "text": "Margin improved."},
            },
        ]
    )

    exit_code, payload, _ = invoke(runner, ["init", str(deck_path)])
    assert exit_code == 0
    assert payload["ok"] is True

    exit_code, payload, _ = invoke(
        runner, ["batch", str(deck_path)], input_text=batch_input
    )
    assert exit_code == 0
    assert payload["data"]["operations"] == 4

    exit_code, payload, _ = invoke(
        runner, ["build", str(deck_path), "-o", str(pptx_path)]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    presentation = Presentation(str(pptx_path))
    slide = presentation.slides[0]
    chart_nodes = [shape for shape in slide.shapes if shape.has_chart]
    text_values = [
        shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
    ]

    assert len(chart_nodes) == 1
    assert "Quarterly results" in text_values
    assert "Margin improved." in text_values
    assert category_data(chart_nodes[0].chart) == (
        ("Q1", "Q2"),
        [("Revenue", (9.0, 11.0))],
    )


def test_pie_chart_rejects_multiple_series(tmp_path: Path) -> None:
    runner = CliRunner()
    deck_path = tmp_path / "deck.json"
    invalid_data = {
        "categories": ["East", "West"],
        "series": [
            {"name": "Revenue", "values": [5.0, 7.0]},
            {"name": "Cost", "values": [2.0, 3.0]},
        ],
    }

    invoke(runner, ["init", str(deck_path)])
    invoke(runner, ["slide", "add", str(deck_path), "--layout", "two_col"])

    result = runner.invoke(
        cli,
        [
            "chart",
            "add",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "col1",
            "--type",
            "pie",
            "--data",
            chart_data_payload(invalid_data),
        ],
    )
    payload = json.loads(result.stderr)

    assert result.exit_code == 1
    assert payload["ok"] is False
    assert payload["error"]["code"] == CHART_DATA_ERROR
    assert (
        payload["error"]["message"]
        == "Invalid chart data: pie charts support exactly one series"
    )
