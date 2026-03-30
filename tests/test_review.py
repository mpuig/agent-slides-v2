from __future__ import annotations

import json
from pathlib import Path

from click.testing import CliRunner
from pptx import Presentation

from agent_slides.cli import cli
from agent_slides.io import write_computed_deck
from agent_slides.model import (
    ChartSpec,
    Counters,
    Deck,
    Node,
    NodeContent,
    Slide,
    TextBlock,
)
from tests.image_helpers import write_png
from tests.test_e2e_template import create_test_template, parse_last_json_line


def write_deck(path: Path, deck: Deck) -> None:
    payload = json.loads(deck.model_dump_json(by_alias=True))
    for slide in payload["slides"]:
        slide.pop("computed", None)
    path.write_text(f"{json.dumps(payload, indent=2)}\n", encoding="utf-8")
    write_computed_deck(str(path), deck)


def fake_render_factory() -> object:
    def fake_render_pptx_to_pngs(
        pptx_path: Path, output_dir: Path, *, dpi: int
    ) -> tuple[Path, list[Path]]:
        assert dpi == 200
        pdf_dir = output_dir / "pdf"
        pdf_dir.mkdir(parents=True, exist_ok=True)
        pdf_path = pdf_dir / f"{pptx_path.stem}.pdf"
        pdf_path.write_bytes(b"%PDF-1.4")

        slides_dir = output_dir / "slides"
        slides_dir.mkdir(parents=True, exist_ok=True)
        presentation = Presentation(str(pptx_path))
        screenshots: list[Path] = []
        for index, _ in enumerate(presentation.slides, start=1):
            screenshot = slides_dir / f"slide-{index:02d}.png"
            write_png(screenshot, width=8, height=8, color=(32 * index, 80, 120))
            screenshots.append(screenshot)
        return pdf_path, screenshots

    return fake_render_pptx_to_pngs


def test_review_command_generates_visual_report(monkeypatch, tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    output_dir = tmp_path / "review-artifacts"
    deck = Deck(
        deck_id="deck-review",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="heading",
                        type="text",
                        content="Q2 revenue review",
                    ),
                    Node(
                        node_id="n-2",
                        slot_binding="subheading",
                        type="text",
                        content="March 2026",
                    ),
                ],
            ),
            Slide(
                slide_id="s-2",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-3",
                        slot_binding="heading",
                        type="text",
                        content="Market Overview",
                    ),
                    Node(
                        node_id="n-4",
                        slot_binding="body",
                        type="text",
                        content=NodeContent(
                            blocks=[
                                TextBlock(
                                    type="bullet",
                                    text="Automation cuts weekly reporting time by 60%",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="Sales coverage expanded into three new segments",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="NPS improved after the onboarding refresh",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="Revenue pipeline quality improved in enterprise",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="Channel mix shifted toward higher-retention accounts",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="Forecast variance narrowed materially",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="Support backlog dropped below the target band",
                                ),
                            ]
                        ),
                    ),
                ],
            ),
            Slide(
                slide_id="s-3",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-5",
                        slot_binding="heading",
                        type="text",
                        content="Revenue increased across regions",
                    ),
                    Node(
                        node_id="n-6",
                        slot_binding="body",
                        type="chart",
                        chart_spec=ChartSpec.model_validate(
                            {
                                "chart_type": "bar",
                                "categories": ["Q1", "Q2"],
                                "series": [{"name": "Revenue", "values": [12.0, 18.0]}],
                            }
                        ),
                    ),
                ],
            ),
        ],
        counters=Counters(slides=3, nodes=6),
    )
    write_deck(deck_path, deck)
    monkeypatch.setattr(
        "agent_slides.review.render_pptx_to_pngs", fake_render_factory()
    )

    runner = CliRunner()
    result = runner.invoke(
        cli, ["review", str(deck_path), "--output-dir", str(output_dir)]
    )

    assert result.exit_code == 0
    payload = json.loads(result.output)
    assert payload["data"]["slides"] == 3
    assert payload["data"]["fixes_applied"] == 0
    assert payload["data"]["overall_grade"]

    report_path = Path(payload["data"]["report_path"])
    report_json_path = Path(payload["data"]["report_json_path"])
    signals_json_path = output_dir / "signals.json"
    assert report_path.exists()
    assert report_json_path.exists()
    assert signals_json_path.exists()

    report = json.loads(report_json_path.read_text(encoding="utf-8"))
    signals = json.loads(signals_json_path.read_text(encoding="utf-8"))
    active = report["active"]
    assert active["categories"]["Content Quality"]["total"] == 8
    issue_items = [issue["item"] for issue in active["all_issues"]]
    assert "Title is not a topic label" in issue_items
    assert "Charts have titles and clear labels" in issue_items
    assert "Slide stays within 6 bullets" in issue_items
    assert len(signals) == 3
    assert signals[0]["layout_slug"] == "title"
    assert "slide-02.png" in report_path.read_text(encoding="utf-8")


def test_review_fix_rewrites_titles_and_chart_metadata(
    monkeypatch, tmp_path: Path
) -> None:
    deck_path = tmp_path / "deck.json"
    output_dir = tmp_path / "review-artifacts"
    deck = Deck(
        deck_id="deck-fix",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="heading",
                        type="text",
                        content="Market Overview",
                    ),
                    Node(
                        node_id="n-2",
                        slot_binding="body",
                        type="text",
                        content=NodeContent(
                            blocks=[
                                TextBlock(
                                    type="bullet",
                                    text="Automation cuts weekly reporting time by 60%",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="Sales coverage expanded into three new segments",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="NPS improved after the onboarding refresh",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="Revenue pipeline quality improved in enterprise",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="Channel mix shifted toward higher-retention accounts",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="Forecast variance narrowed materially",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="Support backlog dropped below the target band",
                                ),
                                TextBlock(
                                    type="bullet",
                                    text="Renewal risk fell in the SMB cohort",
                                ),
                            ]
                        ),
                    ),
                ],
            ),
            Slide(
                slide_id="s-2",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-3",
                        slot_binding="heading",
                        type="text",
                        content="Regional performance",
                    ),
                    Node(
                        node_id="n-4",
                        slot_binding="body",
                        type="chart",
                        chart_spec=ChartSpec.model_validate(
                            {
                                "chart_type": "bar",
                                "categories": ["Q1", "Q2"],
                                "series": [{"name": "Revenue", "values": [8.0, 12.0]}],
                            }
                        ),
                    ),
                ],
            ),
        ],
        counters=Counters(slides=2, nodes=4),
    )
    write_deck(deck_path, deck)
    monkeypatch.setattr(
        "agent_slides.review.render_pptx_to_pngs", fake_render_factory()
    )

    runner = CliRunner()
    result = runner.invoke(
        cli, ["review", str(deck_path), "--output-dir", str(output_dir), "--fix"]
    )

    assert result.exit_code == 0
    payload = json.loads(result.output)
    assert payload["data"]["fixes_applied"] >= 2

    updated = Deck.model_validate_json(deck_path.read_text(encoding="utf-8"))
    assert len(updated.slides) == 3
    assert (
        updated.slides[0].nodes[0].content.to_plain_text()
        == "Automation cuts weekly reporting time by 60%"
    )
    assert updated.slides[1].nodes[0].content.to_plain_text().endswith("(continued)")
    assert updated.slides[2].nodes[1].chart_spec is not None
    assert (
        updated.slides[2].nodes[0].content.to_plain_text()
        == "Revenue increased from Q1 to Q2"
    )
    assert (
        updated.slides[2].nodes[1].chart_spec.title == "Revenue increased from Q1 to Q2"
    )

    report = json.loads(
        Path(payload["data"]["report_json_path"]).read_text(encoding="utf-8")
    )
    assert report["comparison"]["before_grade"]
    assert report["comparison"]["after_grade"]
    assert report["fixes_applied"]


def test_review_command_fails_cleanly_when_tooling_is_missing(
    monkeypatch, tmp_path: Path
) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(
        deck_path,
        Deck(
            deck_id="deck-missing-tool",
            slides=[
                Slide(
                    slide_id="s-1",
                    layout="title",
                    nodes=[
                        Node(
                            node_id="n-1",
                            slot_binding="heading",
                            type="text",
                            content="Hello",
                        )
                    ],
                )
            ],
            counters=Counters(slides=1, nodes=1),
        ),
    )

    def fake_which(name: str) -> str | None:
        return None if name == "soffice" else f"/usr/bin/{name}"

    monkeypatch.setattr("agent_slides.review.shutil.which", fake_which)

    runner = CliRunner()
    result = runner.invoke(cli, ["review", str(deck_path)])

    assert result.exit_code == 1
    error = json.loads(result.stderr)
    assert error["error"]["code"] == "SCHEMA_ERROR"
    assert "soffice" in error["error"]["message"]


def test_review_command_supports_template_decks(monkeypatch, tmp_path: Path) -> None:
    template_path = tmp_path / "brand-template.pptx"
    deck_path = tmp_path / "deck.json"
    output_dir = tmp_path / "review-artifacts"
    create_test_template(template_path)

    runner = CliRunner()
    learn_result = runner.invoke(cli, ["learn", str(template_path)])
    assert learn_result.exit_code == 0
    manifest_path = tmp_path / "brand-template.manifest.json"
    init_result = runner.invoke(
        cli, ["init", str(deck_path), "--template", str(manifest_path)]
    )
    assert init_result.exit_code == 0
    runner.invoke(
        cli,
        ["slide", "add", str(deck_path), "--layout", "title_slide"],
    )
    runner.invoke(
        cli,
        [
            "slot",
            "set",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "heading",
            "--text",
            "Template deck review",
        ],
    )
    monkeypatch.setattr(
        "agent_slides.review.render_pptx_to_pngs", fake_render_factory()
    )

    result = runner.invoke(
        cli, ["review", str(deck_path), "--output-dir", str(output_dir)]
    )

    assert result.exit_code == 0
    payload = parse_last_json_line(result.output)
    report = json.loads(
        Path(payload["data"]["report_json_path"]).read_text(encoding="utf-8")
    )
    signals = json.loads((output_dir / "signals.json").read_text(encoding="utf-8"))
    assert report["active"]["deck"]["template"] == "template"
    assert len(signals) == 1
    assert signals[0]["layout_slug"] == "title_slide"
