from __future__ import annotations

import hashlib
import json
from pathlib import Path
from xml.etree import ElementTree as ET
from zipfile import ZipFile

import pytest
from click.testing import CliRunner, Result
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from agent_slides.cli import cli
from agent_slides.errors import AgentSlidesError, FILE_NOT_FOUND, SCHEMA_ERROR
from agent_slides.io import read_template_manifest
from agent_slides.io import template_reader

PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"p": PML_NS, "a": DML_NS}

ET.register_namespace("a", DML_NS)
ET.register_namespace("p", PML_NS)
ET.register_namespace(
    "r", "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
)


def invoke_cli(args: list[str]) -> Result:
    runner = CliRunner()
    return runner.invoke(cli, args)


def parse_last_json_line(output: str) -> dict[str, object]:
    return json.loads(output.strip().splitlines()[-1])


def patch_template(path: Path) -> None:
    updated_path = path.with_suffix(".tmp")

    with ZipFile(path) as source_zip, ZipFile(updated_path, "w") as dest_zip:
        for info in source_zip.infolist():
            data = source_zip.read(info.filename)

            if info.filename == "ppt/slideLayouts/slideLayout2.xml":
                data = _rewrite_xml(
                    data,
                    lambda root: (
                        _set_layout_name(root, "Agenda"),
                        _set_placeholder_type(root, 1, "body"),
                    ),
                )
            elif info.filename == "ppt/slideLayouts/slideLayout4.xml":
                data = _rewrite_xml(
                    data,
                    lambda root: (
                        _set_layout_name(root, "Agenda"),
                        _set_placeholder_type(root, 1, "body"),
                        _set_placeholder_type(root, 2, "body"),
                    ),
                )
            elif info.filename == "ppt/slideLayouts/slideLayout5.xml":
                data = _rewrite_xml(
                    data,
                    lambda root: (
                        _set_layout_name(root, "Comparison Lab"),
                        _set_placeholder_type(root, 2, "chart"),
                        _set_placeholder_type(root, 4, "media"),
                    ),
                )
            elif info.filename == "ppt/slideLayouts/slideLayout8.xml":
                data = _rewrite_xml(
                    data,
                    lambda root: (
                        _set_layout_name(root, "Captioned Content"),
                        _set_placeholder_type(root, 1, "tbl"),
                    ),
                )
            elif info.filename == "ppt/slideLayouts/slideLayout9.xml":
                data = _rewrite_xml(
                    data,
                    lambda root: _set_layout_name(root, "Photo Story"),
                )
            elif info.filename == "ppt/theme/theme1.xml":
                data = _rewrite_xml(data, _set_theme_values)

            dest_zip.writestr(info, data)

    updated_path.replace(path)


def _rewrite_xml(xml_bytes: bytes, update: object) -> bytes:
    root = ET.fromstring(xml_bytes)
    update(root)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _set_layout_name(root: ET.Element, name: str) -> None:
    common_slide = root.find("./p:cSld", NS)
    assert common_slide is not None
    common_slide.set("name", name)


def _set_placeholder_type(root: ET.Element, idx: int, placeholder_type: str) -> None:
    placeholder = root.find(f".//p:ph[@idx='{idx}']", NS)
    assert placeholder is not None
    placeholder.set("type", placeholder_type)


def _set_theme_values(root: ET.Element) -> None:
    for color_name, value in {
        "dk1": "101010",
        "lt1": "FAFAFA",
        "dk2": "202020",
        "lt2": "E5E7EB",
        "accent1": "112233",
        "accent2": "445566",
        "accent3": "778899",
    }.items():
        _set_scheme_color(root, color_name, value)

    _set_latin_font(root, "majorFont", "Aptos Display")
    _set_latin_font(root, "minorFont", "Aptos")


def _set_scheme_color(root: ET.Element, color_name: str, value: str) -> None:
    color = root.find(f".//a:clrScheme/a:{color_name}", NS)
    assert color is not None

    srgb = color.find("./a:srgbClr", NS)
    if srgb is not None:
        srgb.set("val", value)
        return

    system = color.find("./a:sysClr", NS)
    assert system is not None
    system.set("lastClr", value)


def _set_latin_font(root: ET.Element, scheme_name: str, typeface: str) -> None:
    font = root.find(f".//a:fontScheme/a:{scheme_name}/a:latin", NS)
    assert font is not None
    font.set("typeface", typeface)


def make_template_fixture(tmp_path: Path) -> Path:
    template_path = tmp_path / "corporate-template.pptx"
    Presentation().save(template_path)
    patch_template(template_path)
    return template_path


def manifest_layouts(manifest: dict[str, object]) -> dict[str, dict[str, object]]:
    layouts: dict[str, dict[str, object]] = {}
    for master in manifest["slide_masters"]:
        for layout in master["layouts"]:
            layouts[layout["slug"]] = layout
    return layouts


def test_learn_command_writes_manifest_and_extracts_expected_structure(
    tmp_path: Path,
) -> None:
    template_path = make_template_fixture(tmp_path)

    result = invoke_cli(["learn", str(template_path)])

    assert result.exit_code == 0
    payload = parse_last_json_line(result.output)
    manifest_path = tmp_path / "corporate-template.manifest.json"
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    layouts = manifest_layouts(manifest)

    assert payload == {
        "ok": True,
        "data": {
            "source": "corporate-template.pptx",
            "layouts_found": 11,
            "usable_layouts": 11,
        },
    }
    assert manifest["source"] == "corporate-template.pptx"
    assert (
        manifest["source_hash"]
        == hashlib.sha256(template_path.read_bytes()).hexdigest()
    )
    assert manifest["theme"] == {
        "colors": {
            "primary": "#112233",
            "secondary": "#445566",
            "accent": "#778899",
            "text": "#101010",
            "heading_text": "#202020",
            "subtle_text": "#E5E7EB",
            "background": "#FAFAFA",
        },
        "fonts": {
            "heading": "Aptos Display",
            "body": "Aptos",
        },
        "spacing": {
            "base_unit": 10.0,
            "margin": 60.0,
            "gutter": 20.0,
        },
    }

    assert layouts["title_slide"]["slot_mapping"] == {"heading": 0, "subheading": 1}
    assert layouts["agenda"]["slot_mapping"] == {"heading": 0, "body": 1}
    assert layouts["agenda_2"]["slot_mapping"] == {"heading": 0, "col1": 1, "col2": 2}
    assert layouts["comparison_lab"]["slot_mapping"] == {
        "heading": 0,
        "col1": 1,
        "col2": 3,
    }
    assert layouts["photo_story"]["slot_mapping"] == {
        "heading": 0,
        "body": 2,
        "image": 1,
    }
    assert layouts["blank"]["usable"] is True
    assert layouts["blank"]["placeholders"] == []
    assert layouts["blank"]["color_zones"] == [
        {
            "region": "full_slide",
            "left": 0.0,
            "width": 960.0,
            "bg_color": "FAFAFA",
            "text_color": "333333",
        }
    ]
    assert layouts["blank"]["editable_regions"] == [
        {
            "name": "content_area",
            "left": 0.0,
            "top": 0.0,
            "width": 960.0,
            "height": 540.0,
            "source": "visual_inference_no_placeholders",
        }
    ]
    assert layouts["agenda"]["placeholders"][1]["idx"] == 1
    assert layouts["agenda"]["placeholders"][1]["type"] == "BODY"
    assert layouts["agenda"]["placeholders"][1]["name"] == "Content Placeholder 2"
    assert layouts["agenda"]["placeholders"][1]["bounds"] == {
        "x": 36.0,
        "y": 126.0,
        "w": 648.0,
        "h": 356.375,
    }
    assert layouts["agenda"]["placeholders"][1]["shape_kind"] == "placeholder"
    assert layouts["agenda"]["placeholders"][1]["suggested_slot"] == "body"
    assert layouts["agenda"]["editable_regions"] == [
        {
            "name": "content_box",
            "left": 36.0,
            "top": 126.0,
            "width": 648.0,
            "height": 356.375,
            "source": "placeholder_union",
        }
    ]
    agenda_2_body_placeholders = [
        placeholder
        for placeholder in layouts["agenda_2"]["placeholders"]
        if placeholder["type"] == "BODY"
    ]
    assert layouts["agenda_2"]["editable_regions"] == [
        {
            "name": "content_box",
            "left": min(
                placeholder["bounds"]["x"] for placeholder in agenda_2_body_placeholders
            ),
            "top": min(
                placeholder["bounds"]["y"] for placeholder in agenda_2_body_placeholders
            ),
            "width": max(
                placeholder["bounds"]["x"] + placeholder["bounds"]["w"]
                for placeholder in agenda_2_body_placeholders
            )
            - min(
                placeholder["bounds"]["x"] for placeholder in agenda_2_body_placeholders
            ),
            "height": max(
                placeholder["bounds"]["y"] + placeholder["bounds"]["h"]
                for placeholder in agenda_2_body_placeholders
            )
            - min(
                placeholder["bounds"]["y"] for placeholder in agenda_2_body_placeholders
            ),
            "source": "placeholder_union",
        }
    ]

    assert (
        "Warning: layout 'Comparison Lab': skipped unsupported media_clip placeholder 'Content Placeholder 5'"
        in result.output
    )
    assert (
        "Warning: layout 'Comparison Lab': skipped unsupported chart placeholder 'Content Placeholder 3'"
        not in result.output
    )
    assert (
        "Warning: layout 'Captioned Content': skipped unsupported table placeholder 'Content Placeholder 2'"
        not in result.output
    )


def test_learn_uses_relative_source_when_output_path_changes(tmp_path: Path) -> None:
    template_path = make_template_fixture(tmp_path)
    output_path = tmp_path / "nested" / "template.manifest.json"
    output_path.parent.mkdir()

    result = read_template_manifest(template_path, output_path)

    manifest = json.loads(output_path.read_text(encoding="utf-8"))
    assert result.source == "../corporate-template.pptx"
    assert manifest["source"] == "../corporate-template.pptx"


def test_learn_command_returns_file_not_found_error(tmp_path: Path) -> None:
    result = invoke_cli(["learn", str(tmp_path / "missing-template.pptx")])

    assert result.exit_code == 1
    assert parse_last_json_line(result.output) == {
        "ok": False,
        "error": {
            "code": FILE_NOT_FOUND,
            "message": f"Template file not found: {tmp_path / 'missing-template.pptx'}",
        },
    }


def test_learn_command_returns_schema_error_for_invalid_pptx(tmp_path: Path) -> None:
    template_path = tmp_path / "invalid.pptx"
    template_path.write_text("not a pptx", encoding="utf-8")

    result = invoke_cli(["learn", str(template_path)])

    assert result.exit_code == 1
    assert parse_last_json_line(result.output) == {
        "ok": False,
        "error": {
            "code": SCHEMA_ERROR,
            "message": "not a valid PPTX file",
        },
    }


def test_learn_command_returns_schema_error_for_password_protected_file(
    tmp_path: Path,
) -> None:
    template_path = tmp_path / "protected.pptx"
    template_path.write_bytes(
        bytes.fromhex("D0CF11E0A1B11AE1") + b"encrypted-office-data"
    )

    result = invoke_cli(["learn", str(template_path)])

    assert result.exit_code == 1
    assert parse_last_json_line(result.output) == {
        "ok": False,
        "error": {
            "code": SCHEMA_ERROR,
            "message": "password-protected files not supported",
        },
    }


def test_read_template_manifest_raises_schema_error_for_zero_layouts(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    template_path = tmp_path / "template.pptx"
    template_path.write_bytes(b"placeholder")

    class FakeSlideMaster:
        slide_layouts: list[object] = []

    class FakePresentation:
        slide_masters = [FakeSlideMaster()]

    monkeypatch.setattr(
        template_reader, "_open_presentation", lambda _: FakePresentation()
    )

    with pytest.raises(
        AgentSlidesError, match="template has no slide layouts"
    ) as exc_info:
        read_template_manifest(template_path)

    assert exc_info.value.code == SCHEMA_ERROR


def test_read_template_manifest_warns_when_no_layouts_are_usable(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    template_path = tmp_path / "template.pptx"
    template_path.write_bytes(b"placeholder")

    class FakePlaceholderFormat:
        idx = 10
        type = object()

    class FakePlaceholder:
        name = "Footer Placeholder 1"
        left = 0
        top = 0
        width = 100
        height = 100
        placeholder_format = FakePlaceholderFormat()

    class FakeLayout:
        name = "Footer Only"
        placeholders = [FakePlaceholder()]

    class FakeSlideMaster:
        slide_layouts = [FakeLayout()]

    class FakePresentation:
        slide_masters = [FakeSlideMaster()]

    monkeypatch.setattr(
        template_reader, "_open_presentation", lambda _: FakePresentation()
    )
    monkeypatch.setattr(
        template_reader,
        "_extract_theme",
        lambda _: {"colors": {}, "fonts": {}, "spacing": {}},
    )

    result = read_template_manifest(template_path)
    manifest = json.loads(result.manifest_path.read_text(encoding="utf-8"))

    assert result.layouts_found == 1
    assert result.usable_layouts == 0
    assert result.warnings == ["template has 0 usable layouts"]
    assert manifest["slide_masters"][0]["layouts"][0]["slot_mapping"] == {}


def test_read_template_manifest_extracts_non_placeholder_content_shapes(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    template_path = tmp_path / "template.pptx"
    template_path.write_bytes(b"placeholder")

    class FakeTextFrame:
        def __init__(self, text: str) -> None:
            self.text = text

    class FakeTextShape:
        def __init__(self, *, shape_id: int, name: str, text: str, top: int) -> None:
            self.shape_id = shape_id
            self.name = name
            self.left = 0
            self.top = top
            self.width = 100
            self.height = 40 * 12700  # 40pt in EMU
            self.is_placeholder = False
            self.has_text_frame = True
            self.has_table = False
            self.has_chart = False
            self.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
            self.text_frame = FakeTextFrame(text)

    class FakeTable:
        rows = [object(), object(), object()]
        columns = [object(), object()]

    class FakeTableShape:
        def __init__(self) -> None:
            self.shape_id = 3
            self.name = "Agenda Table"
            self.left = 0
            self.top = 120 * 12700
            self.width = 200 * 12700
            self.height = 100 * 12700
            self.is_placeholder = False
            self.has_text_frame = False
            self.has_table = True
            self.has_chart = False
            self.shape_type = MSO_SHAPE_TYPE.TABLE
            self.table = FakeTable()

    class FakeChart:
        chart_type = 7

    class FakeChartShape:
        def __init__(self) -> None:
            self.shape_id = 4
            self.name = "Score Chart"
            self.left = 0
            self.top = 240 * 12700
            self.width = 200 * 12700
            self.height = 100 * 12700
            self.is_placeholder = False
            self.has_text_frame = False
            self.has_table = False
            self.has_chart = True
            self.shape_type = MSO_SHAPE_TYPE.CHART
            self.chart = FakeChart()

    class FakeGroupShape:
        def __init__(self) -> None:
            self.shape_id = 5
            self.name = "Agenda Group"
            self.left = 0
            self.top = 360 * 12700
            self.width = 200 * 12700
            self.height = 100 * 12700
            self.is_placeholder = False
            self.has_text_frame = False
            self.has_table = False
            self.has_chart = False
            self.shape_type = MSO_SHAPE_TYPE.GROUP
            self.shapes = [object(), object()]

    class FakeLayout:
        name = "Quote Agenda"
        placeholders: list[object] = []
        shapes = [
            FakeTextShape(shape_id=1, name="Quote", text="Stay hungry", top=0),
            FakeTextShape(shape_id=2, name="Attribution", text="Steve Jobs", top=60),
            FakeTableShape(),
            FakeChartShape(),
            FakeGroupShape(),
        ]

    class FakeSlideMaster:
        slide_layouts = [FakeLayout()]

    class FakePresentation:
        slide_masters = [FakeSlideMaster()]

    monkeypatch.setattr(
        template_reader, "_open_presentation", lambda _: FakePresentation()
    )
    monkeypatch.setattr(
        template_reader,
        "_extract_theme",
        lambda _: {"colors": {}, "fonts": {}, "spacing": {}},
    )

    result = read_template_manifest(template_path)
    layout = json.loads(result.manifest_path.read_text(encoding="utf-8"))[
        "slide_masters"
    ][0]["layouts"][0]

    assert result.usable_layouts == 1
    assert layout["usable"] is True
    assert [placeholder["type"] for placeholder in layout["placeholders"]] == [
        "TEXT_BOX",
        "TEXT_BOX",
        "TABLE",
        "CHART",
        "GROUP",
    ]
    assert layout["placeholders"][0]["suggested_slot"] == "quote"
    assert layout["placeholders"][1]["suggested_slot"] == "attribution"
    assert layout["placeholders"][2]["table"] == {"rows": 3, "cols": 2}
    assert layout["placeholders"][3]["chart"] == {"chart_type": 7}
    assert layout["placeholders"][4]["group"] == {"children": 2}
    assert layout["slot_mapping"]["quote"] == 1_000_001
    assert layout["slot_mapping"]["attribution"] == 1_000_002
    # TABLE, CHART, GROUP are at different vertical positions (not same row),
    # so the first body-like shape wins as "body", not col1/col2/col3
    assert layout["slot_mapping"]["body"] == 1_000_003


def test_read_template_manifest_extracts_color_zones_from_large_filled_shapes(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    template_path = tmp_path / "template.pptx"
    template_path.write_bytes(b"placeholder")

    class FakePlaceholderFormat:
        idx = 0
        type = PP_PLACEHOLDER.TITLE

    class FakeTitlePlaceholder:
        name = "Title Placeholder 1"
        left = 60 * 12700
        top = 40 * 12700
        width = 360 * 12700
        height = 50 * 12700
        shape_id = 1
        placeholder_format = FakePlaceholderFormat()

    class FakeForeColor:
        def __init__(self, rgb: str) -> None:
            self.rgb = rgb

    class FakeFill:
        def __init__(self, rgb: str) -> None:
            self.type = MSO_FILL.SOLID
            self.fore_color = FakeForeColor(rgb)

    class FakePanelShape:
        is_placeholder = False
        has_text_frame = False
        has_table = False
        has_chart = False
        shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE

        def __init__(
            self, *, shape_id: int, name: str, left_pt: float, width_pt: float, rgb: str
        ) -> None:
            self.shape_id = shape_id
            self.name = name
            self.left = int(left_pt * 12700)
            self.top = 0
            self.width = int(width_pt * 12700)
            self.height = int(template_reader.SLIDE_HEIGHT_PT * 12700)
            self.fill = FakeFill(rgb)

    class FakeDecorativeShape:
        is_placeholder = False
        has_text_frame = False
        has_table = False
        has_chart = False
        shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
        shape_id = 4
        name = "Small Accent"
        left = 200 * 12700
        top = 120 * 12700
        width = 100 * 12700
        height = 60 * 12700
        fill = FakeFill("FFAA00")

    layout_root = ET.fromstring(
        f"""
        <p:sldLayout xmlns:p="{PML_NS}" xmlns:a="{DML_NS}">
          <p:cSld name="Green Highlight">
            <p:bg>
              <p:bgPr>
                <a:solidFill>
                  <a:srgbClr val="FFFFFF" />
                </a:solidFill>
              </p:bgPr>
            </p:bg>
          </p:cSld>
        </p:sldLayout>
        """
    )
    master_root = ET.fromstring(
        f"""
        <p:sldMaster xmlns:p="{PML_NS}" xmlns:a="{DML_NS}">
          <p:cSld>
            <p:bg>
              <p:bgPr>
                <a:solidFill>
                  <a:srgbClr val="F1F5F9" />
                </a:solidFill>
              </p:bgPr>
            </p:bg>
          </p:cSld>
        </p:sldMaster>
        """
    )

    class FakeSlideMaster:
        def __init__(self) -> None:
            self._element = master_root
            self.slide_layouts = [FakeLayout(self)]

    class FakeLayout:
        def __init__(self, slide_master: FakeSlideMaster) -> None:
            self.name = "Green Highlight"
            self.slide_master = slide_master
            self.placeholders = [FakeTitlePlaceholder()]
            self.shapes = [
                FakePanelShape(
                    shape_id=2,
                    name="Left Panel",
                    left_pt=0.0,
                    width_pt=420.0,
                    rgb="FFFFFF",
                ),
                FakePanelShape(
                    shape_id=3,
                    name="Right Panel",
                    left_pt=540.0,
                    width_pt=420.0,
                    rgb="00A651",
                ),
                FakeDecorativeShape(),
            ]
            self._element = layout_root

    class FakePresentation:
        slide_masters = [FakeSlideMaster()]

    monkeypatch.setattr(
        template_reader, "_open_presentation", lambda _: FakePresentation()
    )
    monkeypatch.setattr(
        template_reader,
        "_extract_theme",
        lambda _: {
            "colors": {
                "primary": "#112233",
                "secondary": "#445566",
                "accent": "#778899",
                "text": "#111111",
                "heading_text": "#222222",
                "subtle_text": "#E5E7EB",
                "background": "#FAFAFA",
            },
            "fonts": {},
            "spacing": {},
        },
    )

    result = read_template_manifest(template_path)
    layout = json.loads(result.manifest_path.read_text(encoding="utf-8"))[
        "slide_masters"
    ][0]["layouts"][0]

    assert result.usable_layouts == 1
    assert layout["slug"] == "green_highlight"
    assert layout["color_zones"] == [
        {
            "region": "panel_0",
            "left": 0.0,
            "width": 420.0,
            "bg_color": "FFFFFF",
            "text_color": "333333",
            "editable_above": {
                "left": 0.0,
                "top": 0.0,
                "width": 420.0,
                "height": 40.0,
            },
            "editable_below": {
                "left": 0.0,
                "top": 90.0,
                "width": 420.0,
                "height": 450.0,
            },
        },
        {
            "region": "gap_0",
            "left": 420.0,
            "width": 120.0,
            "bg_color": "FFFFFF",
            "text_color": "333333",
        },
        {
            "region": "panel_1",
            "left": 540.0,
            "width": 420.0,
            "bg_color": "00A651",
            "text_color": "FFFFFF",
        },
    ]


def test_read_template_manifest_extracts_editable_regions_from_visual_inference(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    template_path = tmp_path / "template.pptx"
    template_path.write_bytes(b"placeholder")

    class FakePlaceholderFormat:
        idx = 0
        type = PP_PLACEHOLDER.TITLE

    class FakeTitlePlaceholder:
        name = "Title Placeholder 1"
        left = 560 * 12700
        top = 40 * 12700
        width = 300 * 12700
        height = 50 * 12700
        shape_id = 1
        placeholder_format = FakePlaceholderFormat()

    class FakeForeColor:
        def __init__(self, rgb: str) -> None:
            self.rgb = rgb

    class FakeFill:
        def __init__(self, rgb: str) -> None:
            self.type = MSO_FILL.SOLID
            self.fore_color = FakeForeColor(rgb)

    class FakePanelShape:
        is_placeholder = False
        has_text_frame = False
        has_table = False
        has_chart = False
        shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE

        def __init__(
            self, *, shape_id: int, name: str, left_pt: float, width_pt: float, rgb: str
        ) -> None:
            self.shape_id = shape_id
            self.name = name
            self.left = int(left_pt * 12700)
            self.top = 0
            self.width = int(width_pt * 12700)
            self.height = int(template_reader.SLIDE_HEIGHT_PT * 12700)
            self.fill = FakeFill(rgb)

    class FakeTextFrame:
        def __init__(self, text: str) -> None:
            self.text = text

    class FakeFooterShape:
        is_placeholder = False
        has_text_frame = True
        has_table = False
        has_chart = False
        shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        shape_id = 4
        name = "Footer Text"
        left = 620 * 12700
        top = 500 * 12700
        width = 220 * 12700
        height = 16 * 12700
        text_frame = FakeTextFrame("Copyright 2026")

    class FakeDecorativeShape:
        is_placeholder = False
        has_text_frame = False
        has_table = False
        has_chart = False
        shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
        shape_id = 5
        name = "Accent Block"
        left = 820 * 12700
        top = 110 * 12700
        width = 140 * 12700
        height = 60 * 12700
        fill = FakeFill("0A7A3B")

    layout_root = ET.fromstring(
        f"""
        <p:sldLayout xmlns:p="{PML_NS}" xmlns:a="{DML_NS}">
          <p:cSld name="Green Highlight">
            <p:bg>
              <p:bgPr>
                <a:solidFill>
                  <a:srgbClr val="FFFFFF" />
                </a:solidFill>
              </p:bgPr>
            </p:bg>
          </p:cSld>
        </p:sldLayout>
        """
    )
    master_root = ET.fromstring(
        f"""
        <p:sldMaster xmlns:p="{PML_NS}" xmlns:a="{DML_NS}">
          <p:cSld>
            <p:bg>
              <p:bgPr>
                <a:solidFill>
                  <a:srgbClr val="F1F5F9" />
                </a:solidFill>
              </p:bgPr>
            </p:bg>
          </p:cSld>
        </p:sldMaster>
        """
    )

    class FakeSlideMaster:
        def __init__(self) -> None:
            self._element = master_root
            self.slide_layouts = [FakeLayout(self)]

    class FakeLayout:
        def __init__(self, slide_master: FakeSlideMaster) -> None:
            self.name = "Green Highlight"
            self.slide_master = slide_master
            self.placeholders = [FakeTitlePlaceholder()]
            self.shapes = [
                FakePanelShape(
                    shape_id=2,
                    name="Left Panel",
                    left_pt=0.0,
                    width_pt=420.0,
                    rgb="FFFFFF",
                ),
                FakePanelShape(
                    shape_id=3,
                    name="Right Panel",
                    left_pt=540.0,
                    width_pt=420.0,
                    rgb="00A651",
                ),
                FakeDecorativeShape(),
                FakeFooterShape(),
            ]
            self._element = layout_root

    class FakePresentation:
        slide_masters = [FakeSlideMaster()]

    monkeypatch.setattr(
        template_reader, "_open_presentation", lambda _: FakePresentation()
    )
    monkeypatch.setattr(
        template_reader,
        "_extract_theme",
        lambda _: {
            "colors": {
                "primary": "#112233",
                "secondary": "#445566",
                "accent": "#778899",
                "text": "#111111",
                "heading_text": "#222222",
                "subtle_text": "#E5E7EB",
                "background": "#FAFAFA",
            },
            "fonts": {},
            "spacing": {},
        },
    )

    result = read_template_manifest(template_path)
    layout = json.loads(result.manifest_path.read_text(encoding="utf-8"))[
        "slide_masters"
    ][0]["layouts"][0]

    assert result.usable_layouts == 1
    assert layout["slug"] == "green_highlight"
    assert layout["editable_regions"] == [
        {
            "name": "content_area",
            "left": 540.0,
            "top": 170.0,
            "width": 420.0,
            "height": 320.0,
            "source": "visual_inference_no_placeholders",
        }
    ]


def test_read_template_manifest_marks_blank_variants_usable(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    template_path = tmp_path / "template.pptx"
    template_path.write_bytes(b"placeholder")

    class FakeLayout:
        name = "d_blank"
        placeholders: list[object] = []
        shapes: list[object] = []

    class FakeSlideMaster:
        slide_layouts = [FakeLayout()]

    class FakePresentation:
        slide_masters = [FakeSlideMaster()]

    monkeypatch.setattr(
        template_reader, "_open_presentation", lambda _: FakePresentation()
    )
    monkeypatch.setattr(
        template_reader,
        "_extract_theme",
        lambda _: {"colors": {}, "fonts": {}, "spacing": {}},
    )

    result = read_template_manifest(template_path)
    layout = json.loads(result.manifest_path.read_text(encoding="utf-8"))[
        "slide_masters"
    ][0]["layouts"][0]

    assert result.usable_layouts == 1
    assert layout["usable"] is True
    assert layout["placeholders"] == []
