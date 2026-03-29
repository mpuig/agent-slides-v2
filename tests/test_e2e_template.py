from __future__ import annotations

import hashlib
import json
from pathlib import Path
from xml.etree import ElementTree as ET
from zipfile import ZipFile

import pytest
from click.testing import CliRunner, Result
from pptx import Presentation
from pptx.util import Pt

from agent_slides.cli import cli
from agent_slides.errors import TEMPLATE_CHANGED

PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"p": PML_NS, "a": DML_NS}

ET.register_namespace("a", DML_NS)
ET.register_namespace("p", PML_NS)
ET.register_namespace("r", "http://schemas.openxmlformats.org/officeDocument/2006/relationships")


def invoke(args: list[str]) -> Result:
    runner = CliRunner()
    return runner.invoke(cli, args)


def parse_last_json_line(output: str) -> dict[str, object]:
    return json.loads(output.strip().splitlines()[-1])


def collect_text_values(path: Path) -> list[str]:
    presentation = Presentation(str(path))
    values: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                values.append(shape.text_frame.text)
    return values


def layout_by_slug(manifest: dict[str, object], slug: str) -> dict[str, object]:
    for master in manifest["slide_masters"]:
        for layout in master["layouts"]:
            if layout["slug"] == slug:
                return layout
    raise AssertionError(f"Layout {slug} not found")


def create_test_template(path: Path, layouts: dict[int, str] | None = None) -> None:
    """Create a PPTX template with predictable theme and layout names for testing."""

    Presentation().save(path)

    updates: dict[str, object] = {
        "ppt/theme/theme1.xml": _set_theme_values,
    }
    for index, name in (layouts or {}).items():
        updates[f"ppt/slideLayouts/slideLayout{index + 1}.xml"] = _layout_name_updater(name)

    _rewrite_pptx(path, updates)


def patch_template_theme(path: Path, *, accent1: str) -> None:
    _rewrite_pptx(
        path,
        {
            "ppt/theme/theme1.xml": lambda root: _set_scheme_color(root, "accent1", accent1),
        },
    )


def _rewrite_pptx(path: Path, updates: dict[str, object]) -> None:
    updated_path = path.with_suffix(".tmp")
    with ZipFile(path) as source_zip, ZipFile(updated_path, "w") as dest_zip:
        for info in source_zip.infolist():
            data = source_zip.read(info.filename)
            updater = updates.get(info.filename)
            if updater is not None:
                data = _rewrite_xml(data, updater)
            dest_zip.writestr(info, data)
    updated_path.replace(path)


def _rewrite_xml(xml_bytes: bytes, update: object) -> bytes:
    root = ET.fromstring(xml_bytes)
    update(root)
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def _layout_name_updater(name: str):
    def update(root: ET.Element) -> None:
        common_slide = root.find("./p:cSld", NS)
        assert common_slide is not None
        common_slide.set("name", name)

    return update


def _set_theme_values(root: ET.Element) -> None:
    for color_name, value in {
        "dk1": "101010",
        "lt1": "FAFAFA",
        "dk2": "202020",
        "lt2": "E5E7EB",
        "accent1": "112233",
        "accent2": "445566",
        "accent3": "778899",
    }.items():
        _set_scheme_color(root, color_name, value)

    _set_latin_font(root, "majorFont", "Aptos Display")
    _set_latin_font(root, "minorFont", "Aptos")


def _set_scheme_color(root: ET.Element, color_name: str, value: str) -> None:
    color = root.find(f".//a:clrScheme/a:{color_name}", NS)
    assert color is not None

    srgb = color.find("./a:srgbClr", NS)
    if srgb is not None:
        srgb.set("val", value)
        return

    system = color.find("./a:sysClr", NS)
    assert system is not None
    system.set("lastClr", value)


def _set_latin_font(root: ET.Element, scheme_name: str, typeface: str) -> None:
    font = root.find(f".//a:fontScheme/a:{scheme_name}/a:latin", NS)
    assert font is not None
    font.set("typeface", typeface)


def test_template_ingestion_full_flow(tmp_path: Path) -> None:
    template_path = tmp_path / "brand-template.pptx"
    create_test_template(template_path)

    learn_result = invoke(["learn", str(template_path)])
    assert learn_result.exit_code == 0
    learn_payload = parse_last_json_line(learn_result.output)
    manifest_path = tmp_path / "brand-template.manifest.json"
    assert learn_payload == {
        "ok": True,
        "data": {
            "source": "brand-template.pptx",
            "layouts_found": 11,
            "usable_layouts": 11,
        },
    }

    inspect_result = invoke(["inspect", str(manifest_path)])
    assert inspect_result.exit_code == 0
    inspect_payload = json.loads(inspect_result.output)
    assert inspect_payload["data"]["source"] == "brand-template.pptx"
    assert inspect_payload["data"]["layouts_found"] == 11
    assert inspect_payload["data"]["usable_layouts"] >= 3
    assert any(layout["slug"] == "title_slide" for layout in inspect_payload["data"]["layouts"])
    assert any(layout["slug"] == "two_content" for layout in inspect_payload["data"]["layouts"])
    assert any(layout["slug"] == "blank" for layout in inspect_payload["data"]["layouts"])

    deck_path = tmp_path / "deck.json"
    init_result = invoke(["init", str(deck_path), "--template", str(manifest_path)])
    assert init_result.exit_code == 0
    init_payload = json.loads(init_result.output)
    assert init_payload["ok"] is True
    assert init_payload["data"]["template_manifest"] == "brand-template.manifest.json"

    for layout in ["title_slide", "two_content"]:
        slide_add_result = invoke(["slide", "add", str(deck_path), "--layout", layout])
        assert slide_add_result.exit_code == 0
        assert json.loads(slide_add_result.output)["ok"] is True

    slot_updates = [
        ("0", "heading", "Brand Title"),
        ("0", "subheading", "Template subtitle"),
        ("1", "heading", "Agenda"),
        ("1", "col1", "Left column"),
        ("1", "col2", "Right column"),
    ]
    for slide_ref, slot_name, text in slot_updates:
        slot_result = invoke(
            [
                "slot",
                "set",
                str(deck_path),
                "--slide",
                slide_ref,
                "--slot",
                slot_name,
                "--text",
                text,
            ]
        )
        assert slot_result.exit_code == 0
        assert json.loads(slot_result.output)["ok"] is True

    output_path = tmp_path / "deck.pptx"
    build_result = invoke(["build", str(deck_path), "-o", str(output_path)])
    assert build_result.exit_code == 0
    assert json.loads(build_result.output)["ok"] is True
    assert build_result.stderr == ""

    presentation = Presentation(str(output_path))
    assert len(presentation.slides) == 2
    assert presentation.slides[0].slide_layout.name == "Title Slide"
    assert presentation.slides[1].slide_layout.name == "Two Content"

    all_text = collect_text_values(output_path)
    assert "Brand Title" in all_text
    assert "Template subtitle" in all_text
    assert "Agenda" in all_text
    assert "Left column" in all_text
    assert "Right column" in all_text


def test_template_changed_warning(tmp_path: Path) -> None:
    template_path = tmp_path / "brand-template.pptx"
    create_test_template(template_path)

    learn_result = invoke(["learn", str(template_path)])
    assert learn_result.exit_code == 0
    manifest_path = tmp_path / "brand-template.manifest.json"

    deck_path = tmp_path / "deck.json"
    init_result = invoke(["init", str(deck_path), "--template", str(manifest_path)])
    assert init_result.exit_code == 0

    slide_add_result = invoke(["slide", "add", str(deck_path), "--layout", "title_slide"])
    assert slide_add_result.exit_code == 0

    slot_result = invoke(
        [
            "slot",
            "set",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "heading",
            "--text",
            "Changed template warning",
        ]
    )
    assert slot_result.exit_code == 0

    patch_template_theme(template_path, accent1="AA5500")

    output_path = tmp_path / "deck.pptx"
    build_result = invoke(["build", str(deck_path), "-o", str(output_path)])
    assert build_result.exit_code == 0
    assert json.loads(build_result.stdout)["ok"] is True

    warning = json.loads(build_result.stderr)
    assert warning["warning"]["code"] == TEMPLATE_CHANGED
    assert warning["data"]["template"] == str(template_path)
    assert output_path.exists()
    assert collect_text_values(output_path).count("Changed template warning") == 1


def test_multi_layout_template(tmp_path: Path) -> None:
    template_path = tmp_path / "brand-template.pptx"
    create_test_template(
        template_path,
        layouts={
            0: "Title Slide",
            3: "Two Content",
            6: "Blank",
        },
    )

    learn_result = invoke(["learn", str(template_path)])
    assert learn_result.exit_code == 0
    manifest_path = tmp_path / "brand-template.manifest.json"
    inspect_payload = json.loads(invoke(["inspect", str(manifest_path)]).output)
    layout_slugs = {layout["slug"] for layout in inspect_payload["data"]["layouts"]}
    assert {"title_slide", "two_content", "blank"} <= layout_slugs

    deck_path = tmp_path / "deck.json"
    init_result = invoke(["init", str(deck_path), "--template", str(manifest_path)])
    assert init_result.exit_code == 0

    for layout in ["title_slide", "two_content", "blank"]:
        slide_add_result = invoke(["slide", "add", str(deck_path), "--layout", layout])
        assert slide_add_result.exit_code == 0

    for slide_ref, slot_name, text in [
        ("0", "heading", "Multi-layout title"),
        ("1", "heading", "Columns"),
        ("1", "col1", "Column A"),
        ("1", "col2", "Column B"),
    ]:
        slot_result = invoke(
            [
                "slot",
                "set",
                str(deck_path),
                "--slide",
                slide_ref,
                "--slot",
                slot_name,
                "--text",
                text,
            ]
        )
        assert slot_result.exit_code == 0

    output_path = tmp_path / "multi-layout.pptx"
    build_result = invoke(["build", str(deck_path), "-o", str(output_path)])
    assert build_result.exit_code == 0
    assert build_result.stderr == ""

    presentation = Presentation(str(output_path))
    assert [slide.slide_layout.name for slide in presentation.slides] == [
        "Title Slide",
        "Two Content",
        "Blank",
    ]

def test_template_build_swaps_inverted_quote_and_attribution_placeholders(tmp_path: Path) -> None:
    template_path = tmp_path / "quote-template.pptx"
    create_test_template(template_path)

    manifest_path = tmp_path / "quote-template.manifest.json"
    manifest_path.write_text(
        json.dumps(
            {
                "source": "quote-template.pptx",
                "source_hash": hashlib.sha256(template_path.read_bytes()).hexdigest(),
                "theme": {
                    "colors": {
                        "primary": "#112233",
                        "secondary": "#445566",
                        "accent": "#778899",
                        "text": "#101010",
                        "heading_text": "#202020",
                        "subtle_text": "#E5E7EB",
                        "background": "#FAFAFA",
                    },
                    "fonts": {
                        "heading": "Aptos Display",
                        "body": "Aptos",
                    },
                    "spacing": {
                        "base_unit": 10.0,
                        "margin": 60.0,
                        "gutter": 20.0,
                    },
                },
                "slide_masters": [
                    {
                        "index": 0,
                        "name": "Slide Master 1",
                        "layouts": [
                            {
                                "index": 0,
                                "master_index": 0,
                                "name": "Title Slide",
                                "slug": "quote_layout",
                                "usable": True,
                                "placeholders": [
                                    {
                                        "idx": 0,
                                        "type": "BODY",
                                        "name": "Large Quote",
                                        "bounds": {"x": 72, "y": 96, "w": 576, "h": 240},
                                    },
                                    {
                                        "idx": 1,
                                        "type": "BODY",
                                        "name": "Small Attribution",
                                        "bounds": {"x": 72, "y": 360, "w": 576, "h": 48},
                                    },
                                ],
                                "slot_mapping": {
                                    "quote": 1,
                                    "attribution": 0,
                                },
                            }
                        ],
                    }
                ],
            },
            indent=2,
        ),
        encoding="utf-8",
    )

    deck_path = tmp_path / "deck.json"
    init_result = invoke(["init", str(deck_path), "--template", str(manifest_path)])
    assert init_result.exit_code == 0

    slide_add_result = invoke(["slide", "add", str(deck_path), "--layout", "quote_layout"])
    assert slide_add_result.exit_code == 0

    for slot_name, text in [("quote", "Stay hungry, stay foolish."), ("attribution", "Steve Jobs")]:
        slot_result = invoke(
            [
                "slot",
                "set",
                str(deck_path),
                "--slide",
                "0",
                "--slot",
                slot_name,
                "--text",
                text,
            ]
        )
        assert slot_result.exit_code == 0

    output_path = tmp_path / "quote-deck.pptx"
    build_result = invoke(["build", str(deck_path), "-o", str(output_path)])
    assert build_result.exit_code == 0
    assert build_result.stderr == ""

    presentation = Presentation(str(output_path))
    slide = presentation.slides[0]
    assert slide.placeholders[0].text_frame.text == "Stay hungry, stay foolish."
    assert slide.placeholders[1].text_frame.text == "Steve Jobs"


def test_template_build_applies_computed_font_sizes_to_placeholder_text(tmp_path: Path) -> None:
    template_path = tmp_path / "brand-template.pptx"
    create_test_template(template_path)

    learn_result = invoke(["learn", str(template_path)])
    assert learn_result.exit_code == 0
    manifest_path = tmp_path / "brand-template.manifest.json"
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    title_layout = layout_by_slug(manifest, "title_slide")
    two_content_layout = layout_by_slug(manifest, "two_content")

    deck_path = tmp_path / "deck.json"
    init_result = invoke(["init", str(deck_path), "--template", str(manifest_path)])
    assert init_result.exit_code == 0

    for layout in ["title_slide", "two_content"]:
        slide_add_result = invoke(["slide", "add", str(deck_path), "--layout", layout])
        assert slide_add_result.exit_code == 0

    slot_updates = [
        (
            "0",
            "heading",
            "A very long template title that should shrink to stay inside the learned title placeholder bounds",
        ),
        (
            "1",
            "col1",
            "Point A\nPoint B\nPoint C\nPoint D\nPoint E\nPoint F\nPoint G\nPoint H",
        ),
    ]
    for slide_ref, slot_name, text in slot_updates:
        slot_result = invoke(
            [
                "slot",
                "set",
                str(deck_path),
                "--slide",
                slide_ref,
                "--slot",
                slot_name,
                "--text",
                text,
            ]
        )
        assert slot_result.exit_code == 0

    output_path = tmp_path / "deck-font-sizes.pptx"
    build_result = invoke(["build", str(deck_path), "-o", str(output_path)])
    assert build_result.exit_code == 0
    assert json.loads(build_result.output)["ok"] is True

    deck_payload = json.loads(deck_path.read_text(encoding="utf-8"))
    computed_payload = json.loads((tmp_path / "deck.computed.json").read_text(encoding="utf-8"))
    title_node_id = next(
        node["node_id"] for node in deck_payload["slides"][0]["nodes"] if node.get("slot_binding") == "heading"
    )
    body_node_id = next(
        node["node_id"] for node in deck_payload["slides"][1]["nodes"] if node.get("slot_binding") == "col1"
    )
    title_font_size = computed_payload["slides"][0]["computed"][title_node_id]["font_size_pt"]
    body_font_size = computed_payload["slides"][1]["computed"][body_node_id]["font_size_pt"]

    presentation = Presentation(str(output_path))
    title_run = presentation.slides[0].placeholders[title_layout["slot_mapping"]["heading"]].text_frame.paragraphs[0].runs[0]
    body_run = presentation.slides[1].placeholders[two_content_layout["slot_mapping"]["col1"]].text_frame.paragraphs[0].runs[0]

    assert title_run.font.size == Pt(title_font_size)
    assert body_run.font.size == Pt(body_font_size)
    assert title_run.font.size.pt == pytest.approx(title_font_size)
    assert body_run.font.size.pt == pytest.approx(body_font_size)
