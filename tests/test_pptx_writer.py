from __future__ import annotations

import base64
import json
from collections.abc import Iterable
from pathlib import Path
from xml.etree import ElementTree as ET
from zipfile import ZipFile

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from agent_slides.engine.reflow import reflow_deck
from agent_slides.errors import AgentSlidesError, FILE_NOT_FOUND, SCHEMA_ERROR
from agent_slides.io import read_template_manifest
from agent_slides.io.pptx_writer import write_pptx
from agent_slides.model.types import (
    BlockPosition,
    ChartSeries,
    ChartSpec,
    ChartStyle,
    ComputedNode,
    Counters,
    Deck,
    EMU_PER_POINT,
    Node,
    NodeContent,
    ScatterPoint,
    ScatterSeries,
    ShapeSpec,
    Slide,
    TableSpec,
    TextBlock,
    TextRun,
)
from tests.image_helpers import write_png

PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9Wn0K4sAAAAASUVORK5CYII="
)
CHART_NS = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}
DRAWING_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def build_deck() -> Deck:
    return Deck(
        deck_id="deck-1",
        revision=4,
        slides=[
            Slide(
                slide_id="s-1",
                layout="title",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="heading",
                        type="text",
                        content="Heading line 1\nHeading line 2",
                    ),
                    Node(
                        node_id="n-2",
                        slot_binding=None,
                        type="text",
                        content="Should not render",
                    ),
                    Node(
                        node_id="n-3",
                        slot_binding="body",
                        type="text",
                        content="",
                    ),
                ],
                computed={
                    "n-1": ComputedNode(
                        x=72.0,
                        y=54.0,
                        width=576.0,
                        height=96.0,
                        font_size_pt=28.0,
                        font_family="Aptos",
                        color="#112233",
                        bg_color="#F1E2D3",
                        font_bold=True,
                        revision=4,
                    ),
                    "n-2": ComputedNode(
                        x=10.0,
                        y=10.0,
                        width=100.0,
                        height=20.0,
                        font_size_pt=12.0,
                        font_family="Aptos",
                        color="#000000",
                        bg_color=None,
                        font_bold=False,
                        revision=4,
                    ),
                    "n-3": ComputedNode(
                        x=90.0,
                        y=180.0,
                        width=400.0,
                        height=40.0,
                        font_size_pt=18.0,
                        font_family="Calibri",
                        color="#445566",
                        bg_color=None,
                        font_bold=False,
                        revision=4,
                    ),
                },
            ),
            Slide(
                slide_id="s-2",
                layout="body",
                nodes=[
                    Node(
                        node_id="n-4",
                        slot_binding="body",
                        type="text",
                        content="Computed missing is skipped",
                    )
                ],
                computed={},
            ),
        ],
        counters=Counters(slides=2, nodes=4),
    )


def open_presentation(path: Path) -> Presentation:
    return Presentation(path)


def write_test_png(path: Path) -> None:
    path.write_bytes(PNG_1X1)


def chart_computed(
    *,
    revision: int = 1,
    x: float = 72.0,
    y: float = 108.0,
    width: float = 576.0,
    height: float = 288.0,
) -> ComputedNode:
    return ComputedNode(
        x=x,
        y=y,
        width=width,
        height=height,
        font_size_pt=0.0,
        font_family="Aptos",
        color="#112233",
        bg_color=None,
        font_bold=False,
        revision=revision,
        content_type="chart",
    )


def shape_computed(
    *,
    revision: int = 1,
    x: float = 60.0,
    y: float = 120.0,
    width: float = 180.0,
    height: float = 120.0,
) -> ComputedNode:
    return ComputedNode(
        x=x,
        y=y,
        width=width,
        height=height,
        font_size_pt=0.0,
        font_family="",
        color="#000000",
        bg_color=None,
        font_bold=False,
        revision=revision,
        content_type="shape",
    )


def first_chart_shape(presentation: Presentation):
    return next(shape for shape in presentation.slides[0].shapes if shape.shape_type == MSO_SHAPE_TYPE.CHART)


def first_table_shape(presentation: Presentation):
    return next(shape for shape in presentation.slides[0].shapes if shape.shape_type == MSO_SHAPE_TYPE.TABLE)


def read_chart_xml(path: Path, *, index: int = 1) -> ET.Element:
    with ZipFile(path) as archive:
        payload = archive.read(f"ppt/charts/chart{index}.xml")
    return ET.fromstring(payload)


def read_slide_xml(path: Path, *, index: int = 1) -> ET.Element:
    with ZipFile(path) as archive:
        payload = archive.read(f"ppt/slides/slide{index}.xml")
    return ET.fromstring(payload)


def create_template_manifest(tmp_path: Path) -> tuple[Path, Path, dict[str, object]]:
    template_path = tmp_path / "template.pptx"
    manifest_path = tmp_path / "template.manifest.json"

    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[0])
    presentation.slides.add_slide(presentation.slide_layouts[1])
    presentation.slides.add_slide(presentation.slide_layouts[1])
    presentation.save(template_path)

    result = read_template_manifest(template_path, manifest_path)
    return template_path, result.manifest_path, result.manifest


def find_layout(manifest: dict[str, object], required_slots: Iterable[str]) -> dict[str, object]:
    required = set(required_slots)
    for master in manifest["slide_masters"]:
        for layout in master["layouts"]:
            slot_mapping = layout["slot_mapping"]
            if required.issubset(slot_mapping):
                return layout
    raise AssertionError(f"No layout found with slots {sorted(required)}")


@pytest.mark.parametrize(
    ("shape_type", "expected_auto_shape"),
    [
        ("rectangle", MSO_SHAPE.RECTANGLE),
        ("rounded_rectangle", MSO_SHAPE.ROUNDED_RECTANGLE),
        ("oval", MSO_SHAPE.OVAL),
        ("arrow", MSO_SHAPE.RIGHT_ARROW),
        ("chevron", MSO_SHAPE.CHEVRON),
    ],
)
def test_write_pptx_renders_auto_shape_primitives(tmp_path: Path, shape_type: str, expected_auto_shape) -> None:
    output_path = tmp_path / f"{shape_type}.pptx"
    deck = Deck(
        deck_id=f"deck-{shape_type}",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title",
                nodes=[
                    Node(
                        node_id="n-shape",
                        type="shape",
                        shape_spec=ShapeSpec(
                            shape_type=shape_type,
                            fill_color="#F2F2F2",
                            line_color="#1A73E8",
                            line_width=2.0,
                            corner_radius=8.0,
                            opacity=0.8,
                        ),
                        style_overrides={"x": 60.0, "y": 120.0, "width": 180.0, "height": 120.0, "z_index": -1},
                    )
                ],
                computed={"n-shape": shape_computed()},
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    shape = presentation.slides[0].shapes[0]

    assert shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    assert shape.auto_shape_type == expected_auto_shape
    assert shape.left == int(60.0 * EMU_PER_POINT)
    assert shape.top == int(120.0 * EMU_PER_POINT)
    assert shape.width == int(180.0 * EMU_PER_POINT)
    assert shape.height == int(120.0 * EMU_PER_POINT)
    assert shape.fill.fore_color.rgb == RGBColor.from_string("F2F2F2")
    assert shape.line.color.rgb == RGBColor.from_string("1A73E8")
    if shape_type == "rounded_rectangle":
        assert shape.adjustments[0] == pytest.approx((8.0 * 2.0) / 120.0, abs=1e-4)


def test_write_pptx_renders_line_shape_primitive(tmp_path: Path) -> None:
    output_path = tmp_path / "line.pptx"
    deck = Deck(
        deck_id="deck-line",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title",
                nodes=[
                    Node(
                        node_id="n-line",
                        type="shape",
                        shape_spec=ShapeSpec(
                            shape_type="line",
                            line_color="#333333",
                            line_width=1.5,
                            dash="dash",
                        ),
                        style_overrides={"x": 60.0, "y": 300.0, "width": 600.0, "height": 0.0, "z_index": -1},
                    )
                ],
                computed={
                    "n-line": shape_computed(x=60.0, y=300.0, width=600.0, height=0.0),
                },
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    shape = presentation.slides[0].shapes[0]

    assert shape.shape_type == MSO_SHAPE_TYPE.LINE
    assert shape.begin_x == int(60.0 * EMU_PER_POINT)
    assert shape.end_x == int(660.0 * EMU_PER_POINT)
    assert shape.line.color.rgb == RGBColor.from_string("333333")
    assert shape.line.width.pt == pytest.approx(1.5)
    assert shape.line.dash_style == MSO_LINE_DASH_STYLE.DASH


def test_write_pptx_renders_shapes_before_text_nodes(tmp_path: Path) -> None:
    output_path = tmp_path / "shape-z-order.pptx"
    deck = Deck(
        deck_id="deck-z-order",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title",
                nodes=[
                    Node(
                        node_id="n-text",
                        slot_binding="heading",
                        type="text",
                        content="Foreground title",
                    ),
                    Node(
                        node_id="n-shape",
                        type="shape",
                        shape_spec=ShapeSpec(
                            shape_type="rectangle",
                            fill_color="#F2F2F2",
                            line_color="#CCCCCC",
                            shadow=True,
                            opacity=0.9,
                        ),
                        style_overrides={"x": 48.0, "y": 96.0, "width": 624.0, "height": 180.0, "z_index": -1},
                    ),
                ],
                computed={
                    "n-text": ComputedNode(
                        x=72.0,
                        y=54.0,
                        width=576.0,
                        height=80.0,
                        font_size_pt=28.0,
                        font_family="Aptos",
                        color="#112233",
                        bg_color=None,
                        font_bold=True,
                        revision=1,
                    ),
                    "n-shape": shape_computed(x=48.0, y=96.0, width=624.0, height=180.0),
                },
            )
        ],
        counters=Counters(slides=1, nodes=2),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    assert presentation.slides[0].shapes[0].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    assert presentation.slides[0].shapes[1].has_text_frame is True


def test_write_pptx_renders_pattern_nodes_as_shape_and_text_primitives(tmp_path: Path) -> None:
    output_path = tmp_path / "pattern.pptx"
    deck = Deck(
        deck_id="deck-pattern",
        theme="default",
        design_rules="default",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title_content",
                nodes=[
                    Node(node_id="n-1", slot_binding="heading", type="text", content="Executive summary"),
                    Node(
                        node_id="n-2",
                        slot_binding="body",
                        type="pattern",
                        pattern_spec={
                            "pattern_type": "kpi-row",
                            "data": [
                                {"value": "87%", "label": "Adoption"},
                                {"value": "3.2x", "label": "ROI"},
                            ],
                        },
                    ),
                ],
                computed={},
            )
        ],
        counters=Counters(slides=1, nodes=2),
    )
    reflow_deck(deck)

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    shapes = list(presentation.slides[0].shapes)

    assert len(shapes) >= 5
    assert any(shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE for shape in shapes)
    assert any(shape.has_text_frame and shape.text_frame.text == "87%" for shape in shapes)
    assert any(shape.has_text_frame and shape.text_frame.text == "Adoption" for shape in shapes)

    slide_xml = read_slide_xml(output_path)
    assert slide_xml.find(".//a:effectLst/a:outerShdw", DRAWING_NS) is not None


@pytest.mark.parametrize(
    ("chart_type", "pptx_type", "series"),
    [
        (
            "bar",
            XL_CHART_TYPE.BAR_CLUSTERED,
            [
                ("North", [12.0, 18.0, 16.0]),
                ("South", [9.0, 14.0, 20.0]),
            ],
        ),
        ("column", XL_CHART_TYPE.COLUMN_CLUSTERED, [("Revenue", [8.0, 13.0, 21.0])]),
        (
            "line",
            XL_CHART_TYPE.LINE,
            [
                ("Plan", [10.0, 11.0, 12.0]),
                ("Actual", [9.5, 12.0, 13.5]),
                ("Stretch", [11.0, 13.0, 15.0]),
            ],
        ),
        ("pie", XL_CHART_TYPE.PIE, [("Share", [55.0, 30.0, 15.0])]),
        ("area", XL_CHART_TYPE.AREA, [("Pipeline", [18.0, 16.0, 22.0])]),
        ("doughnut", XL_CHART_TYPE.DOUGHNUT, [("Mix", [40.0, 35.0, 25.0])]),
    ],
)
def test_write_pptx_renders_native_category_charts(tmp_path: Path, chart_type: str, pptx_type, series) -> None:
    output_path = tmp_path / f"{chart_type}.pptx"
    categories = ["Q1", "Q2", "Q3"]
    chart_node = Node(
        node_id="n-chart",
        slot_binding="body",
        type="chart",
        chart_spec=ChartSpec(
            chart_type=chart_type,
            title=f"{chart_type.title()} chart",
            categories=categories,
            series=[ChartSeries(name=name, values=values) for name, values in series],
            style=ChartStyle(
                has_legend=chart_type != "doughnut",
                series_colors=["#FF6B35", "#004E89", "#2A9D8F"],
            ),
        ),
    )
    deck = Deck(
        deck_id=f"deck-{chart_type}",
        slides=[
            Slide(
                slide_id="s-1",
                layout="content",
                nodes=[chart_node],
                computed={"n-chart": chart_computed()},
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    chart_shape = first_chart_shape(presentation)
    chart = chart_shape.chart

    assert chart_shape.left == int(72.0 * EMU_PER_POINT)
    assert chart_shape.top == int(108.0 * EMU_PER_POINT)
    assert chart_shape.width == int(576.0 * EMU_PER_POINT)
    assert chart_shape.height == int(288.0 * EMU_PER_POINT)
    assert chart.chart_type == pptx_type
    assert chart.has_title is True
    assert chart.chart_title.text_frame.text == f"{chart_type.title()} chart"
    assert chart.has_legend is (chart_type != "doughnut")
    assert list(chart.plots[0].categories) == categories
    assert [series_item.name for series_item in chart.series] == [name for name, _ in series]
    assert [list(series_item.values) for series_item in chart.series] == [values for _, values in series]
    assert chart.series[0].format.fill.fore_color.rgb == RGBColor.from_string("FF6B35")
    if len(chart.series) > 1:
        assert chart.series[1].format.fill.fore_color.rgb == RGBColor.from_string("004E89")


def test_write_pptx_renders_native_scatter_chart_with_read_back_data(tmp_path: Path) -> None:
    output_path = tmp_path / "scatter.pptx"
    deck = Deck(
        deck_id="deck-scatter",
        slides=[
            Slide(
                slide_id="s-1",
                layout="content",
                nodes=[
                    Node(
                        node_id="n-chart",
                        slot_binding="body",
                        type="chart",
                        chart_spec=ChartSpec(
                            chart_type="scatter",
                            title="Correlation",
                            scatter_series=[
                                ScatterSeries(
                                    name="Observed",
                                    points=[
                                        ScatterPoint(x=0.7, y=2.7),
                                        ScatterPoint(x=1.8, y=3.2),
                                        ScatterPoint(x=2.6, y=0.8),
                                    ],
                                ),
                                ScatterSeries(
                                    name="Projected",
                                    points=[
                                        ScatterPoint(x=1.3, y=3.7),
                                        ScatterPoint(x=2.7, y=2.3),
                                        ScatterPoint(x=1.6, y=1.8),
                                    ],
                                ),
                            ],
                            style=ChartStyle(has_legend=False, series_colors=["#AA0000", "#0055CC"]),
                        ),
                    )
                ],
                computed={"n-chart": chart_computed(x=96.0, y=120.0, width=500.0, height=260.0)},
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    chart_shape = first_chart_shape(presentation)
    chart = chart_shape.chart
    xml = read_chart_xml(output_path)

    assert chart_shape.left == int(96.0 * EMU_PER_POINT)
    assert chart_shape.top == int(120.0 * EMU_PER_POINT)
    assert chart_shape.width == int(500.0 * EMU_PER_POINT)
    assert chart_shape.height == int(260.0 * EMU_PER_POINT)
    assert chart.chart_type == XL_CHART_TYPE.XY_SCATTER
    assert chart.has_title is True
    assert chart.chart_title.text_frame.text == "Correlation"
    assert chart.has_legend is False
    assert [series_item.name for series_item in chart.series] == ["Observed", "Projected"]
    assert [list(series_item.values) for series_item in chart.series] == [
        [2.7, 3.2, 0.8],
        [3.7, 2.3, 1.8],
    ]
    assert chart.series[0].format.fill.fore_color.rgb == RGBColor.from_string("AA0000")
    assert chart.series[1].format.fill.fore_color.rgb == RGBColor.from_string("0055CC")
    assert [point.text for point in xml.findall(".//c:xVal//c:numCache/c:pt/c:v", CHART_NS)] == [
        "0.7",
        "1.8",
        "2.6",
        "1.3",
        "2.7",
        "1.6",
    ]
    assert [point.text for point in xml.findall(".//c:yVal//c:numCache/c:pt/c:v", CHART_NS)] == [
        "2.7",
        "3.2",
        "0.8",
        "3.7",
        "2.3",
        "1.8",
    ]


def test_write_pptx_renders_native_table_with_alignment_and_striping(tmp_path: Path) -> None:
    output_path = tmp_path / "table.pptx"
    deck = Deck(
        deck_id="deck-table",
        theme="default",
        slides=[
            Slide(
                slide_id="s-1",
                layout="content",
                nodes=[
                    Node(
                        node_id="n-table",
                        slot_binding="body",
                        type="table",
                        table_spec=TableSpec(
                            headers=["Metric", "Q1", "Q2"],
                            rows=[
                                ["Revenue", "$100K", "$150K"],
                                ["Users", "1000", "1500"],
                            ],
                            header_color="#1F4E79",
                            stripe=True,
                        ),
                    )
                ],
                computed={
                    "n-table": ComputedNode(
                        x=72.0,
                        y=108.0,
                        width=576.0,
                        height=216.0,
                        font_size_pt=0.0,
                        font_family="Aptos",
                        color="#112233",
                        bg_color="#FFFFFF",
                        font_bold=False,
                        revision=1,
                        content_type="table",
                    )
                },
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    table_shape = first_table_shape(presentation)
    table = table_shape.table

    assert table_shape.left == int(72.0 * EMU_PER_POINT)
    assert table_shape.top == int(108.0 * EMU_PER_POINT)
    assert table_shape.width == int(576.0 * EMU_PER_POINT)
    assert table_shape.height == int(216.0 * EMU_PER_POINT)
    assert table.cell(0, 0).text == "Metric"
    assert table.cell(1, 0).text == "Revenue"
    assert table.columns[0].width > table.columns[1].width
    assert table.cell(0, 0).text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
    assert table.cell(1, 0).text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT
    assert table.cell(1, 1).text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT
    assert table.cell(0, 0).fill.fore_color.rgb == RGBColor.from_string("1F4E79")
    assert table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold is True
    assert table.cell(2, 0).fill.fore_color.rgb != table.cell(1, 0).fill.fore_color.rgb


def test_write_pptx_applies_status_colors_to_table_cells(tmp_path: Path) -> None:
    output_path = tmp_path / "table-status.pptx"
    deck = Deck(
        deck_id="deck-table-status",
        theme="default",
        design_rules="default",
        slides=[
            Slide(
                slide_id="s-1",
                layout="content",
                nodes=[
                    Node(
                        node_id="n-table",
                        slot_binding="body",
                        type="table",
                        table_spec=TableSpec(
                            headers=["Item", "Status"],
                            rows=[
                                ["Migration", "Blocked"],
                                ["Launch", "In Progress"],
                                ["Ops", "Complete"],
                            ],
                        ),
                    )
                ],
                computed={
                    "n-table": ComputedNode(
                        x=72.0,
                        y=108.0,
                        width=576.0,
                        height=216.0,
                        font_size_pt=0.0,
                        font_family="Aptos",
                        color="#112233",
                        bg_color="#FFFFFF",
                        font_bold=False,
                        revision=1,
                        content_type="table",
                    )
                },
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    table = first_table_shape(open_presentation(output_path)).table

    assert table.cell(1, 1).fill.fore_color.rgb == RGBColor.from_string("FDE3E3")
    assert table.cell(2, 1).fill.fore_color.rgb == RGBColor.from_string("FFF1C7")
    assert table.cell(3, 1).fill.fore_color.rgb == RGBColor.from_string("DDF4E4")
    assert table.cell(1, 1).text_frame.paragraphs[0].runs[0].font.bold is True


def test_write_pptx_renders_expected_slides_and_shapes(tmp_path: Path) -> None:
    output_path = tmp_path / "deck.pptx"

    write_pptx(build_deck(), str(output_path))

    presentation = open_presentation(output_path)

    assert len(presentation.slides) == 2
    assert presentation.slide_width == Inches(10)
    assert presentation.slide_height == Inches(7.5)

    first_slide = presentation.slides[0]
    assert len(first_slide.shapes) == 2

    heading_box = first_slide.shapes[0]
    assert heading_box.left == int(72.0 * EMU_PER_POINT)
    assert heading_box.top == int(54.0 * EMU_PER_POINT)
    assert heading_box.width == int(576.0 * EMU_PER_POINT)
    assert heading_box.height == int(96.0 * EMU_PER_POINT)
    assert heading_box.fill.fore_color.rgb == RGBColor.from_string("F1E2D3")

    text_frame = heading_box.text_frame
    assert text_frame.word_wrap is True
    assert text_frame.auto_size == MSO_AUTO_SIZE.NONE
    assert [paragraph.text for paragraph in text_frame.paragraphs] == [
        "Heading line 1",
        "Heading line 2",
    ]

    first_run = text_frame.paragraphs[0].runs[0]
    assert first_run.font.name == "Aptos"
    assert first_run.font.size == Pt(28)
    assert first_run.font.bold is True
    assert first_run.font.color.rgb == RGBColor.from_string("112233")

    empty_box = first_slide.shapes[1]
    assert empty_box.text == ""
    empty_run = empty_box.text_frame.paragraphs[0].runs[0]
    assert empty_run.font.name == "Calibri"
    assert empty_run.font.size == Pt(18)
    assert empty_run.font.bold is False
    assert empty_run.font.color.rgb == RGBColor.from_string("445566")

    second_slide = presentation.slides[1]
    assert len(second_slide.shapes) == 0


def test_write_pptx_creates_valid_empty_presentation(tmp_path: Path) -> None:
    output_path = tmp_path / "empty.pptx"
    deck = Deck(deck_id="empty-deck")

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)

    assert output_path.exists()
    assert len(presentation.slides) == 0


def test_write_pptx_renders_structured_headings_and_bullets(tmp_path: Path) -> None:
    output_path = tmp_path / "structured.pptx"
    deck = Deck(
        deck_id="deck-structured",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="body",
                        type="text",
                        content=NodeContent(
                            blocks=[
                                TextBlock(type="heading", text="Highlights"),
                                TextBlock(type="bullet", text="First takeaway"),
                                TextBlock(type="bullet", text="Nested takeaway", level=1),
                            ]
                        ),
                    )
                ],
                computed={
                    "n-1": ComputedNode(
                        x=72.0,
                        y=54.0,
                        width=400.0,
                        height=140.0,
                        font_size_pt=20.0,
                        font_family="Aptos",
                        color="#112233",
                        bg_color=None,
                        font_bold=False,
                        revision=1,
                    )
                },
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    text_frame = presentation.slides[0].shapes[0].text_frame
    slide_xml = read_slide_xml(output_path)
    paragraphs = slide_xml.findall(".//a:p", DRAWING_NS)
    texts = [text.text for text in slide_xml.findall(".//a:t", DRAWING_NS)]

    assert [paragraph.text for paragraph in text_frame.paragraphs] == [
        "Highlights",
        "First takeaway",
        "Nested takeaway",
    ]
    assert text_frame.paragraphs[0].runs[0].font.size == Pt(27)
    assert text_frame.paragraphs[1].level == 0
    assert text_frame.paragraphs[2].level == 1
    assert texts == ["Highlights", "First takeaway", "Nested takeaway"]
    assert paragraphs[0].find("./a:pPr/a:buNone", DRAWING_NS) is not None
    assert paragraphs[0].find("./a:pPr/a:buChar", DRAWING_NS) is None
    assert paragraphs[1].find("./a:pPr/a:buChar", DRAWING_NS) is not None
    assert paragraphs[1].find("./a:pPr/a:buNone", DRAWING_NS) is None
    assert paragraphs[1].find("./a:pPr", DRAWING_NS).attrib["marL"] == str(Pt(18))
    assert paragraphs[1].find("./a:pPr", DRAWING_NS).attrib["indent"] == str(Pt(-18))
    assert paragraphs[2].find("./a:pPr/a:buChar", DRAWING_NS) is not None
    assert paragraphs[2].find("./a:pPr", DRAWING_NS).attrib["lvl"] == "1"
    assert paragraphs[2].find("./a:pPr", DRAWING_NS).attrib["marL"] == str(Pt(36))
    assert paragraphs[2].find("./a:pPr", DRAWING_NS).attrib["indent"] == str(Pt(-18))

def test_write_pptx_uses_block_positions_for_precise_text_boxes(tmp_path: Path) -> None:
    output_path = tmp_path / "block-positions.pptx"
    deck = Deck(
        deck_id="deck-block-positions",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="body",
                        type="text",
                        content=NodeContent(
                            blocks=[
                                TextBlock(type="heading", text="Highlights"),
                                TextBlock(type="bullet", text="First takeaway"),
                            ]
                        ),
                    )
                ],
                computed={
                    "n-1": ComputedNode(
                        x=100.0,
                        y=120.0,
                        width=260.0,
                        height=140.0,
                        font_size_pt=18.0,
                        font_family="Aptos",
                        color="#112233",
                        bg_color="#F1E2D3",
                        font_bold=False,
                        revision=1,
                        block_positions=[
                            BlockPosition(
                                block_index=0,
                                x=108.0,
                                y=128.0,
                                width=244.0,
                                height=35.2,
                                font_size_pt=32.0,
                            ),
                            BlockPosition(
                                block_index=1,
                                x=108.0,
                                y=173.2,
                                width=244.0,
                                height=21.6,
                                font_size_pt=18.0,
                            ),
                        ],
                    )
                },
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    slide = presentation.slides[0]
    slide_xml = read_slide_xml(output_path)
    paragraphs = slide_xml.findall(".//a:p", DRAWING_NS)

    assert len(slide.shapes) == 1
    text_box = slide.shapes[0]
    assert text_box.fill.fore_color.rgb == RGBColor.from_string("F1E2D3")
    assert text_box.left == int(100.0 * EMU_PER_POINT)
    assert text_box.top == int(120.0 * EMU_PER_POINT)
    assert text_box.width == int(260.0 * EMU_PER_POINT)
    assert text_box.height == int(140.0 * EMU_PER_POINT)
    assert text_box.text_frame.margin_left == int(8.0 * EMU_PER_POINT)
    assert text_box.text_frame.margin_top == int(8.0 * EMU_PER_POINT)
    assert text_box.text_frame.paragraphs[0].runs[0].font.size == Pt(32)
    assert text_box.text_frame.paragraphs[1].text == "First takeaway"
    assert text_box.text_frame.paragraphs[1].runs[0].font.size == Pt(18)
    assert paragraphs[0].find("./a:pPr/a:spcAft/a:spcPts", DRAWING_NS).attrib["val"] == "1000"
    assert any(paragraph.find("./a:pPr/a:buChar", DRAWING_NS) is not None for paragraph in paragraphs)


def test_write_pptx_renders_icon_nodes_and_icon_bullets(tmp_path: Path) -> None:
    output_path = tmp_path / "icons.pptx"
    deck = Deck(
        deck_id="deck-icons",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="body",
                        type="text",
                        content=NodeContent(
                            blocks=[
                                TextBlock(type="bullet", text="On track for Q3", icon="checkmark"),
                            ]
                        ),
                    ),
                    Node(
                        node_id="n-2",
                        type="icon",
                        icon_name="warning",
                        x=420.0,
                        y=120.0,
                        size=18.0,
                        color="#D93025",
                    ),
                ],
                computed={
                    "n-1": ComputedNode(
                        x=72.0,
                        y=100.0,
                        width=320.0,
                        height=90.0,
                        font_size_pt=18.0,
                        font_family="Aptos",
                        color="#1A73E8",
                        bg_color=None,
                        revision=1,
                    ),
                    "n-2": ComputedNode(
                        x=420.0,
                        y=120.0,
                        width=18.0,
                        height=18.0,
                        font_size_pt=0.0,
                        font_family="Aptos",
                        color="#D93025",
                        bg_color=None,
                        revision=1,
                        content_type="icon",
                        icon_svg_path="M12 2 L22 21 L2 21 Z M11 8 L13 8 L13 15 L11 15 Z M11 17 L13 17 L13 19 L11 19 Z",
                    ),
                },
            )
        ],
        counters=Counters(slides=1, nodes=2),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    slide = presentation.slides[0]
    text_frame = slide.shapes[0].text_frame
    slide_xml = read_slide_xml(output_path)
    custom_geometry = slide_xml.findall(".//a:custGeom", DRAWING_NS)
    paragraphs = slide_xml.findall(".//a:p", DRAWING_NS)
    fills = [slide.shapes[index].fill.fore_color.rgb for index in range(1, len(slide.shapes))]

    assert len(slide.shapes) >= 3
    assert len(custom_geometry) >= 2
    assert text_frame.paragraphs[0].text == "On track for Q3"
    assert text_frame.paragraphs[0].runs[0].font.color.rgb == RGBColor.from_string("1A73E8")
    assert RGBColor.from_string("1A73E8") in fills
    assert RGBColor.from_string("D93025") in fills
    assert paragraphs[0].find("./a:pPr/a:buNone", DRAWING_NS) is not None
    assert paragraphs[0].find("./a:pPr/a:buChar", DRAWING_NS) is None
    assert int(paragraphs[0].find("./a:pPr", DRAWING_NS).attrib["marL"]) > 0


def test_write_pptx_renders_inline_text_runs(tmp_path: Path) -> None:
    output_path = tmp_path / "inline-runs.pptx"
    deck = Deck(
        deck_id="deck-inline-runs",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="body",
                        type="text",
                        content=NodeContent(
                            blocks=[
                                TextBlock(
                                    type="paragraph",
                                    runs=[
                                        TextRun(text="Revenue grew "),
                                        TextRun(text="23%", bold=True, color="#1A73E8"),
                                        TextRun(text=" driven by "),
                                        TextRun(text="premium segment", italic=True, font_size=24, underline=True),
                                    ],
                                )
                            ]
                        ),
                    )
                ],
                computed={
                    "n-1": ComputedNode(
                        x=72.0,
                        y=54.0,
                        width=400.0,
                        height=80.0,
                        font_size_pt=18.0,
                        font_family="Aptos",
                        color="#112233",
                        bg_color=None,
                        font_bold=False,
                        revision=1,
                    )
                },
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    paragraph = presentation.slides[0].shapes[0].text_frame.paragraphs[0]
    slide_xml = read_slide_xml(output_path)
    runs = paragraph.runs
    xml_runs = slide_xml.findall(".//a:r", DRAWING_NS)

    assert paragraph.text == "Revenue grew 23% driven by premium segment"
    assert [run.text for run in runs] == [
        "Revenue grew ",
        "23%",
        " driven by ",
        "premium segment",
    ]
    assert runs[0].font.size == Pt(18)
    assert runs[1].font.bold is True
    assert runs[1].font.color.rgb == RGBColor.from_string("1A73E8")
    assert runs[3].font.italic is True
    assert runs[3].font.size == Pt(24)
    assert runs[3].font.underline is True
    assert xml_runs[3].find("./a:rPr", DRAWING_NS).attrib["u"] == "sng"

def test_write_pptx_applies_conditional_text_colors_from_design_rules(tmp_path: Path) -> None:
    output_path = tmp_path / "conditional-text.pptx"
    deck = Deck(
        deck_id="deck-conditional-text",
        design_rules="default",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="body",
                        type="text",
                        content=NodeContent(
                            blocks=[TextBlock(type="paragraph", text="Revenue +23% is urgent but margin -5%")]
                        ),
                    )
                ],
                computed={
                    "n-1": ComputedNode(
                        x=72.0,
                        y=54.0,
                        width=400.0,
                        height=80.0,
                        font_size_pt=18.0,
                        font_family="Aptos",
                        color="#112233",
                        bg_color=None,
                        font_bold=False,
                        revision=1,
                    )
                },
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    runs = open_presentation(output_path).slides[0].shapes[0].text_frame.paragraphs[0].runs

    assert any(run.text == "+23%" and run.font.color.rgb == RGBColor.from_string("1B8A2D") for run in runs)
    assert any(run.text == "-5%" and run.font.color.rgb == RGBColor.from_string("D32F2F") for run in runs)
    assert any(run.text == "urgent" and run.font.bold is True and run.font.color.rgb == RGBColor.from_string("C98E48") for run in runs)


def test_write_pptx_applies_conditional_point_colors_to_single_series_charts(tmp_path: Path) -> None:
    output_path = tmp_path / "conditional-chart.pptx"
    deck = Deck(
        deck_id="deck-conditional-chart",
        theme="default",
        design_rules="default",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-chart",
                        slot_binding="body",
                        type="chart",
                        chart_spec=ChartSpec(
                            chart_type="column",
                            categories=["Q1", "Q2", "Q3"],
                            series=[ChartSeries(name="Revenue", values=[12.0, -4.0, 18.0])],
                            style=ChartStyle(color_by_value=True),
                        ),
                    )
                ],
                computed={"n-chart": chart_computed()},
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    chart = first_chart_shape(open_presentation(output_path)).chart

    assert chart.series[0].points[0].format.fill.fore_color.rgb == RGBColor.from_string("1B8A2D")
    assert chart.series[0].points[1].format.fill.fore_color.rgb == RGBColor.from_string("D32F2F")


def test_write_pptx_applies_highlight_and_muted_chart_point_colors(tmp_path: Path) -> None:
    output_path = tmp_path / "highlight-chart.pptx"
    deck = Deck(
        deck_id="deck-highlight-chart",
        theme="default",
        design_rules="default",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-chart",
                        slot_binding="body",
                        type="chart",
                        chart_spec=ChartSpec(
                            chart_type="pie",
                            categories=["A", "B", "C"],
                            series=[ChartSeries(name="Share", values=[40.0, 35.0, 25.0])],
                            style=ChartStyle(highlight_index=1),
                        ),
                    )
                ],
                computed={"n-chart": chart_computed()},
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    chart = first_chart_shape(open_presentation(output_path)).chart

    assert chart.series[0].points[0].format.fill.fore_color.rgb == RGBColor.from_string("CFC8BD")
    assert chart.series[0].points[1].format.fill.fore_color.rgb == RGBColor.from_string("C98E48")
def test_write_pptx_renders_image_nodes_with_contain_fit(tmp_path: Path) -> None:
    image_path = write_png(tmp_path / "photo.png", width=200, height=100)
    output_path = tmp_path / "images.pptx"
    deck = Deck(
        deck_id="deck-images",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title_content",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="body",
                        type="image",
                        image_path=str(image_path),
                    )
                ],
                computed={
                    "n-1": ComputedNode(
                        x=100.0,
                        y=120.0,
                        width=240.0,
                        height=180.0,
                        revision=1,
                        image_fit="contain",
                    )
                },
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    picture = presentation.slides[0].shapes[0]

    assert picture.shape_type == MSO_SHAPE_TYPE.PICTURE
    assert picture.left == int(100.0 * EMU_PER_POINT)
    assert picture.top == int(150.0 * EMU_PER_POINT)
    assert picture.width == int(240.0 * EMU_PER_POINT)
    assert picture.height == int(120.0 * EMU_PER_POINT)
    assert picture.image.size == (200, 100)


def test_write_pptx_renders_image_nodes_with_stretch_fit(tmp_path: Path) -> None:
    image_path = write_png(tmp_path / "photo.png", width=80, height=160)
    output_path = tmp_path / "stretch.pptx"
    deck = Deck(
        deck_id="deck-images-stretch",
        slides=[
            Slide(
                slide_id="s-1",
                layout="content",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="body",
                        type="image",
                        image_path=str(image_path),
                        image_fit="stretch",
                    )
                ],
                computed={
                    "n-1": ComputedNode(
                        x=80.0,
                        y=90.0,
                        width=300.0,
                        height=140.0,
                        revision=1,
                        image_fit="stretch",
                    )
                },
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    picture = presentation.slides[0].shapes[0]

    assert picture.left == int(80.0 * EMU_PER_POINT)
    assert picture.top == int(90.0 * EMU_PER_POINT)
    assert picture.width == int(300.0 * EMU_PER_POINT)
    assert picture.height == int(140.0 * EMU_PER_POINT)


def test_write_pptx_renders_image_layout_variants(tmp_path: Path) -> None:
    output_path = tmp_path / "image-layouts.pptx"
    image_path = tmp_path / "pixel.png"
    write_test_png(image_path)

    deck = Deck(
        deck_id="deck-images",
        slides=[
            Slide(
                slide_id="s-1",
                layout="image_left",
                nodes=[
                    Node(node_id="n-1", slot_binding="image", type="image", image_path=str(image_path)),
                    Node(node_id="n-2", slot_binding="heading", type="text", content="Image left"),
                    Node(node_id="n-3", slot_binding="body", type="text", content="Caption copy"),
                ],
            ),
            Slide(
                slide_id="s-2",
                layout="image_right",
                nodes=[
                    Node(node_id="n-4", slot_binding="heading", type="text", content="Image right"),
                    Node(node_id="n-5", slot_binding="body", type="text", content="Mirrored split"),
                    Node(node_id="n-6", slot_binding="image", type="image", image_path=str(image_path)),
                ],
            ),
            Slide(
                slide_id="s-3",
                layout="hero_image",
                nodes=[
                    Node(node_id="n-7", slot_binding="image", type="image", image_path=str(image_path)),
                    Node(node_id="n-8", slot_binding="heading", type="text", content="Hero"),
                    Node(node_id="n-9", slot_binding="subheading", type="text", content="Overlay copy"),
                ],
            ),
            Slide(
                slide_id="s-4",
                layout="gallery",
                nodes=[
                    Node(node_id="n-10", slot_binding="heading", type="text", content="Gallery"),
                    Node(node_id="n-11", slot_binding="img1", type="image", image_path=str(image_path)),
                    Node(node_id="n-12", slot_binding="img2", type="image", image_path=str(image_path)),
                    Node(node_id="n-13", slot_binding="img3", type="image", image_path=str(image_path)),
                    Node(node_id="n-14", slot_binding="img4", type="image", image_path=str(image_path)),
                ],
            ),
        ],
        counters=Counters(slides=4, nodes=14),
    )

    reflow_deck(deck)
    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)

    assert len(presentation.slides) == 4
    assert sum(1 for shape in presentation.slides[0].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE) == 1
    assert sum(1 for shape in presentation.slides[1].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE) == 1
    assert sum(1 for shape in presentation.slides[2].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE) == 1
    assert sum(1 for shape in presentation.slides[3].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE) == 4

    hero_picture = next(
        shape for shape in presentation.slides[2].shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    assert hero_picture.left == 0
    assert hero_picture.top == 0
    assert hero_picture.width == int(720.0 * EMU_PER_POINT)
    assert hero_picture.height == int(540.0 * EMU_PER_POINT)


def test_write_pptx_clones_template_and_fills_native_placeholders(tmp_path: Path) -> None:
    template_path, manifest_path, manifest = create_template_manifest(tmp_path)
    output_path = tmp_path / "template-output.pptx"
    title_layout = find_layout(manifest, {"heading", "subheading"})
    body_layout = find_layout(manifest, {"heading", "body"})

    deck = Deck(
        deck_id="template-deck",
        template_manifest=str(manifest_path),
        slides=[
            Slide(
                slide_id="s-1",
                layout=title_layout["slug"],
                nodes=[
                    Node(node_id="n-1", slot_binding="heading", type="text", content="Template heading"),
                    Node(node_id="n-2", slot_binding="subheading", type="text", content="Line 1\nLine 2"),
                ],
            ),
            Slide(
                slide_id="s-2",
                layout=body_layout["slug"],
                nodes=[
                    Node(node_id="n-3", slot_binding="heading", type="text", content="Agenda"),
                    Node(node_id="n-4", slot_binding="body", type="text", content="Point A\nPoint B"),
                ],
            ),
        ],
        counters=Counters(slides=2, nodes=4),
    )

    write_pptx(deck, str(output_path))

    template = open_presentation(template_path)
    presentation = open_presentation(output_path)

    assert len(presentation.slides) == 2
    assert len(presentation.slide_masters) == len(template.slide_masters)
    assert presentation.slides[0].slide_layout.name == title_layout["name"]
    assert presentation.slides[1].slide_layout.name == body_layout["name"]

    title_placeholder = presentation.slides[0].placeholders[title_layout["slot_mapping"]["heading"]]
    subtitle_placeholder = presentation.slides[0].placeholders[title_layout["slot_mapping"]["subheading"]]
    body_placeholder = presentation.slides[1].placeholders[body_layout["slot_mapping"]["body"]]

    assert [paragraph.text for paragraph in subtitle_placeholder.text_frame.paragraphs] == ["Line 1", "Line 2"]
    assert [paragraph.text for paragraph in body_placeholder.text_frame.paragraphs] == ["Point A", "Point B"]
    assert title_placeholder.text_frame.paragraphs[0].runs[0].font.name is None
    assert title_placeholder.text_frame.paragraphs[0].runs[0].font.size is None

    with ZipFile(output_path) as archive:
        slide_parts = sorted(
            name
            for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
    assert slide_parts == ["ppt/slides/slide1.xml", "ppt/slides/slide2.xml"]


def test_write_pptx_applies_computed_font_styles_to_template_placeholders(tmp_path: Path) -> None:
    _, manifest_path, manifest = create_template_manifest(tmp_path)
    output_path = tmp_path / "template-font-styles.pptx"
    body_layout = find_layout(manifest, {"heading", "body"})

    deck = Deck(
        deck_id="template-font-styles",
        template_manifest=str(manifest_path),
        slides=[
            Slide(
                slide_id="s-1",
                layout=body_layout["slug"],
                nodes=[
                    Node(node_id="n-1", slot_binding="heading", type="text", content="A longer template heading"),
                    Node(
                        node_id="n-2",
                        slot_binding="body",
                        type="text",
                        content=NodeContent(
                            blocks=[
                                TextBlock(
                                    type="paragraph",
                                    runs=[
                                        TextRun(text="Revenue grew "),
                                        TextRun(text="23%", bold=True, color="#1A73E8"),
                                        TextRun(text=" driven by "),
                                        TextRun(text="premium segment", italic=True, font_size=24, underline=True),
                                    ],
                                )
                            ]
                        ),
                    ),
                ],
                computed={
                    "n-1": ComputedNode(
                        x=72.0,
                        y=54.0,
                        width=576.0,
                        height=72.0,
                        font_size_pt=26.0,
                        font_family="Aptos Display",
                        color="#112233",
                        bg_color=None,
                        font_bold=True,
                        revision=1,
                    ),
                    "n-2": ComputedNode(
                        x=72.0,
                        y=144.0,
                        width=576.0,
                        height=180.0,
                        font_size_pt=18.0,
                        font_family="Aptos",
                        color="#112233",
                        bg_color=None,
                        font_bold=False,
                        revision=1,
                    ),
                },
            )
        ],
        counters=Counters(slides=1, nodes=2),
    )

    write_pptx(deck, str(output_path))

    presentation = open_presentation(output_path)
    slide = presentation.slides[0]
    heading_placeholder = slide.placeholders[body_layout["slot_mapping"]["heading"]]
    body_placeholder = slide.placeholders[body_layout["slot_mapping"]["body"]]

    heading_run = heading_placeholder.text_frame.paragraphs[0].runs[0]
    body_runs = body_placeholder.text_frame.paragraphs[0].runs

    # Template runs inherit native placeholder formatting; font.name and font.size
    # are None at run level (inherited from the placeholder/layout master).
    assert heading_run.font.name is None
    assert heading_run.font.size is None
    assert body_runs[0].font.name is None
    assert body_runs[0].font.size is None
    # Explicit TextRun overrides ARE applied:
    assert body_runs[1].font.bold is True
    assert body_runs[1].font.color.rgb == RGBColor.from_string("1A73E8")
    assert body_runs[3].font.italic is True
    assert body_runs[3].font.size == Pt(24)
    assert body_runs[3].font.underline is True


def test_write_pptx_warns_when_template_hash_changes(tmp_path: Path, capsys: pytest.CaptureFixture[str]) -> None:
    _, manifest_path, manifest = create_template_manifest(tmp_path)
    output_path = tmp_path / "hash-mismatch.pptx"
    title_layout = find_layout(manifest, {"heading", "subheading"})

    payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    payload["source_hash"] = "0" * 64
    manifest_path.write_text(f"{json.dumps(payload, indent=2)}\n", encoding="utf-8")

    deck = Deck(
        deck_id="template-deck",
        template_manifest=str(manifest_path),
        slides=[
            Slide(
                slide_id="s-1",
                layout=title_layout["slug"],
                nodes=[Node(node_id="n-1", slot_binding="heading", type="text", content="Updated title")],
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )

    write_pptx(deck, str(output_path))

    warning = json.loads(capsys.readouterr().err)
    assert warning["warning"]["code"] == "TEMPLATE_CHANGED"
    assert output_path.exists()


def test_write_pptx_reports_missing_template_file_from_manifest(tmp_path: Path) -> None:
    manifest_path = tmp_path / "missing-template.manifest.json"
    manifest_path.write_text(
        json.dumps(
            {
                "source": "missing-template.pptx",
                "source_hash": "abc123",
                "slide_masters": [
                    {
                        "layouts": [
                            {
                                "name": "Title Slide",
                                "slug": "title_slide",
                                "index": 0,
                                "master_index": 0,
                                "slot_mapping": {"heading": 0},
                            }
                        ]
                    }
                ],
            },
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )

    deck = Deck(
        deck_id="missing-template",
        template_manifest=str(manifest_path),
        slides=[Slide(slide_id="s-1", layout="title_slide")],
    )

    with pytest.raises(AgentSlidesError) as exc_info:
        write_pptx(deck, str(tmp_path / "missing-template-output.pptx"))

    assert exc_info.value.code == FILE_NOT_FOUND


def test_write_pptx_reports_manifest_layout_index_out_of_range(tmp_path: Path) -> None:
    _, manifest_path, manifest = create_template_manifest(tmp_path)
    title_layout = find_layout(manifest, {"heading", "subheading"})
    payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    payload["slide_masters"][0]["layouts"][0]["index"] = 99
    manifest_path.write_text(f"{json.dumps(payload, indent=2)}\n", encoding="utf-8")

    deck = Deck(
        deck_id="bad-layout-index",
        template_manifest=str(manifest_path),
        slides=[Slide(slide_id="s-1", layout=title_layout["slug"])],
    )

    with pytest.raises(AgentSlidesError) as exc_info:
        write_pptx(deck, str(tmp_path / "bad-layout-index.pptx"))

    assert exc_info.value.code == SCHEMA_ERROR
