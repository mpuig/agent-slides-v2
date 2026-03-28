from __future__ import annotations

import contextlib
import json
import os
import signal
import socket
import threading
import time
import urllib.error
import urllib.request
from pathlib import Path
from uuid import UUID

import pytest
from click.testing import CliRunner, Result
from pptx import Presentation

from agent_slides.cli import cli
from agent_slides.errors import (
    CHART_DATA_ERROR,
    FILE_EXISTS,
    FILE_NOT_FOUND,
    INVALID_CHART_TYPE,
    INVALID_LAYOUT,
    INVALID_NODE_TYPE,
    INVALID_SLIDE,
    INVALID_SLOT,
    SCHEMA_ERROR,
    THEME_NOT_FOUND,
    UNBOUND_NODES,
)
from agent_slides.io import computed_sidecar_path, init_deck, read_computed_deck, read_deck, write_computed_deck
from agent_slides.model import BuiltinLayoutProvider, Counters, Deck, Node, Slide, get_layout, list_layouts
from agent_slides.model.design_rules import load_design_rules
from agent_slides.model.types import ComputedNode
from tests.image_helpers import write_png


def read_payload(path: Path) -> dict[str, object]:
    return json.loads(path.read_text(encoding="utf-8"))


def write_json(path: Path, payload: dict[str, object]) -> None:
    path.write_text(f"{json.dumps(payload, indent=2)}\n", encoding="utf-8")


def write_deck(path: Path, deck: Deck) -> None:
    payload = json.loads(deck.model_dump_json(by_alias=True))
    for slide in payload["slides"]:
        slide.pop("computed", None)
    path.write_text(f"{json.dumps(payload, indent=2)}\n", encoding="utf-8")
    write_computed_deck(str(path), deck)


def build_slide(slide_id: str, layout: str, slots: list[str], *, start_node: int = 1) -> Slide:
    return Slide(
        slide_id=slide_id,
        layout=layout,
        nodes=[
            Node(
                node_id=f"n-{start_node + offset}",
                slot_binding=slot_name,
                type="text",
                content=f"{slide_id}:{slot_name}",
            )
            for offset, slot_name in enumerate(slots)
        ],
    )


def make_empty_deck() -> Deck:
    return Deck(
        deck_id="deck-1",
        revision=0,
        theme="default",
        design_rules="default",
        slides=[],
        counters=Counters(),
    )


def make_bar_chart_data(*, values: list[float] | None = None) -> dict[str, object]:
    return {
        "categories": ["Q1", "Q2"] if values is None else [f"Q{index + 1}" for index in range(len(values))],
        "series": [
            {
                "name": "Revenue",
                "values": values or [12.0, 18.0],
            }
        ],
    }


def invoke_cli(args: list[str], *, input: str | None = None) -> Result:
    runner = CliRunner()
    return runner.invoke(cli, args, input=input)


def find_free_port() -> int:
    with contextlib.closing(socket.socket(socket.AF_INET, socket.SOCK_STREAM)) as sock:
        sock.bind(("127.0.0.1", 0))
        return int(sock.getsockname()[1])


def wait_for_http(url: str, *, timeout: float = 5.0) -> tuple[int, str]:
    deadline = time.monotonic() + timeout
    while time.monotonic() < deadline:
        try:
            with urllib.request.urlopen(url, timeout=0.25) as response:
                return response.status, response.read().decode("utf-8")
        except (ConnectionError, TimeoutError, urllib.error.URLError):
            time.sleep(0.05)

    raise AssertionError(f"Timed out waiting for {url}")


def make_computed_node(font_size_pt: float, *, overflow: bool = False) -> ComputedNode:
    return ComputedNode(
        x=72.0,
        y=54.0,
        width=576.0,
        height=80.0,
        font_size_pt=font_size_pt,
        font_family="Aptos",
        color="#333333",
        bg_color="#FFFFFF",
        font_bold=False,
        text_overflow=overflow,
        revision=1,
    )


def make_clean_deck() -> Deck:
    return Deck(
        deck_id="deck-clean",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="heading",
                        type="text",
                        content="Deck title",
                    )
                ],
                computed={"n-1": make_computed_node(28.0)},
            ),
            Slide(
                slide_id="s-2",
                layout="closing",
                nodes=[
                    Node(
                        node_id="n-2",
                        slot_binding="body",
                        type="text",
                        content="Thanks",
                    )
                ],
                computed={"n-2": make_computed_node(14.0)},
            ),
        ],
    )


def make_theme_switch_deck(theme_name: str = "default") -> Deck:
    return Deck(
        deck_id="deck-theme-switch",
        theme=theme_name,
        slides=[
            Slide(
                slide_id="s-1",
                layout="two_col",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="heading",
                        type="text",
                        content="Theme heading",
                    ),
                    Node(
                        node_id="n-2",
                        slot_binding="col1",
                        type="text",
                        content="Theme body left",
                    ),
                    Node(
                        node_id="n-3",
                        slot_binding="col2",
                        type="text",
                        content="Theme body right",
                    ),
                ],
            )
        ],
        counters=Counters(slides=1, nodes=3),
    )


def render_slide_signature(path: Path) -> tuple[tuple[int, int, int, str, str, str], ...]:
    presentation = Presentation(path)
    slide = presentation.slides[0]
    signature: list[tuple[int, int, int, str, str, str]] = []
    for shape in slide.shapes:
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        signature.append(
            (
                shape.left,
                shape.top,
                shape.width,
                run.font.name or "",
                str(run.font.color.rgb),
                str(shape.fill.fore_color.rgb),
            )
        )
    return tuple(signature)


def make_overflow_deck() -> Deck:
    rules = load_design_rules("default")
    return Deck(
        deck_id="deck-overflow",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="heading",
                        type="text",
                        content="Overflow heading",
                    )
                ],
                computed={
                    "n-1": make_computed_node(
                        float(rules.overflow_policy.min_font_size),
                        overflow=True,
                    )
                },
            ),
            Slide(
                slide_id="s-2",
                layout="closing",
                nodes=[
                    Node(
                        node_id="n-2",
                        slot_binding="body",
                        type="text",
                        content="Bye",
                    )
                ],
                computed={"n-2": make_computed_node(14.0)},
            ),
        ],
    )


def test_init_creates_valid_deck_file_and_reports_success_json(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    computed_path = computed_sidecar_path(deck_path)

    result = invoke_cli(["init", str(deck_path)])

    assert result.exit_code == 0
    response = json.loads(result.output)
    payload = read_payload(deck_path)

    assert response["ok"] is True
    assert response["data"] == {
        "deck_id": payload["deck_id"],
        "theme": "default",
        "design_rules": "default",
    }
    assert UUID(str(payload["deck_id"]))
    assert payload["revision"] == 0
    assert payload["slides"] == []
    assert payload["theme"] == "default"
    assert payload["design_rules"] == "default"
    assert payload["version"] == 2
    assert payload["_counters"] == {"slides": 0, "nodes": 0}
    assert computed_path.exists()
    assert read_computed_deck(str(deck_path)).slides == []


def test_init_applies_explicit_theme_and_rules_options(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"

    result = invoke_cli(["init", str(deck_path), "--theme", "default", "--rules", "default"])

    assert result.exit_code == 0
    payload = read_payload(deck_path)

    assert payload["theme"] == "default"
    assert payload["design_rules"] == "default"
    assert computed_sidecar_path(deck_path).exists()


def test_init_with_template_sets_relative_manifest_and_uses_manifest_layouts(tmp_path: Path) -> None:
    deck_path = tmp_path / "decks" / "deck.json"
    manifest_path = tmp_path / "templates" / "client-brand.manifest.json"
    template_path = tmp_path / "templates" / "client-brand.pptx"
    deck_path.parent.mkdir(parents=True)
    manifest_path.parent.mkdir(parents=True)
    template_path.write_bytes(b"pptx")
    write_json(
        manifest_path,
        {
            "source": "client-brand.pptx",
            "source_hash": "abc123",
            "theme": {
                "colors": {"primary": "#112233"},
                "fonts": {"heading": "Aptos Display", "body": "Aptos"},
                "spacing": {"base_unit": 10, "margin": 60, "gutter": 20},
            },
            "layouts": [
                {
                    "slug": "title_slide",
                    "usable": True,
                    "slot_mapping": {
                        "heading": {
                            "role": "heading",
                            "bounds": {"x": 72, "y": 72, "width": 560, "height": 96},
                        },
                        "subheading": {
                            "role": "body",
                            "bounds": {"x": 72, "y": 192, "width": 560, "height": 180},
                        },
                    },
                },
                {
                    "slug": "two_content",
                    "usable": True,
                    "slot_mapping": {
                        "heading": {
                            "role": "heading",
                            "bounds": {"x": 72, "y": 72, "width": 560, "height": 72},
                        },
                        "col1": {
                            "role": "body",
                            "bounds": {"x": 72, "y": 168, "width": 260, "height": 240},
                        },
                        "col2": {
                            "role": "body",
                            "bounds": {"x": 360, "y": 168, "width": 260, "height": 240},
                        },
                    },
                }
            ],
        },
    )

    result = invoke_cli(["init", str(deck_path), "--template", str(manifest_path)])

    assert result.exit_code == 0
    response = json.loads(result.output)
    payload = read_payload(deck_path)

    assert response["data"]["theme"] == "extracted-client-brand"
    assert response["data"]["template"] == "../templates/client-brand.manifest.json"
    assert payload["theme"] == "extracted-client-brand"
    assert payload["template_manifest"] == "../templates/client-brand.manifest.json"

    slide_add = invoke_cli(["slide", "add", str(deck_path), "--layout", "two_content"])

    assert slide_add.exit_code == 0
    updated = read_deck(str(deck_path))
    assert updated.slides[0].layout == "two_content"
    assert {node.slot_binding for node in updated.slides[0].nodes} == {"heading", "col1", "col2"}

    invalid_layout = invoke_cli(["slide", "add", str(deck_path), "--layout", "two_col"])

    assert invalid_layout.exit_code == 1
    assert json.loads(invalid_layout.stderr) == {
        "ok": False,
        "error": {
            "code": INVALID_LAYOUT,
            "message": "Unknown layout 'two_col'. Available layouts: title_slide, two_content",
        },
    }


def test_init_returns_error_when_theme_and_template_are_both_provided(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    manifest_path = tmp_path / "manifest.json"
    write_json(
        manifest_path,
        {
            "source": "client-brand.pptx",
            "source_hash": "abc123",
            "theme": {
                "colors": {"primary": "#112233"},
                "fonts": {"heading": "Aptos Display", "body": "Aptos"},
                "spacing": {"base_unit": 10, "margin": 60, "gutter": 20},
            },
            "layouts": [],
        },
    )

    result = invoke_cli(["init", str(deck_path), "--theme", "default", "--template", str(manifest_path)])

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": SCHEMA_ERROR,
            "message": "`--theme` and `--template` are mutually exclusive.",
        },
    }
    assert not deck_path.exists()


def test_init_returns_file_exists_error_when_target_exists(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck_path.write_text("{}", encoding="utf-8")

    result = invoke_cli(["init", str(deck_path)])

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": FILE_EXISTS,
            "message": f"Deck file already exists: {deck_path}",
        },
    }


def test_init_force_overwrites_existing_file(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck_path.write_text(
        json.dumps({"deck_id": "old", "revision": 9, "slides": ["stale"]}),
        encoding="utf-8",
    )

    result = invoke_cli(["init", str(deck_path), "--force"])

    assert result.exit_code == 0
    payload = read_payload(deck_path)

    assert payload["revision"] == 0
    assert payload["slides"] == []
    assert payload["theme"] == "default"
    assert payload["design_rules"] == "default"
    assert payload["deck_id"] != "old"
    assert computed_sidecar_path(deck_path).exists()


def test_init_returns_theme_validation_error_for_invalid_theme(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"

    result = invoke_cli(["init", str(deck_path), "--theme", "missing"])

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": THEME_NOT_FOUND,
            "message": "Theme 'missing' was not found.",
        },
    }
    assert not deck_path.exists()


def test_init_returns_rules_validation_error_for_invalid_rules(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"

    result = invoke_cli(["init", str(deck_path), "--rules", "missing"])

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": FILE_NOT_FOUND,
            "message": "Design rules profile 'missing' was not found.",
        },
    }
    assert not deck_path.exists()


def test_theme_list_returns_available_themes_as_json() -> None:
    result = invoke_cli(["theme", "list"])

    assert result.exit_code == 0
    assert result.stderr == ""
    assert json.loads(result.output) == {
        "ok": True,
        "data": {
            "themes": ["academic", "corporate", "dark", "default", "startup"],
        },
    }


def test_theme_apply_switches_theme_and_reflows_computed_layout(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_theme_switch_deck())

    result = invoke_cli(["theme", "apply", str(deck_path), "--theme", "startup"])

    assert result.exit_code == 0
    assert result.stderr == ""
    assert json.loads(result.output) == {
        "ok": True,
        "data": {
            "theme": "startup",
            "previous": "default",
        },
    }

    updated_deck = read_deck(str(deck_path))
    computed_deck = read_computed_deck(str(deck_path))
    heading = updated_deck.slides[0].computed["n-1"]
    left_column = updated_deck.slides[0].computed["n-2"]
    right_column = updated_deck.slides[0].computed["n-3"]

    assert updated_deck.theme == "startup"
    assert updated_deck.revision == 1
    assert computed_deck.revision == 1
    assert heading.x == 72.0
    assert heading.y == 72.0
    assert heading.font_family == "Helvetica"
    assert heading.color == "#0F172A"
    assert left_column.font_family == "Arial"
    assert left_column.bg_color == "#FFFDF8"
    assert right_column.x > left_column.x
    assert right_column.x == 388.0
    assert right_column.width == 288.0


def test_theme_apply_returns_error_for_invalid_theme(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_theme_switch_deck())

    result = invoke_cli(["theme", "apply", str(deck_path), "--theme", "missing"])

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": THEME_NOT_FOUND,
            "message": "Theme 'missing' was not found.",
        },
    }
    assert read_deck(str(deck_path)).theme == "default"


def test_suggest_layout_returns_ranked_suggestions_for_inline_content() -> None:
    result = invoke_cli(
        [
            "suggest-layout",
            "--content",
            json.dumps(
                {
                    "blocks": [
                        {"type": "heading", "text": "Quarterly update"},
                        {"type": "bullet", "text": "Revenue up 24%"},
                        {"type": "bullet", "text": "Margin improved"},
                        {"type": "bullet", "text": "Hiring plan"},
                        {"type": "bullet", "text": "Upsell motion"},
                        {"type": "bullet", "text": "Pricing refresh"},
                        {"type": "bullet", "text": "Channel rollout"},
                    ]
                }
            ),
        ]
    )

    assert result.exit_code == 0
    assert result.stderr == ""

    payload = json.loads(result.output)
    suggestions = payload["data"]["suggestions"]

    assert payload["ok"] is True
    assert suggestions[0]["layout"] == "two_col"
    assert suggestions[0]["score"] == 0.7
    assert [item["score"] for item in suggestions] == sorted(
        [item["score"] for item in suggestions],
        reverse=True,
    )


def test_suggest_layout_accepts_file_reference(tmp_path: Path) -> None:
    content_path = tmp_path / "content.json"
    write_json(
        content_path,
        {
            "blocks": [
                {"type": "heading", "text": "Launch plan"},
                {"type": "paragraph", "text": "Sequence the rollout in a single narrative arc."},
            ]
        },
    )

    result = invoke_cli(["suggest-layout", "--content", f"@{content_path}"])

    assert result.exit_code == 0
    payload = json.loads(result.output)

    assert payload["ok"] is True
    assert any(item["layout"] == "title_content" for item in payload["data"]["suggestions"])


def test_suggest_layout_image_count_surfaces_image_layouts() -> None:
    result = invoke_cli(
        [
            "suggest-layout",
            "--content",
            json.dumps(
                {
                    "blocks": [
                        {"type": "heading", "text": "Customer story"},
                        {"type": "paragraph", "text": "Show the product in use with supporting context."},
                    ]
                }
            ),
            "--image-count",
            "2",
        ]
    )

    assert result.exit_code == 0
    payload = json.loads(result.output)
    layouts = [item["layout"] for item in payload["data"]["suggestions"]]

    assert payload["ok"] is True
    assert set(layouts) <= {"gallery", "title", "title_content"}
    assert layouts[0] == "gallery"


def test_suggest_layout_returns_json_error_for_invalid_content() -> None:
    result = invoke_cli(["suggest-layout", "--content", "{bad json"])

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": SCHEMA_ERROR,
            "message": "Content JSON is not valid: `--content`",
        },
    }


def test_theme_apply_changes_subsequent_build_output(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    default_output = tmp_path / "default.pptx"
    corporate_output = tmp_path / "corporate.pptx"
    write_deck(deck_path, make_theme_switch_deck())

    initial_build = invoke_cli(["build", str(deck_path), "-o", str(default_output)])
    theme_result = invoke_cli(["theme", "apply", str(deck_path), "--theme", "corporate"])
    updated_build = invoke_cli(["build", str(deck_path), "-o", str(corporate_output)])

    assert initial_build.exit_code == 0
    assert theme_result.exit_code == 0
    assert updated_build.exit_code == 0
    assert render_slide_signature(default_output) != render_slide_signature(corporate_output)


def test_slot_set_updates_slot_text_using_aliases(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[build_slide("s-1", "title", ["heading", "subheading"], start_node=1)],
        counters=Counters(slides=1, nodes=2),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(
        ["slot", "set", str(deck_path), "--slide", "s-1", "--slot", "title", "--text", "New Title"]
    )
    payload = json.loads(result.stdout)
    updated = read_deck(str(deck_path))

    assert result.exit_code == 0
    assert payload == {
        "ok": True,
        "data": {
            "slide_id": "s-1",
            "slot": "heading",
            "node_id": "n-1",
            "type": "text",
            "text": "New Title",
            "content": {
                "blocks": [
                    {
                        "type": "paragraph",
                        "text": "New Title",
                        "level": 0,
                    }
                ]
            },
            "image_path": None,
            "image_fit": "contain",
            "font_size": None,
        },
    }
    assert updated.slides[0].nodes[0].content.to_plain_text() == "New Title"


def test_slot_set_accepts_structured_content_json(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[build_slide("s-1", "two_col", ["heading", "col1", "col2"], start_node=1)],
        counters=Counters(slides=1, nodes=3),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(
        [
            "slot",
            "set",
            str(deck_path),
            "--slide",
            "s-1",
            "--slot",
            "left",
            "--content",
            json.dumps(
                {
                    "blocks": [
                        {"type": "heading", "text": "Highlights"},
                        {"type": "bullet", "text": "Point one"},
                        {"type": "bullet", "text": "Point two", "level": 1},
                    ]
                }
            ),
        ]
    )
    payload = json.loads(result.stdout)
    updated = read_deck(str(deck_path))

    assert result.exit_code == 0
    assert payload == {
        "ok": True,
        "data": {
            "slide_id": "s-1",
            "slot": "col1",
            "node_id": "n-2",
            "type": "text",
            "text": "Highlights\nPoint one\nPoint two",
            "content": {
                "blocks": [
                    {"type": "heading", "text": "Highlights", "level": 0},
                    {"type": "bullet", "text": "Point one", "level": 0},
                    {"type": "bullet", "text": "Point two", "level": 1},
                ]
            },
            "image_path": None,
            "image_fit": "contain",
            "font_size": None,
        },
    }
    assert updated.slides[0].nodes[1].slot_binding == "col1"
    assert updated.slides[0].nodes[1].content.model_dump(mode="json") == {
        "blocks": [
            {"type": "heading", "text": "Highlights", "level": 0},
            {"type": "bullet", "text": "Point one", "level": 0},
            {"type": "bullet", "text": "Point two", "level": 1},
        ]
    }


@pytest.mark.parametrize(
    ("args",),
    [
        (
            [
                "--text",
                "Hello",
                "--content",
                json.dumps({"blocks": [{"type": "paragraph", "text": "Intro"}]}),
            ],
        ),
        (
            [
                "--text",
                "Hello",
                "--image",
                "photo.png",
            ],
        ),
        (
            [
                "--content",
                json.dumps({"blocks": [{"type": "paragraph", "text": "Intro"}]}),
                "--image",
                "photo.png",
            ],
        ),
    ],
)
def test_slot_set_rejects_mutually_exclusive_content_options(tmp_path: Path, args: list[str]) -> None:
    deck_path = tmp_path / "deck.json"
    write_png(tmp_path / "photo.png", width=12, height=12)
    write_deck(deck_path, make_empty_deck())

    result = invoke_cli(["slot", "set", str(deck_path), "--slide", "0", "--slot", "body", *args])

    assert result.exit_code == 1
    assert json.loads(result.stderr)["error"] == {
        "code": SCHEMA_ERROR,
        "message": "Options '--text', '--content', and '--image' are mutually exclusive; provide exactly one",
    }


def test_slot_set_supports_image_nodes(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    image_path = write_png(tmp_path / "photo.png", width=30, height=20)
    deck = Deck(
        deck_id="deck-1",
        slides=[build_slide("s-1", "title_content", ["heading", "body"], start_node=1)],
        counters=Counters(slides=1, nodes=2),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(
        ["slot", "set", str(deck_path), "--slide", "s-1", "--slot", "body", "--image", image_path.name]
    )
    payload = json.loads(result.stdout)
    updated = read_deck(str(deck_path))

    assert result.exit_code == 0
    assert payload == {
        "ok": True,
        "data": {
            "slide_id": "s-1",
            "slot": "body",
            "node_id": "n-2",
            "type": "image",
            "text": "",
            "content": {"blocks": []},
            "image_path": "photo.png",
            "image_fit": "contain",
            "font_size": None,
        },
    }
    body_node = next(node for node in updated.slides[0].nodes if node.slot_binding == "body")
    assert body_node.type == "image"
    assert body_node.image_path == "photo.png"


def test_slot_set_rejects_missing_image_path(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[build_slide("s-1", "title_content", ["heading", "body"], start_node=1)],
        counters=Counters(slides=1, nodes=2),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(
        ["slot", "set", str(deck_path), "--slide", "0", "--slot", "body", "--image", "missing.png"]
    )

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": FILE_NOT_FOUND,
            "message": f"Image file not found: {tmp_path / 'missing.png'}",
        },
    }


def test_slot_clear_removes_bound_nodes(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[build_slide("s-1", "two_col", ["heading", "col1", "col2"], start_node=1)],
        counters=Counters(slides=1, nodes=3),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(["slot", "clear", str(deck_path), "--slide", "0", "--slot", "right"])
    payload = json.loads(result.stdout)
    updated = read_deck(str(deck_path))

    assert result.exit_code == 0
    assert payload == {
        "ok": True,
        "data": {
            "slide_id": "s-1",
            "slot": "col2",
            "removed_node_ids": ["n-3"],
        },
    }
    assert [node.slot_binding for node in updated.slides[0].nodes] == ["heading", "col1"]


def test_slot_bind_rebinds_existing_node_to_valid_slot(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[
            Slide(
                slide_id="s-1",
                layout="title",
                nodes=[Node(node_id="n-1", slot_binding=None, type="text", content="Loose title")],
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(["slot", "bind", str(deck_path), "--node", "n-1", "--slot", "title"])
    payload = json.loads(result.stdout)
    updated = read_deck(str(deck_path))

    assert result.exit_code == 0
    assert payload == {
        "ok": True,
        "data": {
            "slide_id": "s-1",
            "slot": "heading",
            "node_id": "n-1",
        },
    }
    assert updated.slides[0].nodes[0].slot_binding == "heading"


def test_chart_add_creates_chart_node_with_inline_data_and_title(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[build_slide("s-1", "title_content", ["heading", "body"], start_node=1)],
        counters=Counters(slides=1, nodes=2),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(
        [
            "chart",
            "add",
            str(deck_path),
            "--slide",
            "s-1",
            "--slot",
            "body",
            "--type",
            "bar",
            "--data",
            json.dumps(make_bar_chart_data()),
            "--title",
            "Revenue by Quarter",
        ]
    )
    payload = json.loads(result.stdout)
    updated = read_deck(str(deck_path))
    chart_node = updated.slides[0].nodes[1]

    assert result.exit_code == 0
    assert payload == {
        "ok": True,
        "data": {
            "slide_id": "s-1",
            "slot": "body",
            "node_id": "n-2",
            "chart_type": "bar",
        },
    }
    assert chart_node.type == "chart"
    assert chart_node.chart_spec is not None
    assert chart_node.chart_spec.title == "Revenue by Quarter"
    assert chart_node.chart_spec.categories == ["Q1", "Q2"]
    assert chart_node.chart_spec.series is not None
    assert chart_node.chart_spec.series[0].values == [12.0, 18.0]


def test_chart_add_supports_data_file(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    data_path = tmp_path / "chart-data.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[build_slide("s-1", "title_content", ["heading", "body"], start_node=1)],
        counters=Counters(slides=1, nodes=2),
    )
    write_deck(deck_path, deck)
    write_json(data_path, make_bar_chart_data(values=[5.0, 9.0]))

    result = invoke_cli(
        [
            "chart",
            "add",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "body",
            "--type",
            "column",
            "--data-file",
            str(data_path),
        ]
    )

    assert result.exit_code == 0
    payload = json.loads(result.stdout)
    chart_node = read_deck(str(deck_path)).slides[0].nodes[1]
    assert payload["data"]["chart_type"] == "column"
    assert chart_node.chart_spec is not None
    assert chart_node.chart_spec.chart_type == "column"
    assert chart_node.chart_spec.series is not None
    assert chart_node.chart_spec.series[0].values == [5.0, 9.0]


def test_chart_update_updates_existing_chart_data(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[
            Slide(
                slide_id="s-1",
                layout="closing",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="body",
                        type="chart",
                        chart_spec={
                            "chart_type": "bar",
                            "title": "Revenue",
                            **make_bar_chart_data(values=[10.0, 15.0]),
                        },
                    )
                ],
            )
        ],
        counters=Counters(slides=1, nodes=1),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(
        [
            "chart",
            "update",
            str(deck_path),
            "--node",
            "n-1",
            "--data",
            json.dumps(make_bar_chart_data(values=[20.0, 24.0, 30.0])),
        ]
    )
    payload = json.loads(result.stdout)
    updated_node = read_deck(str(deck_path)).slides[0].nodes[0]

    assert result.exit_code == 0
    assert payload == {
        "ok": True,
        "data": {
            "node_id": "n-1",
            "chart_type": "bar",
            "updated": True,
        },
    }
    assert updated_node.chart_spec is not None
    assert updated_node.chart_spec.title == "Revenue"
    assert updated_node.chart_spec.categories == ["Q1", "Q2", "Q3"]
    assert updated_node.chart_spec.series is not None
    assert updated_node.chart_spec.series[0].values == [20.0, 24.0, 30.0]


def test_chart_update_rejects_non_chart_nodes(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[Slide(slide_id="s-1", layout="closing", nodes=[Node(node_id="n-1", slot_binding="body", type="text")])],
        counters=Counters(slides=1, nodes=1),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(
        [
            "chart",
            "update",
            str(deck_path),
            "--node",
            "n-1",
            "--data",
            json.dumps(make_bar_chart_data(values=[4.0, 8.0])),
        ]
    )

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": INVALID_NODE_TYPE,
            "message": "Node 'n-1' is not a chart node",
            "node_id": "n-1",
            "expected_type": "chart",
            "actual_type": "text",
        },
    }


def test_chart_add_invalid_chart_type_returns_valid_types(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(
        deck_path,
        Deck(
            deck_id="deck-1",
            slides=[build_slide("s-1", "title_content", ["heading", "body"], start_node=1)],
            counters=Counters(slides=1, nodes=2),
        ),
    )

    result = invoke_cli(
        [
            "chart",
            "add",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "body",
            "--type",
            "radar",
            "--data",
            json.dumps(make_bar_chart_data()),
        ]
    )
    payload = json.loads(result.stderr)

    assert result.exit_code == 1
    assert payload["error"]["code"] == INVALID_CHART_TYPE
    assert payload["error"]["chart_type"] == "radar"
    assert payload["error"]["valid_types"] == ["bar", "column", "line", "pie", "scatter", "area", "doughnut"]


def test_chart_add_invalid_chart_data_returns_json_error(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(
        deck_path,
        Deck(
            deck_id="deck-1",
            slides=[build_slide("s-1", "title_content", ["heading", "body"], start_node=1)],
            counters=Counters(slides=1, nodes=2),
        ),
    )

    invalid_data = {
        "categories": ["Q1", "Q2"],
        "series": [{"name": "Revenue", "values": [10.0]}],
    }
    result = invoke_cli(
        [
            "chart",
            "add",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "body",
            "--type",
            "bar",
            "--data",
            json.dumps(invalid_data),
        ]
    )
    payload = json.loads(result.stderr)

    assert result.exit_code == 1
    assert payload["error"]["code"] == CHART_DATA_ERROR
    assert "series 'Revenue' has 1 values for 2 categories" in payload["error"]["message"]
    assert payload["error"]["validation_errors"] == [
        {
            "ctx": {
                "category_count": 2,
                "series_name": "Revenue",
                "value_count": 1,
            },
            "input": {
                "categories": ["Q1", "Q2"],
                "chart_type": "bar",
                "series": [{"name": "Revenue", "values": [10.0]}],
            },
            "loc": [],
            "msg": "series 'Revenue' has 1 values for 2 categories",
            "type": CHART_DATA_ERROR,
        }
    ]


def test_slide_add_creates_layout_bound_nodes(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    init_deck(str(deck_path), theme="default", design_rules="default", force=False)

    result = invoke_cli(["slide", "add", str(deck_path), "--layout", "title"])
    payload = json.loads(result.stdout)
    deck = read_deck(str(deck_path))
    layout = get_layout("title")

    assert result.exit_code == 0
    assert payload == {
        "ok": True,
        "data": {
            "slide_index": 0,
            "slide_id": "s-1",
            "layout": "title",
        },
    }
    assert [node.slot_binding for node in deck.slides[0].nodes] == list(layout.slots)
    assert [node.node_id for node in deck.slides[0].nodes] == ["n-1", "n-2"]


def test_slide_add_invalid_layout_returns_json_error_with_available_layouts(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    init_deck(str(deck_path), theme="default", design_rules="default", force=False)

    result = invoke_cli(["slide", "add", str(deck_path), "--layout", "nope"])
    payload = json.loads(result.stderr)

    assert result.exit_code == 1
    assert payload["ok"] is False
    assert payload["error"]["code"] == INVALID_LAYOUT
    for layout_name in list_layouts():
        assert layout_name in payload["error"]["message"]


def test_slide_add_auto_layout_selects_layout_and_prefills_slots(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    init_deck(str(deck_path), theme="default", design_rules="default", force=False)

    content = json.dumps(
        {
            "blocks": [
                {"type": "heading", "text": "Highlights"},
                {"type": "paragraph", "text": "Left column summary"},
                {"type": "paragraph", "text": "Right column summary"},
            ]
        }
    )

    result = invoke_cli(["slide", "add", str(deck_path), "--auto-layout", "--content", content])
    payload = json.loads(result.stdout)
    deck = read_deck(str(deck_path))

    assert result.exit_code == 0
    assert payload == {
        "ok": True,
        "data": {
            "slide_index": 0,
            "slide_id": "s-1",
            "layout": "two_col",
            "auto_selected": True,
            "reason": "Two balanced content blocks",
        },
    }
    assert deck.slides[0].layout == "two_col"
    assert [
        (node.slot_binding, node.content.model_dump(mode="json"))
        for node in deck.slides[0].nodes
        if node.type == "text"
    ] == [
        ("heading", {"blocks": [{"type": "heading", "text": "Highlights", "level": 0}]}),
        ("col1", {"blocks": [{"type": "paragraph", "text": "Left column summary", "level": 0}]}),
        ("col2", {"blocks": [{"type": "paragraph", "text": "Right column summary", "level": 0}]}),
    ]


def test_slide_add_auto_layout_rejects_layout_override(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    init_deck(str(deck_path), theme="default", design_rules="default", force=False)

    result = invoke_cli(
        [
            "slide",
            "add",
            str(deck_path),
            "--layout",
            "title",
            "--auto-layout",
            "--content",
            json.dumps({"blocks": [{"type": "heading", "text": "Conflicting"}]}),
        ]
    )
    payload = json.loads(result.stderr)

    assert result.exit_code == 1
    assert payload == {
        "ok": False,
        "error": {
            "code": SCHEMA_ERROR,
            "message": "`--auto-layout` and `--layout` are mutually exclusive.",
        },
    }


def test_slide_remove_by_index_shifts_remaining_slides(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[
            build_slide("s-1", "title", ["heading", "subheading"], start_node=1),
            build_slide("s-2", "title", ["heading", "subheading"], start_node=3),
        ],
        counters=Counters(slides=2, nodes=4),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(["slide", "remove", str(deck_path), "--slide", "0"])
    payload = json.loads(result.stdout)
    updated = read_deck(str(deck_path))

    assert result.exit_code == 0
    assert payload == {
        "ok": True,
        "data": {
            "removed": "s-1",
            "slide_count": 1,
        },
    }
    assert [slide.slide_id for slide in updated.slides] == ["s-2"]
    assert updated.get_slide(0).slide_id == "s-2"


def test_slide_remove_last_slide_leaves_empty_deck(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[build_slide("s-1", "title", ["heading", "subheading"], start_node=1)],
        counters=Counters(slides=1, nodes=2),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(["slide", "remove", str(deck_path), "--slide", "0"])
    payload = json.loads(result.stdout)
    updated = read_deck(str(deck_path))

    assert result.exit_code == 0
    assert payload["data"] == {"removed": "s-1", "slide_count": 0}
    assert updated.slides == []


def test_slide_remove_invalid_ref_returns_invalid_slide(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(deck_id="deck-1")
    write_deck(deck_path, deck)

    result = invoke_cli(["slide", "remove", str(deck_path), "--slide", "9"])
    payload = json.loads(result.stderr)

    assert result.exit_code == 1
    assert payload["error"]["code"] == INVALID_SLIDE


def test_slide_remove_accepts_slide_id_reference(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = Deck(
        deck_id="deck-1",
        slides=[
            build_slide("s-1", "title", ["heading", "subheading"], start_node=1),
            build_slide("s-2", "title", ["heading", "subheading"], start_node=3),
        ],
        counters=Counters(slides=2, nodes=4),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(["slide", "remove", str(deck_path), "--slide", "s-1"])
    payload = json.loads(result.stdout)
    updated = read_deck(str(deck_path))

    assert result.exit_code == 0
    assert payload["data"] == {"removed": "s-1", "slide_count": 1}
    assert [slide.slide_id for slide in updated.slides] == ["s-2"]


def test_slide_set_layout_rebinds_content_and_warns_on_unbound_nodes(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    slide = build_slide("s-1", "three_col", ["heading", "col1", "col2", "col3"], start_node=1)
    deck = Deck(
        deck_id="deck-1",
        slides=[slide],
        counters=Counters(slides=1, nodes=4),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(
        ["slide", "set-layout", str(deck_path), "--slide", "s-1", "--layout", "two_col"]
    )
    payload = json.loads(result.stdout)
    warning = json.loads(result.stderr)
    updated = read_deck(str(deck_path))

    assert result.exit_code == 0
    assert payload == {
        "ok": True,
        "data": {
            "slide_id": "s-1",
            "layout": "two_col",
            "unbound_nodes": ["n-4"],
        },
    }
    assert warning == {
        "ok": True,
        "warning": {
            "code": UNBOUND_NODES,
            "message": "1 node(s) became unbound during slot rebinding.",
        },
        "data": {
            "slide_id": "s-1",
            "unbound_nodes": ["n-4"],
        },
    }
    assert updated.slides[0].layout == "two_col"
    assert [node.slot_binding for node in updated.slides[0].nodes] == [
        "heading",
        "col1",
        "col2",
        None,
    ]


def test_slide_set_layout_to_same_layout_has_no_unbound_nodes(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    slide = build_slide("s-1", "two_col", ["heading", "col1", "col2"], start_node=1)
    deck = Deck(
        deck_id="deck-1",
        slides=[slide],
        counters=Counters(slides=1, nodes=3),
    )
    write_deck(deck_path, deck)

    result = invoke_cli(
        ["slide", "set-layout", str(deck_path), "--slide", "0", "--layout", "two_col"]
    )
    payload = json.loads(result.stdout)
    updated = read_deck(str(deck_path))

    assert result.exit_code == 0
    assert result.stderr == ""
    assert payload == {
        "ok": True,
        "data": {
            "slide_id": "s-1",
            "layout": "two_col",
            "unbound_nodes": [],
        },
    }
    assert [node.slot_binding for node in updated.slides[0].nodes] == list(get_layout("two_col").slots)
    assert len(updated.slides[0].nodes) == 3


def test_slide_set_layout_invalid_slide_returns_invalid_slide(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    init_deck(str(deck_path), theme="default", design_rules="default", force=False)

    result = invoke_cli(
        ["slide", "set-layout", str(deck_path), "--slide", "s-999", "--layout", "title"]
    )
    payload = json.loads(result.stderr)

    assert result.exit_code == 1
    assert payload["error"]["code"] == INVALID_SLIDE


def test_slide_set_layout_invalid_layout_returns_invalid_layout(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    init_deck(str(deck_path), theme="default", design_rules="default", force=False)

    result = invoke_cli(["slide", "set-layout", str(deck_path), "--slide", "0", "--layout", "bad"])
    payload = json.loads(result.stderr)

    assert result.exit_code == 1
    assert payload["error"]["code"] == INVALID_LAYOUT

def test_batch_applies_multiple_operations_atomically(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_empty_deck())

    result = invoke_cli(
        ["batch", str(deck_path)],
        input=json.dumps(
            [
                {"command": "slide_add", "args": {"layout": "title"}},
                {"command": "slot_set", "args": {"slide": 0, "slot": "title", "text": "Hello"}},
                {
                    "command": "slot_set",
                    "args": {"slide": 0, "slot": "subtitle", "text": "World", "font_size": 18},
                },
                {"command": "slide_add", "args": {"layout": "two_col"}},
                {"command": "slot_set", "args": {"slide": 1, "slot": "title", "text": "Key Points"}},
            ]
        ),
    )

    assert result.exit_code == 0
    payload = json.loads(result.output)
    assert payload["ok"] is True
    assert payload["data"]["operations"] == 5
    assert payload["data"]["results"][0] == {
        "slide_index": 0,
        "slide_id": "s-1",
        "layout": "title",
    }

    deck = read_deck(str(deck_path))
    assert deck.revision == 1
    assert [slide.layout for slide in deck.slides] == ["title", "two_col"]
    assert [(node.slot_binding, node.content.to_plain_text()) for node in deck.slides[0].nodes] == [
        ("heading", "Hello"),
        ("subheading", "World"),
    ]
    assert deck.slides[0].nodes[1].style_overrides["font_size"] == 18.0
    assert [(node.slot_binding, node.content.to_plain_text()) for node in deck.slides[1].nodes] == [
        ("heading", "Key Points"),
        ("col1", ""),
        ("col2", ""),
    ]


def test_batch_rolls_back_and_reports_operation_index_on_failure(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_empty_deck())
    original_payload = deck_path.read_text(encoding="utf-8")

    result = invoke_cli(
        ["batch", str(deck_path)],
        input=json.dumps(
            [
                {"command": "slide_add", "args": {"layout": "title"}},
                {"command": "slot_set", "args": {"slide": 0, "slot": "body", "text": "Nope"}},
            ]
        ),
    )

    assert result.exit_code == 1
    error = json.loads(result.stderr)["error"]
    assert error["code"] == INVALID_SLOT
    assert error["operation_index"] == 1
    assert deck_path.read_text(encoding="utf-8") == original_payload
    assert read_deck(str(deck_path)).slides == []


def test_batch_slot_set_accepts_structured_content_objects(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_empty_deck())

    result = invoke_cli(
        ["batch", str(deck_path)],
        input=json.dumps(
            [
                {"command": "slide_add", "args": {"layout": "two_col"}},
                {
                    "command": "slot_set",
                    "args": {
                        "slide": 0,
                        "slot": "left",
                        "content": {
                            "blocks": [
                                {"type": "heading", "text": "Highlights"},
                                {"type": "bullet", "text": "Point one"},
                                {"type": "bullet", "text": "Point two", "level": 1},
                            ]
                        },
                    },
                },
            ]
        ),
    )

    assert result.exit_code == 0
    payload = json.loads(result.output)
    assert payload["data"]["results"][1]["content"] == {
        "blocks": [
            {"type": "heading", "text": "Highlights", "level": 0},
            {"type": "bullet", "text": "Point one", "level": 0},
            {"type": "bullet", "text": "Point two", "level": 1},
        ]
    }

    deck = read_deck(str(deck_path))
    assert deck.slides[0].nodes[1].slot_binding == "col1"
    assert deck.slides[0].nodes[1].content.model_dump(mode="json") == {
        "blocks": [
            {"type": "heading", "text": "Highlights", "level": 0},
            {"type": "bullet", "text": "Point one", "level": 0},
            {"type": "bullet", "text": "Point two", "level": 1},
        ]
    }


def test_batch_slide_add_auto_layout_selects_layout_and_prefills_slots(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_empty_deck())

    result = invoke_cli(
        ["batch", str(deck_path)],
        input=json.dumps(
            [
                {
                    "command": "slide_add",
                    "args": {
                        "auto_layout": True,
                        "content": {
                            "blocks": [
                                {"type": "heading", "text": "Highlights"},
                                {"type": "paragraph", "text": "Left column summary"},
                                {"type": "paragraph", "text": "Right column summary"},
                            ]
                        },
                    },
                }
            ]
        ),
    )

    assert result.exit_code == 0
    payload = json.loads(result.output)
    assert payload == {
        "ok": True,
        "data": {
            "operations": 1,
            "results": [
                {
                    "slide_index": 0,
                    "slide_id": "s-1",
                    "layout": "two_col",
                    "auto_selected": True,
                    "reason": "Two balanced content blocks",
                }
            ],
        },
    }

    deck = read_deck(str(deck_path))
    assert deck.slides[0].layout == "two_col"
    assert [
        (node.slot_binding, node.content.to_plain_text())
        for node in deck.slides[0].nodes
        if node.type == "text"
    ] == [
        ("heading", "Highlights"),
        ("col1", "Left column summary"),
        ("col2", "Right column summary"),
    ]


def test_batch_empty_array_is_a_successful_no_op(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_empty_deck())

    result = invoke_cli(["batch", str(deck_path)], input="[]")

    assert result.exit_code == 0
    assert json.loads(result.output) == {
        "ok": True,
        "data": {"operations": 0, "results": []},
    }
    assert read_deck(str(deck_path)).revision == 0


def test_batch_invalid_json_returns_schema_error(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_empty_deck())

    result = invoke_cli(["batch", str(deck_path)], input="{not json")

    assert result.exit_code == 1
    error = json.loads(result.stderr)["error"]
    assert error["code"] == SCHEMA_ERROR
    assert "Invalid JSON batch payload" in error["message"]


def test_batch_supports_chart_add_and_chart_update(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_empty_deck())

    result = invoke_cli(
        ["batch", str(deck_path)],
        input=json.dumps(
            [
                {"command": "slide_add", "args": {"layout": "title_content"}},
                {
                    "command": "chart_add",
                    "args": {
                        "slide": 0,
                        "slot": "body",
                        "type": "bar",
                        "data": make_bar_chart_data(values=[8.0, 13.0]),
                    },
                },
                {
                    "command": "chart_update",
                    "args": {
                        "node": "n-2",
                        "data": make_bar_chart_data(values=[8.0, 13.0, 21.0]),
                    },
                },
            ]
        ),
    )

    assert result.exit_code == 0
    payload = json.loads(result.output)
    assert payload["data"]["operations"] == 3
    assert payload["data"]["results"][1] == {
        "slide_id": "s-1",
        "slot": "body",
        "node_id": "n-2",
        "chart_type": "bar",
    }
    assert payload["data"]["results"][2] == {
        "node_id": "n-2",
        "chart_type": "bar",
        "updated": True,
    }

    chart_node = read_deck(str(deck_path)).slides[0].nodes[1]
    assert chart_node.type == "chart"
    assert chart_node.chart_spec is not None
    assert chart_node.chart_spec.categories == ["Q1", "Q2", "Q3"]
    assert chart_node.chart_spec.series is not None
    assert chart_node.chart_spec.series[0].values == [8.0, 13.0, 21.0]


def test_batch_supports_all_mutation_types(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(
        deck_path,
        Deck(
            deck_id="deck-1",
            revision=0,
            theme="default",
            design_rules="default",
            slides=[
                Slide(
                    slide_id="s-1",
                    layout="title",
                    nodes=[Node(node_id="n-1", slot_binding=None, type="text", content="Loose title")],
                )
            ],
            counters=Counters(slides=1, nodes=1),
        ),
    )

    result = invoke_cli(
        ["batch", str(deck_path)],
        input=json.dumps(
            [
                {"command": "slot_bind", "args": {"node": "n-1", "slot": "title"}},
                {"command": "slot_set", "args": {"slide": "s-1", "slot": "subtitle", "text": "Intro"}},
                {"command": "slide_add", "args": {"layout": "two_col"}},
                {"command": "slot_set", "args": {"slide": 1, "slot": "left", "text": "Key Points"}},
                {"command": "slot_clear", "args": {"slide": 1, "slot": "left"}},
                {"command": "slide_set_layout", "args": {"slide": 1, "layout": "three_col"}},
                {"command": "slide_remove", "args": {"slide": 1}},
            ]
        ),
    )

    assert result.exit_code == 0
    payload = json.loads(result.output)
    assert payload["data"]["operations"] == 7
    assert payload["data"]["results"][0] == {
        "slide_id": "s-1",
        "slot": "heading",
        "node_id": "n-1",
    }
    assert payload["data"]["results"][5] == {
        "slide_id": "s-2",
        "layout": "three_col",
        "unbound_nodes": [],
    }
    assert payload["data"]["results"][6] == {
        "removed": "s-2",
        "slide_count": 1,
    }

    deck = read_deck(str(deck_path))
    assert len(deck.slides) == 1
    assert deck.slides[0].nodes[0].node_id == "n-1"
    assert deck.slides[0].nodes[0].slot_binding == "heading"
    assert deck.slides[0].nodes[0].content.to_plain_text() == "Loose title"
    assert deck.slides[0].nodes[1].slot_binding == "subheading"
    assert deck.slides[0].nodes[1].content.to_plain_text() == "Intro"


def test_batch_uses_single_mutate_deck_call(monkeypatch) -> None:
    calls: list[str] = []

    def fake_mutate_deck(path: str, fn):
        calls.append(path)
        results = fn(Deck(deck_id="deck-1"), BuiltinLayoutProvider())
        return Deck(deck_id="deck-1"), results

    monkeypatch.setattr("agent_slides.commands.batch.mutate_deck", fake_mutate_deck)

    result = invoke_cli(
        ["batch", "deck.json"],
        input=json.dumps([{"command": "slide_add", "args": {"layout": "title"}}]),
    )

    assert result.exit_code == 0
    assert calls == ["deck.json"]
    assert json.loads(result.output)["data"]["operations"] == 1


def test_build_command_creates_valid_pptx_and_reports_slide_count(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    output_path = tmp_path / "presentation.pptx"
    write_deck(deck_path, make_clean_deck())

    result = invoke_cli(["build", str(deck_path), "-o", str(output_path)])

    assert result.exit_code == 0
    assert result.stderr == ""
    assert json.loads(result.output) == {
        "ok": True,
        "data": {
            "output": str(output_path),
            "slides": 2,
        },
    }
    presentation = Presentation(output_path)
    assert len(presentation.slides) == 2
    assert read_computed_deck(str(deck_path)).revision == make_clean_deck().revision


def test_build_command_reports_missing_deck_as_file_not_found(tmp_path: Path) -> None:
    result = invoke_cli(
        ["build", str(tmp_path / "missing.json"), "-o", str(tmp_path / "presentation.pptx")]
    )

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": FILE_NOT_FOUND,
            "message": f"Deck file not found: {tmp_path / 'missing.json'}",
        },
    }


def test_validate_command_reports_clean_deck(tmp_path: Path) -> None:
    deck_path = tmp_path / "clean.json"
    write_deck(deck_path, make_clean_deck())

    result = invoke_cli(["validate", str(deck_path)])

    assert result.exit_code == 0
    assert result.stderr == ""
    assert json.loads(result.output) == {
        "ok": True,
        "data": {
            "warnings": [],
            "clean": True,
        },
    }


def test_validate_command_returns_structured_overflow_warning(tmp_path: Path) -> None:
    deck_path = tmp_path / "overflow.json"
    write_deck(deck_path, make_overflow_deck())

    result = invoke_cli(["validate", str(deck_path)])

    assert result.exit_code == 0
    payload = json.loads(result.output)
    assert payload["ok"] is True
    assert payload["data"]["clean"] is False
    assert any(warning["code"] == "OVERFLOW" for warning in payload["data"]["warnings"])
    for warning in payload["data"]["warnings"]:
        assert {"code", "severity", "message"} <= warning.keys()


def test_inspect_command_summarizes_manifest_json(tmp_path: Path) -> None:
    manifest_path = tmp_path / "client-brand.manifest.json"
    write_json(
        manifest_path,
        {
            "source": "client-brand.pptx",
            "source_hash": "abc123",
            "slide_masters": [
                {
                    "layouts": [
                        {
                            "name": "Title Slide",
                            "slug": "title_slide",
                            "slot_mapping": {"heading": 0, "subheading": 1},
                        },
                        {
                            "name": "Two Content",
                            "slug": "two_content",
                            "slot_mapping": {"heading": 0, "col1": 1, "col2": 2},
                        },
                    ]
                },
                {
                    "layouts": [
                        {
                            "name": "Blank",
                            "slug": "blank",
                            "slot_mapping": {},
                        },
                        {
                            "name": "Custom Layout 4",
                            "slug": "custom_layout_4",
                            "slot_mapping": {},
                        },
                    ]
                },
            ],
            "theme": {
                "colors": {"primary": "#112233"},
                "fonts": {"heading": "Aptos Display", "body": "Aptos"},
                "spacing": {"base_unit": 10, "margin": 60, "gutter": 20},
            },
        },
    )

    result = invoke_cli(["inspect", str(manifest_path)])

    assert result.exit_code == 0
    assert result.stderr == ""
    assert json.loads(result.output) == {
        "ok": True,
        "data": {
            "source": "client-brand.pptx",
            "layouts_found": 4,
            "usable_layouts": 3,
            "theme_extracted": True,
            "layouts": [
                {
                    "name": "Title Slide",
                    "slug": "title_slide",
                    "slots": ["heading", "subheading"],
                    "usable": True,
                },
                {
                    "name": "Two Content",
                    "slug": "two_content",
                    "slots": ["heading", "col1", "col2"],
                    "usable": True,
                },
                {
                    "name": "Blank",
                    "slug": "blank",
                    "slots": [],
                    "usable": True,
                },
                {
                    "name": "Custom Layout 4",
                    "slug": "custom_layout_4",
                    "slots": [],
                    "usable": False,
                    "reason": "no typed placeholders",
                },
            ],
        },
    }


def test_inspect_command_reports_theme_not_extracted_when_missing(tmp_path: Path) -> None:
    manifest_path = tmp_path / "manifest.json"
    write_json(
        manifest_path,
        {
            "source": "template.pptx",
            "slide_masters": [
                {
                    "layouts": [
                        {
                            "name": "One Column",
                            "slug": "one_column",
                            "slot_mapping": {"body": 3},
                        }
                    ]
                }
            ],
        },
    )

    result = invoke_cli(["inspect", str(manifest_path)])

    assert result.exit_code == 0
    assert json.loads(result.output) == {
        "ok": True,
        "data": {
            "source": "template.pptx",
            "layouts_found": 1,
            "usable_layouts": 1,
            "theme_extracted": False,
            "layouts": [
                {
                    "name": "One Column",
                    "slug": "one_column",
                    "slots": ["body"],
                    "usable": True,
                }
            ],
        },
    }


def test_inspect_command_reports_missing_manifest_as_file_not_found(tmp_path: Path) -> None:
    missing_path = tmp_path / "missing.manifest.json"

    result = invoke_cli(["inspect", str(missing_path)])

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": FILE_NOT_FOUND,
            "message": f"Manifest file not found: {missing_path}",
        },
    }


def test_inspect_command_reports_invalid_json_as_schema_error(tmp_path: Path) -> None:
    manifest_path = tmp_path / "broken.manifest.json"
    manifest_path.write_text("{not valid json\n", encoding="utf-8")

    result = invoke_cli(["inspect", str(manifest_path)])

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": SCHEMA_ERROR,
            "message": f"Manifest file is not valid JSON: {manifest_path}",
        },
    }


def test_info_command_dumps_indented_deck_json(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    deck = make_clean_deck()
    write_deck(deck_path, deck)

    result = invoke_cli(["info", str(deck_path)])

    assert result.exit_code == 0
    assert result.stderr == ""
    assert result.output == f"{deck.model_dump_json(by_alias=True, indent=2)}\n"
    assert result.output.startswith("{\n  ")
    assert "computed" not in read_payload(deck_path)["slides"][0]


def test_info_command_reports_missing_deck_as_file_not_found(tmp_path: Path) -> None:
    result = invoke_cli(["info", str(tmp_path / "missing.json")])

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": FILE_NOT_FOUND,
            "message": f"Deck file not found: {tmp_path / 'missing.json'}",
        },
    }


def test_preview_command_starts_server_opens_browser_and_stops_cleanly(
    tmp_path: Path,
    monkeypatch,
) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_clean_deck())
    stop_event = threading.Event()
    browser_calls: list[str] = []
    result_box: dict[str, Result] = {}
    runner = CliRunner()
    port = find_free_port()

    monkeypatch.setattr(
        "agent_slides.commands.preview._wait_for_shutdown",
        lambda: stop_event.wait(timeout=5),
    )
    monkeypatch.setattr(
        "agent_slides.commands.preview.webbrowser.open",
        lambda url: browser_calls.append(url),
    )

    def invoke() -> None:
        result_box["result"] = runner.invoke(cli, ["preview", str(deck_path), "--port", str(port)])

    thread = threading.Thread(target=invoke, daemon=True)
    thread.start()

    status, body = wait_for_http(f"http://127.0.0.1:{port}/api/deck")
    assert status == 200
    assert json.loads(body)["deck_id"] == "deck-clean"
    assert thread.is_alive()

    stop_event.set()
    thread.join(timeout=5)

    assert not thread.is_alive()
    result = result_box["result"]
    lines = [line for line in result.output.strip().splitlines() if line]

    assert result.exit_code == 0
    assert result.stderr == ""
    assert browser_calls == [f"http://localhost:{port}"]
    assert json.loads(lines[0]) == {
        "ok": True,
        "data": {
            "url": f"http://localhost:{port}",
            "watching": "deck.json",
        },
    }
    assert json.loads(lines[1]) == {"ok": True, "data": {"stopped": True}}


def test_preview_command_supports_custom_port_and_no_open(
    tmp_path: Path,
    monkeypatch,
) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_clean_deck())
    stop_event = threading.Event()
    browser_calls: list[str] = []
    result_box: dict[str, Result] = {}
    runner = CliRunner()
    port = find_free_port()

    monkeypatch.setattr(
        "agent_slides.commands.preview._wait_for_shutdown",
        lambda: stop_event.wait(timeout=5),
    )
    monkeypatch.setattr(
        "agent_slides.commands.preview.webbrowser.open",
        lambda url: browser_calls.append(url),
    )

    def invoke() -> None:
        result_box["result"] = runner.invoke(
            cli,
            ["preview", str(deck_path), "--port", str(port), "--no-open"],
        )

    thread = threading.Thread(target=invoke, daemon=True)
    thread.start()

    status, body = wait_for_http(f"http://127.0.0.1:{port}/")
    assert status == 200
    assert "<title>agent-slides preview</title>" in body

    stop_event.set()
    thread.join(timeout=5)

    assert not thread.is_alive()
    result = result_box["result"]
    lines = [line for line in result.output.strip().splitlines() if line]

    assert result.exit_code == 0
    assert browser_calls == []
    assert json.loads(lines[0])["data"] == {
        "url": f"http://localhost:{port}",
        "watching": "deck.json",
    }
    assert json.loads(lines[1]) == {"ok": True, "data": {"stopped": True}}


def test_preview_command_supports_background_mode(
    tmp_path: Path,
    monkeypatch,
) -> None:
    deck_path = tmp_path / "deck.json"
    write_deck(deck_path, make_clean_deck())
    browser_calls: list[str] = []
    port = find_free_port()

    monkeypatch.setattr(
        "agent_slides.commands.preview.webbrowser.open",
        lambda url: browser_calls.append(url),
    )

    result = invoke_cli(["preview", str(deck_path), "--port", str(port), "--background"])

    assert result.exit_code == 0
    assert result.stderr == ""
    assert browser_calls == [f"http://localhost:{port}"]

    payload = json.loads(result.output.strip())
    assert payload["ok"] is True
    assert payload["data"]["url"] == f"http://localhost:{port}"
    assert isinstance(payload["data"]["pid"], int)

    status, body = wait_for_http(f"http://127.0.0.1:{port}/api/deck")
    assert status == 200
    assert json.loads(body)["deck_id"] == "deck-clean"

    os.kill(int(payload["data"]["pid"]), signal.SIGTERM)
    os.waitpid(int(payload["data"]["pid"]), 0)


def test_preview_command_reports_missing_deck_as_file_not_found(tmp_path: Path) -> None:
    missing_path = tmp_path / "missing.json"

    result = invoke_cli(["preview", str(missing_path), "--no-open"])

    assert result.exit_code == 1
    assert json.loads(result.stderr) == {
        "ok": False,
        "error": {
            "code": FILE_NOT_FOUND,
            "message": f"Deck file not found: {missing_path}",
        },
    }


def test_cli_help_does_not_list_chat_command() -> None:
    result = invoke_cli(["--help"])

    assert result.exit_code == 0
    assert " chat " not in result.output


def test_chat_command_is_not_registered() -> None:
    result = invoke_cli(["chat"])

    assert result.exit_code != 0
    assert "No such command 'chat'" in result.output
