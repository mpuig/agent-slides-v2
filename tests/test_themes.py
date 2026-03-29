from __future__ import annotations

from importlib import resources
from pathlib import Path

import pytest
from pptx import Presentation

from agent_slides.engine.reflow import reflow_deck
from agent_slides.io.pptx_writer import write_pptx
from agent_slides.errors import AgentSlidesError, THEME_NOT_FOUND
from agent_slides.model.themes import list_themes, load_theme, resolve_style
from agent_slides.model.types import Counters, Deck, Node, Slide, Theme

EXPECTED_THEMES = {
    "academic": {
        "heading_font": "Georgia",
        "body_font": "Times New Roman",
        "margin": 48,
        "gutter": 16,
    },
    "corporate": {
        "heading_font": "Georgia",
        "body_font": "Arial",
        "margin": 64,
        "gutter": 20,
    },
    "dark": {
        "heading_font": "Arial",
        "body_font": "Calibri",
        "margin": 60,
        "gutter": 24,
    },
    "default": {
        "heading_font": "Calibri",
        "body_font": "Calibri",
        "margin": 60,
        "gutter": 20,
    },
    "startup": {
        "heading_font": "Helvetica",
        "body_font": "Arial",
        "margin": 72,
        "gutter": 28,
    },
}


def _channel_luminance(value: str) -> float:
    channel = int(value, 16) / 255.0
    if channel <= 0.03928:
        return channel / 12.92
    return ((channel + 0.055) / 1.055) ** 2.4


def contrast_ratio(foreground: str, background: str) -> float:
    fg = foreground.removeprefix("#")
    bg = background.removeprefix("#")
    fg_luminance = (
        0.2126 * _channel_luminance(fg[0:2])
        + 0.7152 * _channel_luminance(fg[2:4])
        + 0.0722 * _channel_luminance(fg[4:6])
    )
    bg_luminance = (
        0.2126 * _channel_luminance(bg[0:2])
        + 0.7152 * _channel_luminance(bg[2:4])
        + 0.0722 * _channel_luminance(bg[4:6])
    )
    lighter = max(fg_luminance, bg_luminance)
    darker = min(fg_luminance, bg_luminance)
    return (lighter + 0.05) / (darker + 0.05)


def make_theme_deck(theme_name: str) -> Deck:
    return Deck(
        deck_id=f"deck-{theme_name}",
        theme=theme_name,
        slides=[
            Slide(
                slide_id="s-1",
                layout="two_col",
                nodes=[
                    Node(
                        node_id="n-1",
                        slot_binding="heading",
                        type="text",
                        content="Theme heading",
                    ),
                    Node(
                        node_id="n-2",
                        slot_binding="col1",
                        type="text",
                        content="Left column copy",
                    ),
                    Node(
                        node_id="n-3",
                        slot_binding="col2",
                        type="text",
                        content="Right column copy",
                    ),
                ],
                computed={},
            )
        ],
        counters=Counters(slides=1, nodes=3),
    )


def render_signature(path: Path) -> tuple[tuple[int, int, int, str, str, str], ...]:
    presentation = Presentation(str(path))
    slide = presentation.slides[0]
    signature: list[tuple[int, int, int, str, str, str]] = []
    for shape in slide.shapes:
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        signature.append(
            (
                shape.left,
                shape.top,
                shape.width,
                run.font.name or "",
                str(run.font.color.rgb),
                str(shape.fill.fore_color.rgb),
            )
        )
    return tuple(signature)


@pytest.mark.parametrize("theme_name", sorted(EXPECTED_THEMES))
def test_load_theme_returns_expected_theme(theme_name: str) -> None:
    theme = load_theme(theme_name)
    expected = EXPECTED_THEMES[theme_name]

    assert isinstance(theme, Theme)
    assert theme.name == theme_name
    assert theme.fonts.heading == expected["heading_font"]
    assert theme.fonts.body == expected["body_font"]
    assert theme.spacing.margin == expected["margin"]
    assert theme.spacing.gutter == expected["gutter"]


def test_load_theme_raises_for_missing_theme() -> None:
    with pytest.raises(AgentSlidesError) as exc_info:
        load_theme("nonexistent")

    assert exc_info.value.code == THEME_NOT_FOUND


def test_list_themes_returns_all_builtins() -> None:
    assert list_themes() == sorted(EXPECTED_THEMES)


def test_resolve_style_heading() -> None:
    theme = load_theme("default")

    assert resolve_style(theme, "heading") == {
        "font_family": "Calibri",
        "color": "#1a1a2e",
        "font_bold": True,
    }


def test_resolve_style_body() -> None:
    theme = load_theme("default")

    assert resolve_style(theme, "body") == {
        "font_family": "Calibri",
        "color": "#333333",
        "font_bold": False,
    }


def test_theme_resource_is_packaged() -> None:
    for theme_name in EXPECTED_THEMES:
        resource = resources.files("agent_slides.themes").joinpath(f"{theme_name}.json")

        assert resource.is_file()
        assert f'"name": "{theme_name}"' in resource.read_text(encoding="utf-8")


@pytest.mark.parametrize("theme_name", sorted(EXPECTED_THEMES))
def test_theme_text_colors_meet_wcag_aa(theme_name: str) -> None:
    theme = load_theme(theme_name)

    assert contrast_ratio(theme.colors.text, theme.colors.background) >= 4.5
    assert (
        contrast_ratio(theme.colors.heading_text or "", theme.colors.background) >= 4.5
    )
    assert (
        contrast_ratio(theme.colors.subtle_text or "", theme.colors.background) >= 4.5
    )


def test_reflow_uses_theme_spacing() -> None:
    default_deck = make_theme_deck("default")
    startup_deck = make_theme_deck("startup")

    reflow_deck(default_deck)
    reflow_deck(startup_deck)

    default_heading = default_deck.slides[0].computed["n-1"]
    startup_heading = startup_deck.slides[0].computed["n-1"]
    default_left_column = default_deck.slides[0].computed["n-2"]
    startup_left_column = startup_deck.slides[0].computed["n-2"]

    assert default_heading.x == 60
    assert default_heading.y == 60
    assert startup_heading.x == 72
    assert startup_heading.y == 72
    assert startup_left_column.x > default_left_column.x
    assert startup_left_column.width < default_left_column.width


def test_each_theme_builds_distinct_pptx_output(tmp_path: Path) -> None:
    signatures = {}
    for theme_name in sorted(EXPECTED_THEMES):
        deck = make_theme_deck(theme_name)
        reflow_deck(deck)
        output_path = tmp_path / f"{theme_name}.pptx"
        write_pptx(deck, str(output_path))
        signatures[theme_name] = render_signature(output_path)

    assert len(set(signatures.values())) == len(EXPECTED_THEMES)
