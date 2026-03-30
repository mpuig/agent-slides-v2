from __future__ import annotations

import importlib.util
import json
import subprocess
import sys
from pathlib import Path

from click.testing import CliRunner
from pptx import Presentation

from agent_slides.cli import cli
from agent_slides.io import read_template_manifest
from tests.test_e2e_template import create_test_template

ROOT = Path(__file__).resolve().parents[1]
CERT_SCRIPT_PATH = ROOT / "scripts" / "generate_layout_cert_suite.py"
INVENTORY_SCRIPT_PATH = ROOT / "scripts" / "export_layout_inventory.py"
FIXTURES_DIR = ROOT / "fixtures" / "layout_cases"


def _load_script_module(path: Path, name: str):
    spec = importlib.util.spec_from_file_location(name, path)
    assert spec is not None
    assert spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _run_cert_script(*args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [sys.executable, str(CERT_SCRIPT_PATH), *args],
        check=False,
        capture_output=True,
        text=True,
    )


def _write_inventory_for_template(
    template_path: Path, manifest_path: Path, inventory_path: Path
) -> dict[str, object]:
    inventory_module = _load_script_module(
        INVENTORY_SCRIPT_PATH, "export_layout_inventory"
    )
    manifest = read_template_manifest(template_path, manifest_path).manifest
    inventory = inventory_module.build_inventory(manifest)
    inventory_path.write_text(
        json.dumps(inventory, indent=2, sort_keys=True) + "\n", encoding="utf-8"
    )
    return inventory


def _expected_deck_count(inventory: dict[str, object]) -> int:
    fixture_variants = {
        path.stem: json.loads(path.read_text(encoding="utf-8"))
        for path in sorted(FIXTURES_DIR.glob("*.json"))
    }
    return sum(
        len(fixture_variants[str(layout["slot_structure"])])
        for layout in inventory["layouts"]
        if bool(layout.get("testable", False))
    )


def test_generate_layout_cert_suite_cli_writes_expected_decks(tmp_path: Path) -> None:
    template_path = tmp_path / "brand-template.pptx"
    manifest_path = tmp_path / "brand-template.manifest.json"
    inventory_path = tmp_path / "layout-inventory.json"
    output_dir = tmp_path / "cert-suite"
    create_test_template(template_path)
    inventory = _write_inventory_for_template(
        template_path, manifest_path, inventory_path
    )

    result = _run_cert_script(
        "--manifest",
        str(manifest_path),
        "--inventory",
        str(inventory_path),
        "--fixtures",
        str(FIXTURES_DIR),
        "--output",
        str(output_dir),
    )

    assert result.returncode == 0, result.stderr
    summary = json.loads(result.stdout)
    assert summary["template"] == "brand-template"
    assert summary["deck_count"] == _expected_deck_count(inventory)

    deck_paths = sorted(output_dir.glob("**/deck.json"))
    assert len(deck_paths) == summary["deck_count"]

    first_deck = json.loads(deck_paths[0].read_text(encoding="utf-8"))
    assert first_deck["template_manifest"]
    assert len(first_deck["slides"]) == 1

    for deck_path in deck_paths:
        deck = json.loads(deck_path.read_text(encoding="utf-8"))
        slide = deck["slides"][0]
        variant_name = deck_path.parent.name
        layout_slug = deck_path.parent.parent.name
        assert slide["layout"] == layout_slug
        assert len(deck["slides"]) == 1
        assert len({node["node_id"] for node in slide["nodes"]}) == len(slide["nodes"])

        manifest_ref = (deck_path.parent / deck["template_manifest"]).resolve()
        assert manifest_ref == manifest_path.resolve()

        if variant_name == "image_present":
            image_nodes = [node for node in slide["nodes"] if node["type"] == "image"]
            assert image_nodes
            for node in image_nodes:
                image_path = (deck_path.parent / node["image_path"]).resolve()
                assert image_path.is_file()


def test_generate_layout_cert_suite_is_deterministic(tmp_path: Path) -> None:
    template_path = tmp_path / "brand-template.pptx"
    manifest_path = tmp_path / "brand-template.manifest.json"
    inventory_path = tmp_path / "layout-inventory.json"
    first_output_dir = tmp_path / "first"
    second_output_dir = tmp_path / "second"
    create_test_template(template_path)
    _write_inventory_for_template(template_path, manifest_path, inventory_path)

    first = _run_cert_script(
        "--manifest",
        str(manifest_path),
        "--inventory",
        str(inventory_path),
        "--fixtures",
        str(FIXTURES_DIR),
        "--output",
        str(first_output_dir),
    )
    second = _run_cert_script(
        "--manifest",
        str(manifest_path),
        "--inventory",
        str(inventory_path),
        "--fixtures",
        str(FIXTURES_DIR),
        "--output",
        str(second_output_dir),
    )

    assert first.returncode == 0, first.stderr
    assert second.returncode == 0, second.stderr

    first_payloads = {
        path.relative_to(first_output_dir).as_posix(): path.read_text(encoding="utf-8")
        for path in sorted(first_output_dir.glob("**/deck.json"))
    }
    second_payloads = {
        path.relative_to(second_output_dir).as_posix(): path.read_text(encoding="utf-8")
        for path in sorted(second_output_dir.glob("**/deck.json"))
    }
    assert first_payloads == second_payloads


def test_generated_cert_suite_builds_for_multiple_templates_and_counts_all_templates(
    tmp_path: Path,
) -> None:
    runner = CliRunner()
    total_expected = 0
    all_deck_paths: dict[str, list[Path]] = {}

    for index in range(4):
        template_path = tmp_path / f"example-{index + 1}.pptx"
        manifest_path = tmp_path / f"example-{index + 1}.manifest.json"
        inventory_path = tmp_path / f"example-{index + 1}.inventory.json"
        output_dir = tmp_path / "cert-suite"
        Presentation().save(template_path)
        inventory = _write_inventory_for_template(
            template_path, manifest_path, inventory_path
        )
        total_expected += _expected_deck_count(inventory)

        result = _run_cert_script(
            "--manifest",
            str(manifest_path),
            "--inventory",
            str(inventory_path),
            "--fixtures",
            str(FIXTURES_DIR),
            "--output",
            str(output_dir),
        )

        assert result.returncode == 0, result.stderr
        template_slug = f"example-{index + 1}"
        all_deck_paths[template_slug] = sorted(
            (output_dir / template_slug).glob("**/deck.json")
        )
        assert all_deck_paths[template_slug]

    all_generated_decks = sorted((tmp_path / "cert-suite").glob("**/deck.json"))
    assert len(all_generated_decks) == total_expected

    for template_slug in ("example-1", "example-2"):
        for deck_path in all_deck_paths[template_slug]:
            output_path = deck_path.with_suffix(".pptx")
            result = runner.invoke(
                cli, ["build", str(deck_path), "-o", str(output_path)]
            )
            assert result.exit_code == 0, result.output
            presentation = Presentation(str(output_path))
            assert len(presentation.slides) == 1
