from __future__ import annotations

import json
from pathlib import Path
from xml.etree import ElementTree as ET
from zipfile import ZipFile

from click.testing import CliRunner
from pptx import Presentation

from agent_slides.cli import cli
from agent_slides.errors import UNBOUND_NODES
from tests.image_helpers import write_png

DRAWING_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def collect_slide_text(path: Path) -> list[str]:
    presentation = Presentation(str(path))
    text_values: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_values.append(shape.text_frame.text)
    return text_values


def read_slide_xml(path: Path, *, index: int = 1) -> ET.Element:
    with ZipFile(path) as archive:
        return ET.fromstring(archive.read(f"ppt/slides/slide{index}.xml"))


def invoke(
    runner: CliRunner, args: list[str], *, input_text: str | None = None
) -> tuple[int, dict[str, object], str]:
    result = runner.invoke(cli, args, input=input_text)
    payload = json.loads(result.stdout) if result.stdout else {}
    return result.exit_code, payload, result.stderr


def test_full_demo_flow(tmp_path: Path) -> None:
    runner = CliRunner()
    deck_path = tmp_path / "deck.json"
    pptx_path = tmp_path / "output.pptx"

    exit_code, payload, _ = invoke(
        runner, ["init", str(deck_path), "--theme", "default"]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    for layout in ["title", "two_col", "quote"]:
        exit_code, payload, _ = invoke(
            runner, ["slide", "add", str(deck_path), "--layout", layout]
        )
        assert exit_code == 0
        assert payload["ok"] is True

    slot_sets = [
        ("0", "title", "My Presentation"),
        ("0", "subtitle", "March 2026"),
        ("1", "title", "Key Points"),
        ("1", "col1", "First point"),
        ("1", "col2", "Second point"),
        ("2", "quote", "The best way to predict the future is to invent it."),
        ("2", "attribution", "Alan Kay"),
        ("1", "title", "Three Key Points"),
    ]
    for slide_ref, slot_name, text in slot_sets:
        exit_code, payload, _ = invoke(
            runner,
            [
                "slot",
                "set",
                str(deck_path),
                "--slide",
                slide_ref,
                "--slot",
                slot_name,
                "--text",
                text,
            ],
        )
        assert exit_code == 0
        assert payload["ok"] is True

    exit_code, payload, _ = invoke(
        runner, ["build", str(deck_path), "-o", str(pptx_path)]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    presentation = Presentation(str(pptx_path))
    assert len(presentation.slides) == 3

    all_text = collect_slide_text(pptx_path)
    assert "My Presentation" in all_text
    assert "March 2026" in all_text
    assert "Three Key Points" in all_text
    assert "First point" in all_text
    assert "Second point" in all_text
    assert "The best way to predict the future is to invent it." in all_text
    assert "Alan Kay" in all_text


def test_build_renders_auto_detected_text_bullets_as_native_pptx_bullets(
    tmp_path: Path,
) -> None:
    runner = CliRunner()
    deck_path = tmp_path / "deck.json"
    pptx_path = tmp_path / "output.pptx"

    exit_code, payload, _ = invoke(
        runner, ["init", str(deck_path), "--theme", "default"]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    exit_code, payload, _ = invoke(
        runner, ["slide", "add", str(deck_path), "--layout", "title_content"]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    exit_code, payload, _ = invoke(
        runner,
        [
            "slot",
            "set",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "body",
            "--text",
            "Key points:\n- First item\n* Second item",
        ],
    )
    assert exit_code == 0
    assert payload["ok"] is True

    exit_code, payload, _ = invoke(
        runner, ["build", str(deck_path), "-o", str(pptx_path)]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    presentation = Presentation(str(pptx_path))
    text_frame = next(
        shape.text_frame
        for shape in presentation.slides[0].shapes
        if shape.has_text_frame
        and shape.text_frame.text == "Key points:\nFirst item\nSecond item"
    )
    body_shape = next(
        shape
        for shape in read_slide_xml(pptx_path).findall(".//p:sp", DRAWING_NS)
        if [text.text for text in shape.findall(".//a:t", DRAWING_NS)]
        == ["Key points:", "First item", "Second item"]
    )
    paragraphs = body_shape.findall(".//a:p", DRAWING_NS)

    assert [paragraph.text for paragraph in text_frame.paragraphs] == [
        "Key points:",
        "First item",
        "Second item",
    ]
    assert paragraphs[0].find("./a:pPr/a:buChar", DRAWING_NS) is None
    assert paragraphs[1].find("./a:pPr/a:buChar", DRAWING_NS) is not None
    assert paragraphs[2].find("./a:pPr/a:buChar", DRAWING_NS) is not None


def test_build_succeeds_after_slot_set_normalizes_absolute_image_paths(
    tmp_path: Path,
) -> None:
    runner = CliRunner()
    deck_path = tmp_path / "deck.json"
    pptx_path = tmp_path / "output.pptx"
    image_dir = tmp_path / "assets"
    image_dir.mkdir()
    image_path = write_png(image_dir / "photo.png", width=30, height=20)

    exit_code, payload, _ = invoke(
        runner, ["init", str(deck_path), "--theme", "default"]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    exit_code, payload, _ = invoke(
        runner, ["slide", "add", str(deck_path), "--layout", "image_right"]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    exit_code, payload, _ = invoke(
        runner,
        [
            "slot",
            "set",
            str(deck_path),
            "--slide",
            "0",
            "--slot",
            "image",
            "--image",
            str(image_path),
        ],
    )
    assert exit_code == 0
    assert payload["data"]["image_path"] == "assets/photo.png"

    exit_code, payload, _ = invoke(
        runner, ["build", str(deck_path), "-o", str(pptx_path)]
    )
    assert exit_code == 0
    assert payload["ok"] is True
    assert pptx_path.is_file()


def test_layout_switch_content_migration(tmp_path: Path) -> None:
    runner = CliRunner()
    deck_path = tmp_path / "deck.json"

    exit_code, payload, _ = invoke(runner, ["init", str(deck_path)])
    assert exit_code == 0
    assert payload["ok"] is True

    exit_code, payload, _ = invoke(
        runner, ["slide", "add", str(deck_path), "--layout", "three_col"]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    for slot_name, text in [
        ("col1", "Column 1"),
        ("col2", "Column 2"),
        ("col3", "Column 3"),
    ]:
        exit_code, payload, _ = invoke(
            runner,
            [
                "slot",
                "set",
                str(deck_path),
                "--slide",
                "0",
                "--slot",
                slot_name,
                "--text",
                text,
            ],
        )
        assert exit_code == 0
        assert payload["ok"] is True

    result = runner.invoke(
        cli,
        ["slide", "set-layout", str(deck_path), "--slide", "0", "--layout", "two_col"],
    )
    output = json.loads(result.stdout)
    warning = json.loads(result.stderr)

    assert result.exit_code == 0
    assert output["ok"] is True
    assert warning["warning"]["code"] == UNBOUND_NODES

    info_result = runner.invoke(cli, ["info", str(deck_path)])
    deck_data = json.loads(info_result.output)

    slide = deck_data["slides"][0]
    assert slide["layout"] == "two_col"

    bound_slots = {
        node["slot_binding"] for node in slide["nodes"] if node["slot_binding"]
    }
    assert "col1" in bound_slots
    assert "col2" in bound_slots

    unbound = [node for node in slide["nodes"] if node["slot_binding"] is None]
    assert len(unbound) >= 1
    assert any(
        node["content"]["blocks"]
        == [{"type": "paragraph", "text": "Column 3", "level": 0}]
        for node in unbound
    )


def test_batch_creates_full_deck(tmp_path: Path) -> None:
    runner = CliRunner()
    deck_path = tmp_path / "deck.json"
    pptx_path = tmp_path / "batch.pptx"

    exit_code, payload, _ = invoke(runner, ["init", str(deck_path)])
    assert exit_code == 0
    assert payload["ok"] is True

    batch_input = json.dumps(
        [
            {"command": "slide_add", "args": {"layout": "title"}},
            {
                "command": "slot_set",
                "args": {"slide": 0, "slot": "title", "text": "Batch Created"},
            },
            {"command": "slide_add", "args": {"layout": "quote"}},
            {
                "command": "slot_set",
                "args": {
                    "slide": 1,
                    "slot": "quote",
                    "text": "Efficiency is doing things right.",
                },
            },
        ]
    )

    exit_code, payload, _ = invoke(
        runner, ["batch", str(deck_path)], input_text=batch_input
    )
    assert exit_code == 0
    assert payload["ok"] is True

    exit_code, payload, _ = invoke(
        runner, ["build", str(deck_path), "-o", str(pptx_path)]
    )
    assert exit_code == 0
    assert payload["ok"] is True

    presentation = Presentation(str(pptx_path))
    assert len(presentation.slides) == 2
