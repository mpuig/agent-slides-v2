"""Template-backed layout registry loaded from a learned manifest."""

from __future__ import annotations

import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation

from agent_slides.errors import AgentSlidesError, INVALID_LAYOUT, SCHEMA_ERROR
from agent_slides.model.layouts import DEFAULT_TEXT_FITTING
from agent_slides.model.themes import load_theme
from agent_slides.model.types import GridDef, LayoutDef, SlotDef, TextFitting, Theme, ThemeColors, ThemeFonts, ThemeSpacing

_TEMPLATE_GRID = GridDef(
    columns=1,
    rows=1,
    row_heights=[1.0],
    col_widths=[1.0],
    margin=0.0,
    gutter=0.0,
)
_THEME_SLUG_PATTERN = re.compile(r"[^a-z0-9]+")
_ALIGNMENT_EPSILON_PT = 8.0
_PEER_Y_EPSILON_PT = 18.0
_PEER_HEIGHT_RATIO = 0.1
_CONTENT_TYPES = {"text", "image", "chart", "table"}


@dataclass(frozen=True)
class _PlaceholderStyle:
    preferred_font: float | None = None
    text_align: str | None = None
    vertical_align: str | None = None


def _as_dict(value: object, *, context: str) -> dict[str, Any]:
    if not isinstance(value, dict):
        raise AgentSlidesError(SCHEMA_ERROR, f"{context} must be an object")
    return dict(value)


def _require_string(value: object, *, context: str) -> str:
    if not isinstance(value, str) or not value.strip():
        raise AgentSlidesError(SCHEMA_ERROR, f"{context} must be a non-empty string")
    return value.strip()


def _optional_number(mapping: dict[str, Any], *keys: str) -> float | None:
    for key in keys:
        value = mapping.get(key)
        if value is None:
            continue
        if isinstance(value, bool) or not isinstance(value, int | float):
            raise AgentSlidesError(SCHEMA_ERROR, f"Manifest field '{key}' must be numeric")
        return float(value)
    return None


def _optional_int(mapping: dict[str, Any], *keys: str) -> int | None:
    for key in keys:
        value = mapping.get(key)
        if value is None:
            continue
        if isinstance(value, bool) or not isinstance(value, int):
            raise AgentSlidesError(SCHEMA_ERROR, f"Manifest field '{key}' must be an integer")
        return value
    return None


def _optional_string(mapping: dict[str, Any], *keys: str) -> str | None:
    for key in keys:
        value = mapping.get(key)
        if value is None:
            continue
        if not isinstance(value, str):
            raise AgentSlidesError(SCHEMA_ERROR, f"Manifest field '{key}' must be a string")
        normalized = value.strip()
        return normalized or None
    return None


def _optional_string_list(mapping: dict[str, Any], *keys: str) -> list[str] | None:
    for key in keys:
        value = mapping.get(key)
        if value is None:
            continue
        if not isinstance(value, list):
            raise AgentSlidesError(SCHEMA_ERROR, f"Manifest field '{key}' must be an array")
        items: list[str] = []
        for index, item in enumerate(value):
            if not isinstance(item, str):
                raise AgentSlidesError(SCHEMA_ERROR, f"Manifest field '{key}[{index}]' must be a string")
            normalized = item.strip().lower()
            if normalized not in _CONTENT_TYPES:
                raise AgentSlidesError(
                    SCHEMA_ERROR,
                    f"Manifest field '{key}[{index}]' must be one of {sorted(_CONTENT_TYPES)}",
                )
            items.append(normalized)
        return items
    return None


def _infer_role(slot_name: str, slot_mapping: dict[str, Any]) -> str:
    explicit = slot_mapping.get("role")
    if isinstance(explicit, str) and explicit:
        return explicit.lower()

    placeholder_type = slot_mapping.get("type")
    if isinstance(placeholder_type, str):
        normalized = placeholder_type.upper()
        if normalized == "PICTURE":
            return "image"

    lowered = slot_name.lower()
    if lowered in {"heading", "title", "header"}:
        return "heading"
    if lowered in {"subheading", "subtitle"}:
        return "body"
    if lowered in {"quote"}:
        return "quote"
    if lowered in {"attribution", "credit", "citation"}:
        return "attribution"
    if "image" in lowered or lowered.startswith("img"):
        return "image"
    return "body"


def _coerce_slot_mapping(
    slot_name: str,
    raw_slot: object,
    *,
    placeholders_by_idx: dict[int, dict[str, Any]],
) -> dict[str, Any]:
    if isinstance(raw_slot, bool):
        raise AgentSlidesError(SCHEMA_ERROR, f"slot_mapping[{slot_name!r}] must not be a boolean")
    if isinstance(raw_slot, int):
        try:
            placeholder = placeholders_by_idx[raw_slot]
        except KeyError as exc:
            raise AgentSlidesError(
                SCHEMA_ERROR,
                f"slot_mapping[{slot_name!r}] references missing placeholder idx {raw_slot}",
            ) from exc
        return dict(placeholder)
    return _as_dict(raw_slot, context=f"slot_mapping[{slot_name!r}]")


def _build_slot(
    slot_name: str,
    raw_slot: object,
    *,
    placeholders_by_idx: dict[int, dict[str, Any]],
    placeholder_styles: dict[tuple[int, int, int], _PlaceholderStyle],
    layout_ref: tuple[int, int],
) -> SlotDef:
    slot_mapping = _coerce_slot_mapping(
        slot_name,
        raw_slot,
        placeholders_by_idx=placeholders_by_idx,
    )
    raw_bounds = slot_mapping.get("bounds", slot_mapping)
    bounds = _as_dict(raw_bounds, context=f"slot_mapping[{slot_name!r}] bounds")
    role = _infer_role(slot_name, slot_mapping)
    raw_idx = slot_mapping.get("idx")
    placeholder_style = (
        placeholder_styles.get((layout_ref[0], layout_ref[1], raw_idx))
        if isinstance(raw_idx, int) and not isinstance(raw_idx, bool)
        else None
    )
    text_align = _optional_string(slot_mapping, "text_align")
    vertical_align = _optional_string(slot_mapping, "vertical_align")

    return SlotDef(
        grid_row=1,
        grid_col=1,
        role=role,
        peer_group=_optional_string(slot_mapping, "peer_group"),
        alignment_group=_optional_string(slot_mapping, "alignment_group"),
        reading_order=_optional_int(slot_mapping, "reading_order") or 0,
        size_policy=_optional_string(slot_mapping, "size_policy") or _default_size_policy(role),
        allowed_content=_optional_string_list(slot_mapping, "allowed_content") or ["text", "image", "chart", "table"],
        min_font=_optional_number(slot_mapping, "min_font"),
        max_font=_optional_number(slot_mapping, "max_font"),
        preferred_font=(
            _optional_number(slot_mapping, "preferred_font")
            if _optional_number(slot_mapping, "preferred_font") is not None
            else (placeholder_style.preferred_font if placeholder_style is not None else None)
        ),
        text_align=text_align or (placeholder_style.text_align if placeholder_style is not None else None) or "left",
        vertical_align=(
            vertical_align
            or (placeholder_style.vertical_align if placeholder_style is not None else None)
            or "top"
        ),
        padding=_optional_number(slot_mapping, "padding") if _optional_number(slot_mapping, "padding") is not None else 0.0,
        x=_optional_number(bounds, "x", "left"),
        y=_optional_number(bounds, "y", "top"),
        width=_optional_number(bounds, "width", "w"),
        height=_optional_number(bounds, "height", "h"),
        bg_color=slot_mapping.get("bg_color"),
        bg_transparency=float(slot_mapping.get("bg_transparency", 0.0)),
        full_bleed=bool(slot_mapping.get("full_bleed", False)),
        height_mode=str(slot_mapping.get("height_mode", "fixed")),
        width_mode=str(slot_mapping.get("width_mode", "fixed")),
    )


def _default_text_fitting(role: str) -> TextFitting:
    if role == "heading":
        return DEFAULT_TEXT_FITTING["heading"]
    return DEFAULT_TEXT_FITTING["body"]


def _default_size_policy(role: str) -> str:
    if role in {"heading", "attribution"}:
        return "fit_content"
    if role in {"body", "image", "quote"}:
        return "fill_remaining"
    return "fixed"


def _normalize_text_align(value: object) -> str | None:
    name = getattr(value, "name", None)
    if isinstance(name, str):
        normalized = name.casefold()
        if normalized == "ctr":
            return "center"
        if normalized == "just":
            return "justify"
        return normalized
    return None


def _normalize_vertical_align(value: object) -> str | None:
    name = getattr(value, "name", None)
    if isinstance(name, str):
        normalized = name.casefold()
        if normalized in {"mid", "middle", "ctr"}:
            return "middle"
        return normalized
    return None


def _extract_placeholder_style(placeholder: object) -> _PlaceholderStyle:
    try:
        text_frame = placeholder.text_frame
    except Exception:
        return _PlaceholderStyle()

    paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None
    preferred_font: float | None = None
    text_align = _normalize_text_align(getattr(paragraph, "alignment", None))
    vertical_align = _normalize_vertical_align(getattr(text_frame, "vertical_anchor", None))

    font_candidates: list[object] = []
    if paragraph is not None:
        font_candidates.append(paragraph.font)
        if paragraph.runs:
            font_candidates.append(paragraph.runs[0].font)

    for font in font_candidates:
        size = getattr(getattr(font, "size", None), "pt", None)
        if isinstance(size, int | float):
            preferred_font = float(size)
            break

    return _PlaceholderStyle(
        preferred_font=preferred_font,
        text_align=text_align,
        vertical_align=vertical_align,
    )


def _load_placeholder_styles(source_path: str) -> dict[tuple[int, int, int], _PlaceholderStyle]:
    try:
        presentation = Presentation(source_path)
    except Exception:
        return {}

    styles: dict[tuple[int, int, int], _PlaceholderStyle] = {}
    for master_index, slide_master in enumerate(presentation.slide_masters):
        for layout_index, slide_layout in enumerate(slide_master.slide_layouts):
            for placeholder in slide_layout.placeholders:
                key = (master_index, layout_index, int(placeholder.placeholder_format.idx))
                styles[key] = _extract_placeholder_style(placeholder)
    return styles


def _coerce_layouts(raw_layouts: object) -> list[tuple[str, dict[str, Any]]]:
    if isinstance(raw_layouts, dict):
        return [(str(slug), _as_dict(value, context=f"layout[{slug!r}]")) for slug, value in raw_layouts.items()]
    if isinstance(raw_layouts, list):
        items: list[tuple[str, dict[str, Any]]] = []
        for index, value in enumerate(raw_layouts):
            layout = _as_dict(value, context=f"layouts[{index}]")
            slug = _require_string(layout.get("slug") or layout.get("name"), context=f"layouts[{index}].slug")
            items.append((slug, layout))
        return items
    raise AgentSlidesError(SCHEMA_ERROR, "Manifest field 'layouts' must be a list or object")


def _coerce_layout_entries(manifest: dict[str, Any]) -> list[dict[str, Any]]:
    raw_layouts = manifest.get("layouts")
    if raw_layouts is not None:
        entries: list[dict[str, Any]] = []
        for index, (slug, layout) in enumerate(_coerce_layouts(raw_layouts)):
            entry = dict(layout)
            entry.setdefault("slug", slug)
            entry.setdefault("index", index)
            entry.setdefault("master_index", 0)
            entries.append(entry)
        return entries

    raw_masters = manifest.get("slide_masters", [])
    if not isinstance(raw_masters, list):
        raise AgentSlidesError(SCHEMA_ERROR, "Manifest field 'slide_masters' must be an array")

    entries: list[dict[str, Any]] = []
    for master_index, raw_master in enumerate(raw_masters):
        master = _as_dict(raw_master, context=f"slide_masters[{master_index}]")
        manifest_master_index = master.get("index", master_index)
        for layout_index, raw_layout in enumerate(_coerce_layouts(master.get("layouts", []))):
            slug, layout = raw_layout
            entry = dict(layout)
            entry.setdefault("slug", slug)
            entry.setdefault("index", layout_index)
            entry.setdefault("master_index", manifest_master_index)
            entries.append(entry)
    return entries


def _coerce_placeholder_index(raw_placeholders: object, *, slug: str) -> dict[int, dict[str, Any]]:
    if raw_placeholders is None:
        return {}
    if not isinstance(raw_placeholders, list):
        raise AgentSlidesError(SCHEMA_ERROR, f"layout[{slug!r}].placeholders must be an array")

    placeholders: dict[int, dict[str, Any]] = {}
    for index, raw_placeholder in enumerate(raw_placeholders):
        placeholder = _as_dict(raw_placeholder, context=f"layout[{slug!r}].placeholders[{index}]")
        raw_idx = placeholder.get("idx")
        if isinstance(raw_idx, bool) or not isinstance(raw_idx, int):
            raise AgentSlidesError(SCHEMA_ERROR, f"layout[{slug!r}].placeholders[{index}].idx must be an integer")
        placeholders[raw_idx] = placeholder
    return placeholders


def _reading_order(slots: dict[str, SlotDef]) -> list[tuple[str, SlotDef]]:
    return sorted(
        slots.items(),
        key=lambda item: (
            float(item[1].y if item[1].y is not None else 0.0),
            float(item[1].x if item[1].x is not None else 0.0),
            item[0],
        ),
    )


def _height_close(left: SlotDef, right: SlotDef) -> bool:
    if left.height is None or right.height is None:
        return False
    allowed_delta = max(_ALIGNMENT_EPSILON_PT, min(float(left.height), float(right.height)) * _PEER_HEIGHT_RATIO)
    return abs(float(left.height) - float(right.height)) <= allowed_delta


def _infer_peer_groups(slots: dict[str, SlotDef]) -> dict[str, str]:
    ordered = _reading_order(slots)
    assignments: dict[str, str] = {}
    for index, (slot_name, slot) in enumerate(ordered):
        if slot_name in assignments or slot.role not in {"body", "image", "heading"}:
            continue
        if slot.x is None or slot.y is None or slot.height is None:
            continue

        peers = [slot_name]
        for peer_name, peer_slot in ordered[index + 1 :]:
            if peer_name in assignments:
                continue
            if peer_slot.role != slot.role:
                continue
            if peer_slot.x is None or peer_slot.y is None or peer_slot.height is None:
                continue
            if abs(float(slot.y) - float(peer_slot.y)) > _PEER_Y_EPSILON_PT:
                continue
            if not _height_close(slot, peer_slot):
                continue
            peers.append(peer_name)

        if len(peers) < 2:
            continue

        if slot.role == "body":
            group_name = "columns"
        elif slot.role == "image":
            group_name = "media"
        else:
            group_name = "headings"
        for peer_name in peers:
            assignments[peer_name] = group_name

    return assignments


def _infer_alignment_groups(slots: dict[str, SlotDef]) -> dict[str, str]:
    ordered = _reading_order(slots)
    groups: list[list[str]] = []
    anchors: list[float] = []

    for slot_name, slot in ordered:
        if slot.y is None:
            continue
        assigned = False
        for index, anchor in enumerate(anchors):
            if abs(float(slot.y) - anchor) <= _ALIGNMENT_EPSILON_PT:
                groups[index].append(slot_name)
                assigned = True
                break
        if not assigned:
            anchors.append(float(slot.y))
            groups.append([slot_name])

    assignments: dict[str, str] = {}
    content_index = 0
    for group in groups:
        roles = {slots[slot_name].role for slot_name in group}
        if "heading" in roles and "top" not in assignments.values():
            group_name = "top"
        elif "image" in roles and roles == {"image"}:
            group_name = "media"
        else:
            group_name = "content" if content_index == 0 else f"content_{content_index + 1}"
            content_index += 1
        for slot_name in group:
            assignments[slot_name] = group_name
    return assignments


def _infer_slot_metadata(slots: dict[str, SlotDef]) -> dict[str, SlotDef]:
    if not slots:
        return slots

    peer_groups = _infer_peer_groups(slots)
    alignment_groups = _infer_alignment_groups(slots)
    ordered = _reading_order(slots)
    inferred: dict[str, SlotDef] = {}
    for reading_order, (slot_name, slot) in enumerate(ordered):
        inferred[slot_name] = slot.model_copy(
            update={
                "peer_group": slot.peer_group or peer_groups.get(slot_name),
                "alignment_group": slot.alignment_group or alignment_groups.get(slot_name),
                "reading_order": reading_order,
            }
        )
    return inferred


def _variant_text_fitting(slots: dict[str, SlotDef]) -> dict[str, TextFitting]:
    return {
        slot.role: _default_text_fitting(slot.role)
        for slot in slots.values()
        if slot.role != "image"
    }


def _copy_slot(slot: SlotDef, **updates: Any) -> SlotDef:
    return slot.model_copy(update=updates)


def _is_valid_variant(*, heading: SlotDef, bodies: list[SlotDef], theme: Theme) -> bool:
    from agent_slides.engine.constraints import constraints_from_layout, solve
    from agent_slides.engine.layout_validator import validate_layout

    if heading.x is None or heading.y is None or heading.width is None or heading.height is None:
        return False
    if not bodies:
        return False
    if any(slot.x is None or slot.y is None or slot.width is None or slot.height is None for slot in bodies):
        return False
    if len(bodies) > 1:
        first_group = bodies[0].peer_group
        if first_group is None or any(slot.peer_group != first_group for slot in bodies[1:]):
            return False
        first_alignment = bodies[0].alignment_group
        if first_alignment is None or any(slot.alignment_group != first_alignment for slot in bodies[1:]):
            return False
    validation_slots = {
        "heading": heading,
        **{f"body_{index}": body for index, body in enumerate(bodies, start=1)},
    }
    validation_layout = LayoutDef(
        name="variant_validation",
        slots=validation_slots,
        grid=_TEMPLATE_GRID,
        text_fitting=_variant_text_fitting(validation_slots),
    )
    rects = solve(constraints_from_layout(validation_layout, theme), {}, lambda _slot_name, _content, _width: 0.0)
    return not any(violation.severity == "error" for violation in validate_layout(validation_layout, rects))


def _generate_variants(layout: LayoutDef, theme: Theme) -> list[LayoutDef]:
    heading_slots = [slot for _, slot in _reading_order(layout.slots) if slot.role == "heading"]
    body_slots = [slot for _, slot in _reading_order(layout.slots) if slot.role == "body"]
    if len(heading_slots) != 1 or len(body_slots) < 1:
        return []

    heading = heading_slots[0]
    variants: list[LayoutDef] = []

    if len(body_slots) == 2 and _is_valid_variant(heading=heading, bodies=[body_slots[0]], theme=theme):
        title_content_slots = {
            "heading": _copy_slot(heading, peer_group=None, alignment_group="top", reading_order=0, size_policy="fit_content"),
            "body": _copy_slot(
                body_slots[0],
                peer_group=None,
                alignment_group="content",
                reading_order=1,
                size_policy="fill_remaining",
            ),
        }
        variants.append(
            LayoutDef(
                name="title_content",
                slots=title_content_slots,
                grid=_TEMPLATE_GRID,
                text_fitting=_variant_text_fitting(title_content_slots),
            )
        )

    if len(body_slots) >= 2 and _is_valid_variant(heading=heading, bodies=body_slots[:2], theme=theme):
        two_col_slots = {
            "heading": _copy_slot(heading, peer_group=None, alignment_group="top", reading_order=0, size_policy="fit_content"),
            "col1": _copy_slot(
                body_slots[0],
                peer_group="columns",
                alignment_group="content",
                reading_order=1,
                size_policy="fill_remaining",
            ),
            "col2": _copy_slot(
                body_slots[1],
                peer_group="columns",
                alignment_group="content",
                reading_order=2,
                size_policy="fill_remaining",
            ),
        }
        variants.append(
            LayoutDef(
                name="two_col",
                slots=two_col_slots,
                grid=_TEMPLATE_GRID,
                text_fitting=_variant_text_fitting(two_col_slots),
            )
        )

    if len(body_slots) >= 3 and _is_valid_variant(heading=heading, bodies=body_slots[:3], theme=theme):
        three_col_slots = {
            "heading": _copy_slot(heading, peer_group=None, alignment_group="top", reading_order=0, size_policy="fit_content"),
            "col1": _copy_slot(
                body_slots[0],
                peer_group="columns",
                alignment_group="content",
                reading_order=1,
                size_policy="fill_remaining",
            ),
            "col2": _copy_slot(
                body_slots[1],
                peer_group="columns",
                alignment_group="content",
                reading_order=2,
                size_policy="fill_remaining",
            ),
            "col3": _copy_slot(
                body_slots[2],
                peer_group="columns",
                alignment_group="content",
                reading_order=3,
                size_policy="fill_remaining",
            ),
        }
        variants.append(
            LayoutDef(
                name="three_col",
                slots=three_col_slots,
                grid=_TEMPLATE_GRID,
                text_fitting=_variant_text_fitting(three_col_slots),
            )
        )

    return variants


class TemplateLayoutRegistry:
    """Loads layouts from a template manifest and implements LayoutProvider."""

    def __init__(self, manifest_path: str):
        self._manifest_path = Path(manifest_path).expanduser().resolve(strict=False)
        try:
            payload = json.loads(self._manifest_path.read_text(encoding="utf-8"))
        except FileNotFoundError as exc:
            raise AgentSlidesError(SCHEMA_ERROR, f"Template manifest not found: {self._manifest_path}") from exc
        except json.JSONDecodeError as exc:
            raise AgentSlidesError(
                SCHEMA_ERROR,
                (
                    f"Invalid JSON in template manifest {self._manifest_path}: "
                    f"{exc.msg} at line {exc.lineno} column {exc.colno}"
                ),
            ) from exc

        manifest = _as_dict(payload, context="template manifest")
        self._source_path = self._resolve_source_path(manifest)
        self._source_hash = self._resolve_source_hash(manifest)
        self._theme = self._resolve_theme(manifest)
        self._placeholder_styles = _load_placeholder_styles(self._source_path)
        self._layouts: dict[str, LayoutDef] = {}
        self._slot_names: dict[str, list[str]] = {}
        self._layout_refs: dict[str, tuple[int, int]] = {}
        self._variants: dict[str, list[LayoutDef]] = {}
        self._usable_layouts: list[str] = []
        self._load_layouts(manifest)

    def _load_layouts(self, manifest: dict[str, Any]) -> None:
        for raw_layout in _coerce_layout_entries(manifest):
            slug = _require_string(raw_layout.get("slug"), context="layout.slug")
            slot_mapping = _as_dict(raw_layout.get("slot_mapping", {}), context=f"layout[{slug!r}].slot_mapping")
            placeholders_by_idx = _coerce_placeholder_index(raw_layout.get("placeholders"), slug=slug)
            master_index = raw_layout.get("master_index", 0)
            layout_index = raw_layout.get("index", 0)
            if isinstance(master_index, bool) or not isinstance(master_index, int):
                raise AgentSlidesError(SCHEMA_ERROR, f"layout[{slug!r}].master_index must be an integer")
            if isinstance(layout_index, bool) or not isinstance(layout_index, int):
                raise AgentSlidesError(SCHEMA_ERROR, f"layout[{slug!r}].index must be an integer")
            slots = {
                slot_name: _build_slot(
                    slot_name,
                    slot_value,
                    placeholders_by_idx=placeholders_by_idx,
                    placeholder_styles=self._placeholder_styles,
                    layout_ref=(master_index, layout_index),
                )
                for slot_name, slot_value in slot_mapping.items()
            }
            slots = _infer_slot_metadata(slots)
            text_fitting = {
                slot.role: _default_text_fitting(slot.role)
                for slot in slots.values()
                if slot.role != "image"
            }
            self._layouts[slug] = LayoutDef(
                name=slug,
                slots=slots,
                grid=_TEMPLATE_GRID,
                text_fitting=text_fitting,
            )
            self._slot_names[slug] = list(slot_mapping)
            self._layout_refs[slug] = (master_index, layout_index)
            self._variants[slug] = _generate_variants(self._layouts[slug], self._theme)
            if bool(raw_layout.get("usable", raw_layout.get("is_usable", True))):
                self._usable_layouts.append(slug)

    def _resolve_source_path(self, manifest: dict[str, Any]) -> str:
        raw_source = manifest.get("source")
        if isinstance(raw_source, dict):
            raw_source = raw_source.get("path")
        source = _require_string(raw_source, context="manifest source")
        return str((self._manifest_path.parent / source).resolve(strict=False))

    def _resolve_source_hash(self, manifest: dict[str, Any]) -> str:
        raw_source = manifest.get("source")
        if isinstance(raw_source, dict):
            for key in ("sha256", "hash", "source_hash"):
                value = raw_source.get(key)
                if isinstance(value, str) and value.strip():
                    return value.strip()
        for key in ("source_hash", "source_sha256", "sha256"):
            value = manifest.get(key)
            if isinstance(value, str) and value.strip():
                return value.strip()
        raise AgentSlidesError(SCHEMA_ERROR, "Template manifest is missing source hash metadata")

    def _resolve_theme(self, manifest: dict[str, Any]) -> Theme:
        default_theme = load_theme("default")
        raw_theme = manifest.get("theme", {})
        if not isinstance(raw_theme, dict):
            raise AgentSlidesError(SCHEMA_ERROR, "Manifest field 'theme' must be an object")

        raw_name = raw_theme.get("name")
        if isinstance(raw_name, str) and raw_name.strip():
            theme_name = raw_name.strip()
        else:
            stem = self._manifest_path.name.removesuffix(".json")
            if stem.endswith(".manifest"):
                stem = stem.removesuffix(".manifest")
            slug = _THEME_SLUG_PATTERN.sub("-", stem.casefold()).strip("-") or "template"
            theme_name = f"extracted-{slug}"

        raw_colors = raw_theme.get("colors", {})
        raw_fonts = raw_theme.get("fonts", {})
        raw_spacing = raw_theme.get("spacing", {})
        if not isinstance(raw_colors, dict) or not isinstance(raw_fonts, dict) or not isinstance(raw_spacing, dict):
            raise AgentSlidesError(SCHEMA_ERROR, "Template theme colors, fonts, and spacing must be objects")

        return Theme(
            name=theme_name,
            colors=ThemeColors(
                primary=str(raw_colors.get("primary", default_theme.colors.primary)),
                secondary=str(raw_colors.get("secondary", default_theme.colors.secondary)),
                accent=str(raw_colors.get("accent", default_theme.colors.accent)),
                background=str(raw_colors.get("background", default_theme.colors.background)),
                text=str(raw_colors.get("text", raw_colors.get("foreground", default_theme.colors.text))),
                heading_text=(
                    str(raw_colors["heading_text"])
                    if "heading_text" in raw_colors
                    else default_theme.colors.heading_text
                ),
                subtle_text=(
                    str(raw_colors["subtle_text"])
                    if "subtle_text" in raw_colors
                    else default_theme.colors.subtle_text
                ),
            ),
            fonts=ThemeFonts(
                heading=str(raw_fonts.get("heading", default_theme.fonts.heading)),
                body=str(raw_fonts.get("body", default_theme.fonts.body)),
            ),
            spacing=ThemeSpacing(
                base_unit=float(raw_spacing.get("base_unit", default_theme.spacing.base_unit)),
                margin=float(raw_spacing.get("margin", default_theme.spacing.margin)),
                gutter=float(raw_spacing.get("gutter", default_theme.spacing.gutter)),
            ),
        )

    def get_layout(self, slug: str) -> LayoutDef:
        if slug not in self._layouts or slug not in self._usable_layouts:
            available = ", ".join(self.list_layouts())
            raise AgentSlidesError(
                INVALID_LAYOUT,
                f"Unknown layout '{slug}'. Available layouts: {available}",
            )
        return self._layouts[slug]

    def list_layouts(self) -> list[str]:
        return list(self._usable_layouts)

    def get_slot_names(self, slug: str) -> list[str]:
        layout = self.get_layout(slug)
        return list(self._slot_names[layout.name])

    def get_text_fitting(self, slug: str, role: str) -> TextFitting:
        layout = self.get_layout(slug)
        if role in layout.text_fitting:
            return layout.text_fitting[role]
        return _default_text_fitting(role)

    def get_variants(self, slug: str) -> list[LayoutDef]:
        self.get_layout(slug)
        return list(self._variants.get(slug, []))

    def get_layout_ref(self, slug: str) -> tuple[int, int]:
        self.get_layout(slug)
        return self._layout_refs[slug]

    @property
    def source_path(self) -> str:
        return self._source_path

    @property
    def source_hash(self) -> str:
        return self._source_hash

    @property
    def theme(self) -> Theme:
        return self._theme
