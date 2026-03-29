"""PowerPoint writer for scene-graph decks."""

from __future__ import annotations

import hashlib
import json
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any, cast

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.parts.image import Image
from pptx.shapes.shapetree import SlideShapes
from pptx.util import Emu, Inches, Pt

from agent_slides.engine.conditional_formatting import (
    resolve_chart_point_colors,
    resolve_table_cell_style,
    resolved_text_runs,
)
from agent_slides.errors import AgentSlidesError, FILE_NOT_FOUND, SCHEMA_ERROR, TEMPLATE_CHANGED
from agent_slides.icons import ICON_VIEWBOX, require_icon, svg_path_subpaths
from agent_slides.io.assets import resolve_image_path
from agent_slides.model.types import (
    ComputedNode,
    ComputedPatternElement,
    Deck,
    EMU_PER_POINT,
    Node,
    TextBlock,
    TextRun,
    split_text_runs_by_line,
)
from agent_slides.model.design_rules import ConditionalFormatting, load_design_rules
from agent_slides.model.themes import load_theme

BLANK_LAYOUT_INDEX = 6
BULLET_MARGIN_STEP_PT = 18.0
BLOCK_SPACING_FACTOR = 0.3
LINE_HEIGHT_FACTOR = 1.2
ICON_BULLET_SIZE_FACTOR = 0.78
ICON_BULLET_GAP_FACTOR = 0.4
HEADING_SIZE_FACTOR = 1.35
CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
}
SHAPE_TYPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "arrow": MSO_SHAPE.RIGHT_ARROW,
    "chevron": MSO_SHAPE.CHEVRON,
}
DASH_STYLE_MAP = {
    "dash": MSO_LINE_DASH_STYLE.DASH,
    "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
    "dashDot": MSO_LINE_DASH_STYLE.DASH_DOT,
}


@dataclass(frozen=True)
class TemplateLayoutBinding:
    master_index: int
    layout_index: int
    slot_mapping: dict[str, int]


class TemplateLayoutRegistry:
    """Read template manifest metadata needed by the PPTX writer."""

    def __init__(self, manifest_path: str | Path) -> None:
        self.manifest_path = Path(manifest_path)
        payload = _read_manifest_payload(self.manifest_path)
        self._source = _require_string(payload, "source")
        self._source_hash = _require_string(payload, "source_hash")
        self._layouts = _read_layout_bindings(payload)

    @property
    def source_hash(self) -> str:
        return self._source_hash

    @property
    def source_path(self) -> Path:
        return (self.manifest_path.parent / self._source).resolve()

    def binding_for(self, slug: str) -> TemplateLayoutBinding:
        binding = self._layouts.get(slug)
        if binding is None:
            raise AgentSlidesError(
                SCHEMA_ERROR,
                f"Template manifest does not define layout '{slug}'",
                details={"layout": slug, "manifest": str(self.manifest_path)},
            )
        return binding


def points_to_emu(value_pt: float) -> Emu:
    """Convert points to EMU for python-pptx geometry APIs."""

    return Emu(int(round(value_pt * EMU_PER_POINT)))


def hex_to_rgb(value: str) -> RGBColor:
    """Convert a #RRGGBB-style color string into an RGBColor."""

    normalized = value.lstrip("#")
    return RGBColor.from_string(normalized)


def _block_font_size(computed: ComputedNode, block: TextBlock, *, default_font_size: float | None = None) -> float:
    if default_font_size is not None:
        return default_font_size
    if block.type == "heading":
        return computed.font_size_pt * HEADING_SIZE_FACTOR
    return computed.font_size_pt


def _block_line_runs(
    block: TextBlock,
    conditional_formatting: ConditionalFormatting | None = None,
) -> list[list[TextRun]]:
    formatted_block = block.model_copy(update={"runs": resolved_text_runs(block, conditional_formatting)})
    return split_text_runs_by_line(formatted_block)


def _block_lines(block: TextBlock) -> list[str]:
    lines = block.text.splitlines()
    return lines or [""]


def _block_line_height(computed: ComputedNode, block: TextBlock) -> float:
    return _block_font_size(computed, block) * LINE_HEIGHT_FACTOR


def _icon_bullet_indent(block: TextBlock, computed: ComputedNode | None = None) -> float:
    font_size = computed.font_size_pt if computed is not None else 14.0
    icon_size = font_size * ICON_BULLET_SIZE_FACTOR
    icon_gap = font_size * ICON_BULLET_GAP_FACTOR
    return (block.level * BULLET_MARGIN_STEP_PT) + icon_size + icon_gap


def _run_font_size(
    computed: ComputedNode,
    block: TextBlock,
    run_spec: TextRun,
    *,
    default_font_size: float | None = None,
) -> float:
    if run_spec.font_size is not None:
        return run_spec.font_size
    return _block_font_size(computed, block, default_font_size=default_font_size)


def _set_run_strikethrough(run, enabled: bool) -> None:
    rPr = run._r.get_or_add_rPr()
    if enabled:
        rPr.set("strike", "sngStrike")
    elif "strike" in rPr.attrib:
        del rPr.attrib["strike"]


def _apply_template_run_styles(
    run,
    run_spec: TextRun,
    *,
    computed: ComputedNode | None = None,
    block: TextBlock | None = None,
    default_font_size: float | None = None,
) -> None:
    if computed is not None:
        if computed.font_family:
            run.font.name = computed.font_family
        if block is not None:
            run.font.size = Pt(
                _run_font_size(
                    computed,
                    block,
                    run_spec,
                    default_font_size=default_font_size,
                )
            )
    elif run_spec.font_size is not None:
        run.font.size = Pt(run_spec.font_size)
    if run_spec.bold is not None:
        run.font.bold = run_spec.bold
    if run_spec.italic is not None:
        run.font.italic = run_spec.italic
    if run_spec.color:
        run.font.color.rgb = hex_to_rgb(run_spec.color)
    if run_spec.underline:
        run.font.underline = True
    if run_spec.strikethrough:
        _set_run_strikethrough(run, True)


def _apply_text_run_defaults(
    run,
    computed: ComputedNode,
    block: TextBlock,
    run_spec: TextRun,
    *,
    default_font_size: float | None = None,
) -> None:
    run.font.name = computed.font_family
    run.font.size = Pt(_run_font_size(computed, block, run_spec, default_font_size=default_font_size))
    run.font.bold = computed.font_bold if run_spec.bold is None else run_spec.bold
    if block.type == "heading" and run_spec.bold is None:
        run.font.bold = True
    run.font.italic = run_spec.italic or False
    run.font.color.rgb = hex_to_rgb(run_spec.color or computed.color)
    if run_spec.underline:
        run.font.underline = True
    _set_run_strikethrough(run, run_spec.strikethrough)


def _mix_hex_colors(base: str, overlay: str, ratio: float) -> str:
    base_rgb = [int(base.lstrip("#")[index : index + 2], 16) for index in (0, 2, 4)]
    overlay_rgb = [int(overlay.lstrip("#")[index : index + 2], 16) for index in (0, 2, 4)]
    blended = [
        int(round((1.0 - ratio) * base_channel + ratio * overlay_channel))
        for base_channel, overlay_channel in zip(base_rgb, overlay_rgb, strict=True)
    ]
    return "#" + "".join(f"{channel:02X}" for channel in blended)


def _is_dark_color(value: str) -> bool:
    red, green, blue = [int(value.lstrip("#")[index : index + 2], 16) for index in (0, 2, 4)]
    luminance = (0.299 * red) + (0.587 * green) + (0.114 * blue)
    return luminance < 160.0


def _table_alignment(value: str) -> PP_ALIGN:
    if value == "center":
        return PP_ALIGN.CENTER
    if value == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def _fit_image_to_slot(computed: ComputedNode, image_size_px: tuple[int, int]) -> tuple[float, float, float, float]:
    slot_x = computed.x
    slot_y = computed.y
    slot_width = computed.width
    slot_height = computed.height

    if computed.image_fit == "stretch":
        return slot_x, slot_y, slot_width, slot_height

    image_width_px, image_height_px = image_size_px
    if computed.image_fit == "cover":
        scale = max(slot_width / image_width_px, slot_height / image_height_px)
    else:
        scale = min(slot_width / image_width_px, slot_height / image_height_px)
    width = image_width_px * scale
    height = image_height_px * scale
    return (
        slot_x + ((slot_width - width) / 2),
        slot_y + ((slot_height - height) / 2),
        width,
        height,
    )


def _read_manifest_payload(path: Path) -> dict[str, Any]:
    try:
        raw = path.read_text(encoding="utf-8")
    except FileNotFoundError as exc:
        raise AgentSlidesError(FILE_NOT_FOUND, f"Manifest file not found: {path}") from exc
    except OSError as exc:
        raise AgentSlidesError(SCHEMA_ERROR, f"Failed to read manifest file {path}: {exc}") from exc

    try:
        payload = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise AgentSlidesError(SCHEMA_ERROR, f"Manifest file is not valid JSON: {path}") from exc

    if not isinstance(payload, dict):
        raise AgentSlidesError(SCHEMA_ERROR, "Manifest root must be a JSON object.")
    return payload


def _require_string(payload: dict[str, Any], field: str) -> str:
    value = payload.get(field)
    if not isinstance(value, str) or not value:
        raise AgentSlidesError(SCHEMA_ERROR, f"Manifest field '{field}' must be a non-empty string.")
    return value


def _require_list(payload: dict[str, Any], field: str) -> list[Any]:
    value = payload.get(field)
    if not isinstance(value, list):
        raise AgentSlidesError(SCHEMA_ERROR, f"Manifest field '{field}' must be an array.")
    return value


def _require_dict(payload: dict[str, Any], field: str) -> dict[str, Any]:
    value = payload.get(field)
    if not isinstance(value, dict):
        raise AgentSlidesError(SCHEMA_ERROR, f"Manifest field '{field}' must be an object.")
    return value


def _require_int(payload: dict[str, Any], field: str) -> int:
    value = payload.get(field)
    if not isinstance(value, int) or value < 0:
        raise AgentSlidesError(SCHEMA_ERROR, f"Manifest field '{field}' must be a non-negative integer.")
    return value


def _read_layout_bindings(payload: dict[str, Any]) -> dict[str, TemplateLayoutBinding]:
    bindings: dict[str, TemplateLayoutBinding] = {}

    for master_position, master in enumerate(_require_list(payload, "slide_masters")):
        if not isinstance(master, dict):
            raise AgentSlidesError(SCHEMA_ERROR, f"Manifest slide_masters[{master_position}] must be an object.")

        for layout_position, layout in enumerate(_require_list(master, "layouts")):
            if not isinstance(layout, dict):
                raise AgentSlidesError(
                    SCHEMA_ERROR,
                    f"Manifest slide_masters[{master_position}].layouts[{layout_position}] must be an object.",
                )

            slug = _require_string(layout, "slug")
            slot_mapping = _require_dict(layout, "slot_mapping")
            binding = TemplateLayoutBinding(
                master_index=layout.get("master_index", master_position),
                layout_index=layout.get("index", layout_position),
                slot_mapping={slot: _require_int(slot_mapping, slot) for slot in slot_mapping},
            )
            if slug in bindings:
                raise AgentSlidesError(SCHEMA_ERROR, f"Manifest layout slug '{slug}' must be unique.")
            bindings[slug] = binding

    return bindings


def _node_paragraphs(
    node: Node,
    *,
    conditional_formatting: ConditionalFormatting | None = None,
) -> list[tuple[int, TextBlock, list[TextRun]]]:
    if isinstance(node.content, str):
        block = TextBlock(type="paragraph", text=node.content)
        return [(0, block, line_runs) for line_runs in _block_line_runs(block, conditional_formatting)]

    paragraphs: list[tuple[int, TextBlock, list[TextRun]]] = []
    for block_index, block in enumerate(node.content.blocks):
        paragraphs.extend((block_index, block, line_runs) for line_runs in _block_line_runs(block, conditional_formatting))
    return paragraphs or [(0, TextBlock(type="paragraph", text=""), [TextRun(text="")])]


def _configure_paragraph_bullets(paragraph, block: TextBlock, *, computed: ComputedNode | None = None) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    pPr.remove_all("a:buNone", "a:buAutoNum", "a:buChar", "a:buBlip")

    if block.type == "bullet":
        pPr.lvl = block.level
        if block.icon:
            pPr.insert_element_before(OxmlElement("a:buNone"), "a:tabLst", "a:defRPr", "a:extLst")
            pPr.set("marL", str(Pt(_icon_bullet_indent(block, computed))))
            pPr.set("indent", "0")
            return
        buChar = OxmlElement("a:buChar")
        buChar.set("char", "•")
        pPr.insert_element_before(buChar, "a:tabLst", "a:defRPr", "a:extLst")
        pPr.set("marL", str(Pt(BULLET_MARGIN_STEP_PT * (block.level + 1))))
        pPr.set("indent", str(Pt(-BULLET_MARGIN_STEP_PT)))
        return

    pPr.lvl = 0
    for attr_name in ("marL", "indent"):
        if attr_name in pPr.attrib:
            del pPr.attrib[attr_name]
    pPr.insert_element_before(OxmlElement("a:buNone"), "a:tabLst", "a:defRPr", "a:extLst")


def _configure_text_frame(text_frame) -> None:
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0


def _apply_run_style(run, *, font_family: str, font_size_pt: float, bold: bool, color: str) -> None:
    run.font.name = font_family
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _emit_warning(payload: dict[str, object]) -> None:
    sys.stderr.write(f"{json.dumps(payload)}\n")


def _warn_if_template_changed(template_path: Path, expected_hash: str) -> None:
    actual_hash = _sha256(template_path)
    if actual_hash == expected_hash:
        return

    _emit_warning(
        {
            "warning": {
                "code": TEMPLATE_CHANGED,
                "message": "Template source file changed since the manifest was learned.",
            },
            "data": {
                "template": str(template_path),
                "expected_hash": expected_hash,
                "actual_hash": actual_hash,
            },
        }
    )


def _node_z_index(node: Node) -> int:
    z_index = node.style_overrides.get("z_index")
    if isinstance(z_index, int):
        return z_index
    if node.type == "shape":
        return -1
    return 0


def _iter_rendered_nodes(nodes: list[Node]) -> list[Node]:
    return [
        node
        for _, node in sorted(
            enumerate(nodes),
            key=lambda item: (_node_z_index(item[1]), item[0]),
        )
    ]


def _shape_fill_color(node: Node) -> str | None:
    spec = node.shape_spec
    if spec is None:
        return None
    if spec.fill_color is not None:
        return spec.fill_color
    if spec.shape_type in {"arrow", "chevron"}:
        return spec.line_color
    return None


def _apply_shape_shadow(shape) -> None:
    shape.shadow.inherit = False
    effect_list = shape._element.spPr.get_or_add_effectLst()
    effect_list.clear()
    shadow = OxmlElement("a:outerShdw")
    shadow.set("blurRad", str(points_to_emu(3.0)))
    shadow.set("dist", str(points_to_emu(2.0)))
    shadow.set("dir", "5400000")
    shadow.set("algn", "ctr")
    shadow.set("rotWithShape", "0")
    color = OxmlElement("a:srgbClr")
    color.set("val", "000000")
    alpha = OxmlElement("a:alpha")
    alpha.set("val", "18000")
    color.append(alpha)
    shadow.append(color)
    effect_list.append(shadow)


def _apply_shape_line(shape, node: Node) -> None:
    spec = node.shape_spec
    if spec is None:
        return

    color = spec.line_color
    if spec.shape_type == "line":
        color = color or spec.fill_color or "#333333"

    if color is None or spec.line_width == 0:
        shape.line.fill.background()
        return

    shape.line.color.rgb = hex_to_rgb(color)
    shape.line.width = Pt(spec.line_width)
    shape.line.fill.transparency = max(0.0, min(1.0, 1.0 - spec.opacity))
    if spec.dash is not None:
        shape.line.dash_style = DASH_STYLE_MAP[spec.dash]


def _render_shape_node(slide_shape_collection: SlideShapes, node: Node, computed: ComputedNode) -> None:
    spec = node.shape_spec
    if spec is None:
        return

    if spec.shape_type == "line":
        shape = slide_shape_collection.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            points_to_emu(computed.x),
            points_to_emu(computed.y),
            points_to_emu(computed.x + computed.width),
            points_to_emu(computed.y + computed.height),
        )
        _apply_shape_line(shape, node)
        if spec.shadow:
            _apply_shape_shadow(shape)
        return

    shape = slide_shape_collection.add_shape(
        SHAPE_TYPE_MAP[spec.shape_type],
        points_to_emu(computed.x),
        points_to_emu(computed.y),
        points_to_emu(computed.width),
        points_to_emu(computed.height),
    )

    fill_color = _shape_fill_color(node)
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        shape.fill.transparency = max(0.0, min(1.0, 1.0 - spec.opacity))

    _apply_shape_line(shape, node)

    if spec.shape_type == "rounded_rectangle" and len(shape.adjustments) > 0 and computed.height > 0 and computed.width > 0:
        shape.adjustments[0] = min(1.0, max(0.0, (spec.corner_radius * 2.0) / min(computed.width, computed.height)))

    if spec.shadow:
        _apply_shape_shadow(shape)


def _pattern_vertical_anchor(value: str) -> MSO_ANCHOR:
    if value == "middle":
        return MSO_ANCHOR.MIDDLE
    if value == "bottom":
        return MSO_ANCHOR.BOTTOM
    return MSO_ANCHOR.TOP


def _render_pattern_shape_element(slide_shape_collection: SlideShapes, element: ComputedPatternElement) -> None:
    node = Node(
        node_id="pattern-shape",
        type="shape",
        shape_spec={
            "shape_type": element.shape_type,
            "fill_color": element.fill_color,
            "line_color": element.line_color,
            "line_width": element.line_width,
            "corner_radius": element.corner_radius,
            "shadow": element.shadow,
            "dash": element.dash,
            "opacity": element.opacity,
        },
        style_overrides={
            "x": element.x,
            "y": element.y,
            "width": element.width,
            "height": element.height,
            "z_index": element.z_index,
        },
    )
    computed = ComputedNode(
        x=element.x,
        y=element.y,
        width=element.width,
        height=element.height,
        revision=0,
        content_type="shape",
    )
    _render_shape_node(slide_shape_collection, node, computed)


def _render_pattern_text_element(slide_shape_collection: SlideShapes, element: ComputedPatternElement) -> None:
    shape = slide_shape_collection.add_textbox(
        points_to_emu(element.x),
        points_to_emu(element.y),
        points_to_emu(element.width),
        points_to_emu(element.height),
    )
    shape.line.fill.background()
    shape.fill.background()

    text_frame = shape.text_frame
    _configure_text_frame(text_frame)
    text_frame.vertical_anchor = _pattern_vertical_anchor(element.vertical_align)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = _table_alignment(element.text_align)
    run = paragraph.add_run()
    run.text = element.text or ""
    _apply_run_style(
        run,
        font_family=element.font_family,
        font_size_pt=element.font_size_pt,
        bold=element.font_bold,
        color=element.color,
    )


def _render_pattern_node(slide_shape_collection: SlideShapes, computed: ComputedNode) -> None:
    for element in sorted(computed.pattern_elements, key=lambda item: item.z_index):
        if element.kind == "shape":
            _render_pattern_shape_element(slide_shape_collection, element)
        else:
            _render_pattern_text_element(slide_shape_collection, element)


def _delete_all_slides(prs: Presentation) -> None:
    """Remove every slide from a presentation while keeping masters and layouts."""

    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        prs.slides._sldIdLst.remove(slide_id)
        prs.part.drop_rel(slide_id.rId)


def _resolve_layout(prs: Presentation, slide_layout: str, binding: TemplateLayoutBinding):
    try:
        slide_master = prs.slide_masters[binding.master_index]
    except IndexError as exc:
        raise AgentSlidesError(
            SCHEMA_ERROR,
            (
                f"Template layout '{slide_layout}' references missing master index "
                f"{binding.master_index}"
            ),
            details={
                "layout": slide_layout,
                "master_index": binding.master_index,
                "layout_index": binding.layout_index,
            },
        ) from exc

    try:
        return slide_master.slide_layouts[binding.layout_index]
    except IndexError as exc:
        raise AgentSlidesError(
            SCHEMA_ERROR,
            (
                f"Template layout '{slide_layout}' references missing layout index "
                f"{binding.layout_index} in master {binding.master_index}"
            ),
            details={
                "layout": slide_layout,
                "master_index": binding.master_index,
                "layout_index": binding.layout_index,
            },
        ) from exc


def _fill_placeholder(
    node: Node,
    slot_mapping: dict[str, int],
    slide,
    *,
    computed: ComputedNode | None = None,
    conditional_formatting: ConditionalFormatting | None = None,
) -> None:
    if node.slot_binding is None or node.type != "text":
        return

    placeholder_idx = slot_mapping.get(node.slot_binding)
    if placeholder_idx is None:
        return

    placeholder = slide.placeholders[placeholder_idx]
    text_frame = placeholder.text_frame
    text_frame.clear()

    paragraphs = _node_paragraphs(node, conditional_formatting=conditional_formatting)
    positions_by_index = (
        {position.block_index: position for position in computed.block_positions}
        if computed is not None and computed.block_positions
        else {}
    )

    for paragraph_index, (block_index, block, line_runs) in enumerate(paragraphs):
        paragraph = text_frame.paragraphs[0] if paragraph_index == 0 else text_frame.add_paragraph()
        _configure_paragraph_bullets(paragraph, block, computed=computed)
        position = positions_by_index.get(block_index)
        for run_spec in line_runs:
            run = paragraph.add_run()
            run.text = run_spec.text
            _apply_template_run_styles(
                run,
                run_spec,
                computed=computed,
                block=block,
                default_font_size=position.font_size_pt if position is not None else None,
            )


def _write_template_pptx(deck: Deck, output_path: str) -> None:
    registry = TemplateLayoutRegistry(deck.template_manifest)
    template_path = registry.source_path
    conditional_formatting = load_design_rules(deck.design_rules).conditional_formatting

    if not template_path.exists():
        raise AgentSlidesError(FILE_NOT_FOUND, f"Template file not found: {template_path}")

    _warn_if_template_changed(template_path, registry.source_hash)
    presentation = Presentation(template_path)
    _delete_all_slides(presentation)

    for slide in deck.slides:
        binding = registry.binding_for(slide.layout)
        pptx_slide = presentation.slides.add_slide(_resolve_layout(presentation, slide.layout, binding))
        for node in _iter_rendered_nodes(slide.nodes):
            if node.type in {"shape", "pattern"}:
                computed = slide.computed.get(node.node_id)
                if computed is not None:
                    if node.type == "shape":
                        _render_shape_node(pptx_slide.shapes, node, computed)
                    else:
                        _render_pattern_node(pptx_slide.shapes, computed)
        for node in slide.nodes:
            if node.type == "icon":
                computed = slide.computed.get(node.node_id)
                if computed is not None:
                    _render_icon_node(pptx_slide.shapes, node, computed)
                continue
            _fill_placeholder(
                node,
                binding.slot_mapping,
                pptx_slide,
                computed=slide.computed.get(node.node_id),
                conditional_formatting=conditional_formatting,
            )

    presentation.save(Path(output_path))


def _render_text_node(
    slide_shape_collection: SlideShapes,
    node: Node,
    computed: ComputedNode,
    *,
    conditional_formatting: ConditionalFormatting | None = None,
) -> None:
    """Render a single text node as a positioned text box."""

    shape = slide_shape_collection.add_textbox(
        points_to_emu(computed.x),
        points_to_emu(computed.y),
        points_to_emu(computed.width),
        points_to_emu(computed.height),
    )
    shape.line.fill.background()

    if computed.bg_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(computed.bg_color)
        shape.fill.transparency = computed.bg_transparency
    else:
        shape.fill.background()

    text_frame = shape.text_frame
    _configure_text_frame(text_frame)
    blocks = node.content.blocks or [TextBlock(type="paragraph", text="")]
    if computed.block_positions:
        positions_by_index = {position.block_index: position for position in computed.block_positions}
        first_position = computed.block_positions[0]
        last_position = computed.block_positions[-1]
        left_margin = max(first_position.x - computed.x, 0.0)
        top_margin = max(first_position.y - computed.y, 0.0)
        right_margin = max((computed.x + computed.width) - (first_position.x + first_position.width), 0.0)
        bottom_margin = max((computed.y + computed.height) - (last_position.y + last_position.height), 0.0)

        text_frame.margin_left = points_to_emu(left_margin)
        text_frame.margin_right = points_to_emu(right_margin)
        text_frame.margin_top = points_to_emu(top_margin)
        text_frame.margin_bottom = points_to_emu(bottom_margin)

        paragraph_index = 0
        for index, block in enumerate(blocks):
            position = positions_by_index.get(index)
            if position is None:
                continue

            paragraph = text_frame.paragraphs[0] if paragraph_index == 0 else text_frame.add_paragraph()
            _configure_paragraph_bullets(paragraph, block)
            paragraph.space_before = Pt(0)
            next_position = positions_by_index.get(index + 1)
            gap_after = 0.0 if next_position is None else max(next_position.y - (position.y + position.height), 0.0)
            paragraph.space_after = Pt(gap_after)

            for run_spec in block.resolved_runs():
                run = paragraph.add_run()
                run.text = run_spec.text
                _apply_text_run_defaults(
                    run,
                    computed,
                    block,
                    run_spec,
                    default_font_size=position.font_size_pt,
                )
            paragraph_index += 1
        return

    paragraph_index = 0
    for block in blocks:
        for line_runs in _block_line_runs(block, conditional_formatting):
            paragraph = text_frame.paragraphs[0] if paragraph_index == 0 else text_frame.add_paragraph()
            _configure_paragraph_bullets(paragraph, block, computed=computed)
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            for run_spec in line_runs:
                run = paragraph.add_run()
                run.text = run_spec.text
                _apply_text_run_defaults(run, computed, block, run_spec)
            paragraph_index += 1

    _render_text_block_icons(slide_shape_collection, node, computed)


def _render_icon_shape(
    slide_shape_collection: SlideShapes,
    *,
    path_data: str,
    x: float,
    y: float,
    size: float,
    color: str,
) -> None:
    scale = size / ICON_VIEWBOX
    for subpath in svg_path_subpaths(path_data):
        points = [(x + (px * scale), y + (py * scale)) for px, py in subpath]
        if len(points) < 2:
            continue
        is_closed = points[0] == points[-1]
        vertices = points[1:-1] if is_closed else points[1:]
        builder = slide_shape_collection.build_freeform(
            points_to_emu(points[0][0]),
            points_to_emu(points[0][1]),
        )
        if vertices:
            builder.add_line_segments(
                tuple((points_to_emu(px), points_to_emu(py)) for px, py in vertices),
                close=is_closed,
            )
        shape = builder.convert_to_shape()
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(color)
        shape.line.fill.background()


def _render_text_block_icons(slide_shape_collection: SlideShapes, node: Node, computed: ComputedNode) -> None:
    blocks = node.content.blocks or [TextBlock(type="paragraph", text="")]
    current_y = computed.y
    for index, block in enumerate(blocks):
        line_height = _block_line_height(computed, block)
        for _line in _block_lines(block):
            if block.type == "bullet" and block.icon:
                icon_size = computed.font_size_pt * ICON_BULLET_SIZE_FACTOR
                icon_x = computed.x + (block.level * BULLET_MARGIN_STEP_PT)
                icon_y = current_y + ((line_height - icon_size) / 2)
                _render_icon_shape(
                    slide_shape_collection,
                    path_data=require_icon(block.icon),
                    x=icon_x,
                    y=icon_y,
                    size=icon_size,
                    color=computed.color,
                )
            current_y += line_height
        if index < len(blocks) - 1:
            current_y += computed.font_size_pt * BLOCK_SPACING_FACTOR


def _render_image_node(
    slide_shape_collection: SlideShapes,
    node: Node,
    computed: ComputedNode,
    *,
    asset_base_dir: str | Path | None = None,
) -> None:
    """Render a single image node as a positioned picture."""

    if node.image_path is None:
        return

    image_path = resolve_image_path(node.image_path, base_dir=asset_base_dir)
    image = Image.from_file(str(image_path))
    left, top, width, height = _fit_image_to_slot(computed, cast(tuple[int, int], image.size))
    slide_shape_collection.add_picture(
        str(image_path),
        points_to_emu(left),
        points_to_emu(top),
        width=points_to_emu(width),
        height=points_to_emu(height),
    )
def _render_chart_node(
    slide_shape_collection: SlideShapes,
    node: Node,
    computed: ComputedNode,
    *,
    conditional_formatting: ConditionalFormatting | None = None,
) -> None:
    """Render a single chart node as a native editable PowerPoint chart."""

    spec = node.chart_spec
    if spec is None:
        return

    if spec.chart_type == "scatter":
        chart_data = XyChartData()
        for scatter_series in spec.scatter_series:
            series = chart_data.add_series(scatter_series.name)
            for point in scatter_series.points:
                series.add_data_point(point.x, point.y)
        pptx_type = XL_CHART_TYPE.XY_SCATTER
    else:
        chart_data = CategoryChartData()
        chart_data.categories = spec.categories
        for category_series in spec.series:
            chart_data.add_series(category_series.name, category_series.values)
        pptx_type = CHART_TYPE_MAP[spec.chart_type]

    chart_frame = slide_shape_collection.add_chart(
        pptx_type,
        points_to_emu(computed.x),
        points_to_emu(computed.y),
        points_to_emu(computed.width),
        points_to_emu(computed.height),
        chart_data,
    )
    chart = chart_frame.chart

    if spec.title:
        chart.has_title = True
        chart.chart_title.text_frame.text = spec.title
    chart.has_legend = spec.style.has_legend

    # Native PowerPoint charts do not inherit agent-slides theme colors.
    for index, color in enumerate(spec.style.series_colors or []):
        if index >= len(chart.series):
            break
        fill = chart.series[index].format.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(color)

    point_colors = resolve_chart_point_colors(spec, conditional_formatting)
    if point_colors and chart.series:
        series = chart.series[0]
        for index, color in enumerate(point_colors):
            if index >= len(series.points):
                break
            fill = series.points[index].format.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(color)


def _write_table_cell(
    cell,
    text: str,
    *,
    font_family: str,
    font_size: float,
    font_bold: bool,
    font_color: str,
    fill_color: str,
    alignment: str,
) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = hex_to_rgb(fill_color)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    text_frame = cell.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.margin_left = points_to_emu(4.0)
    text_frame.margin_right = points_to_emu(4.0)
    text_frame.margin_top = points_to_emu(2.0)
    text_frame.margin_bottom = points_to_emu(2.0)

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = _table_alignment(alignment)
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_family
    run.font.size = Pt(font_size)
    run.font.bold = font_bold
    run.font.color.rgb = hex_to_rgb(font_color)


def _render_table_node(
    slide_shape_collection: SlideShapes,
    node: Node,
    computed: ComputedNode,
    *,
    theme,
    conditional_formatting: ConditionalFormatting | None = None,
) -> None:
    """Render a single table node as a native editable PowerPoint table."""

    spec = node.table_spec
    if spec is None:
        return

    row_count = len(spec.rows) + 1
    column_count = len(spec.headers)
    table = slide_shape_collection.add_table(
        row_count,
        column_count,
        points_to_emu(computed.x),
        points_to_emu(computed.y),
        points_to_emu(computed.width),
        points_to_emu(computed.height),
    ).table

    total_width = int(points_to_emu(computed.width))
    width_weights = spec.resolved_col_widths()
    weight_total = sum(width_weights) or float(column_count)
    width_allocated = 0
    for column_index, weight in enumerate(width_weights):
        if column_index == column_count - 1:
            column_width = total_width - width_allocated
        else:
            column_width = int(round(total_width * (weight / weight_total)))
            width_allocated += column_width
        table.columns[column_index].width = Emu(column_width)

    total_height = int(points_to_emu(computed.height))
    row_height = total_height // row_count
    height_allocated = 0
    for row_index in range(row_count):
        if row_index == row_count - 1:
            height = total_height - height_allocated
        else:
            height = row_height
            height_allocated += height
        table.rows[row_index].height = Emu(height)

    header_fill = spec.header_color or theme.colors.primary
    header_text_color = "#FFFFFF" if _is_dark_color(header_fill) else theme.colors.text
    stripe_fill = _mix_hex_colors(theme.colors.background, theme.colors.secondary, 0.08)
    alignments = spec.resolved_col_align()

    table.first_row = True
    table.horz_banding = spec.stripe

    for column_index, header in enumerate(spec.headers):
        _write_table_cell(
            table.cell(0, column_index),
            header,
            font_family=computed.font_family,
            font_size=spec.header_font_size,
            font_bold=True,
            font_color=header_text_color,
            fill_color=header_fill,
            alignment="center",
        )

    for row_index, row in enumerate(spec.rows, start=1):
        fill_color = stripe_fill if spec.stripe and row_index % 2 == 0 else theme.colors.background
        for column_index, value in enumerate(row):
            cell_fill, cell_text_color, cell_bold = resolve_table_cell_style(
                value,
                default_fill=fill_color,
                default_text=computed.color,
                conditional_formatting=conditional_formatting,
            )
            _write_table_cell(
                table.cell(row_index, column_index),
                value,
                font_family=computed.font_family,
                font_size=spec.font_size,
                font_bold=cell_bold,
                font_color=cell_text_color,
                fill_color=cell_fill,
                alignment=alignments[column_index],
            )


def _render_icon_node(slide_shape_collection: SlideShapes, node: Node, computed: ComputedNode) -> None:
    path_data = computed.icon_svg_path or require_icon(str(node.icon_name))
    _render_icon_shape(
        slide_shape_collection,
        path_data=path_data,
        x=computed.x,
        y=computed.y,
        size=computed.width,
        color=computed.color,
    )


render_text_node = _render_text_node
render_image_node = _render_image_node
render_icon_node = _render_icon_node


def _write_v0_pptx(deck: Deck, output_path: str, *, asset_base_dir: str | Path | None = None) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[BLANK_LAYOUT_INDEX]
    theme = load_theme(deck.theme)
    conditional_formatting = load_design_rules(deck.design_rules).conditional_formatting

    for slide in deck.slides:
        pptx_slide = presentation.slides.add_slide(blank_layout)
        if not slide.computed:
            continue

        for node in _iter_rendered_nodes(slide.nodes):
            computed = slide.computed.get(node.node_id)
            if computed is None:
                continue

            if node.type == "shape":
                _render_shape_node(pptx_slide.shapes, node, computed)
            elif node.type == "pattern":
                _render_pattern_node(pptx_slide.shapes, computed)
            elif node.type == "chart":
                _render_chart_node(
                    pptx_slide.shapes,
                    node,
                    computed,
                    conditional_formatting=conditional_formatting,
                )
            elif node.type == "icon":
                _render_icon_node(pptx_slide.shapes, node, computed)
            elif node.type == "table":
                _render_table_node(
                    pptx_slide.shapes,
                    node,
                    computed,
                    theme=theme,
                    conditional_formatting=conditional_formatting,
                )
            elif node.type == "image":
                _render_image_node(
                    pptx_slide.shapes,
                    node,
                    computed,
                    asset_base_dir=asset_base_dir,
                )
            elif node.slot_binding is not None:
                _render_text_node(
                    pptx_slide.shapes,
                    node,
                    computed,
                    conditional_formatting=conditional_formatting,
                )

    presentation.save(Path(output_path))


def write_pptx(deck: Deck, output_path: str, *, asset_base_dir: str | Path | None = None) -> None:
    """Write a deck to PowerPoint using either the v0 or template-backed writer."""

    if deck.template_manifest:
        _write_template_pptx(deck, output_path)
        return

    _write_v0_pptx(deck, output_path, asset_base_dir=asset_base_dir)
