"""Read PPTX templates and emit layout manifest JSON."""

from __future__ import annotations

import hashlib
import json
import os
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Mapping
from xml.etree import ElementTree as ET
from zipfile import BadZipFile

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from agent_slides.errors import AgentSlidesError, FILE_NOT_FOUND, SCHEMA_ERROR
from agent_slides.model.layouts import DEFAULT_GUTTER_PT, DEFAULT_MARGIN_PT
from agent_slides.model.types import EMU_PER_POINT
from agent_slides.template_slots import infer_template_slot_role

DEFAULT_BASE_UNIT_PT = 10.0
BODY_ROW_THRESHOLD_PT = 54.0
OLE_MAGIC = bytes.fromhex("D0CF11E0A1B11AE1")
THEME_RELATIONSHIP = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
)
PRESENTATIONML_NS = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
DRAWINGML_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
OOXML_NS = {**PRESENTATIONML_NS, **DRAWINGML_NS}
SLUG_PATTERN = re.compile(r"[^a-z0-9]+")
NON_PLACEHOLDER_IDX_OFFSET = 1_000_000
BODY_SLOT_TYPES = {"BODY", "TEXT_BOX", "TABLE", "CHART", "GROUP"}
SPECIAL_SLOT_TYPES = {"quote", "attribution"}
SLIDE_WIDTH_PT = 960.0
SLIDE_HEIGHT_PT = 540.0
SLIDE_AREA_PT = SLIDE_WIDTH_PT * SLIDE_HEIGHT_PT
PANEL_AREA_THRESHOLD_PT = SLIDE_AREA_PT * 0.15
DEFAULT_DARK_TEXT = "333333"
DEFAULT_LIGHT_TEXT = "FFFFFF"
DEFAULT_BG_COLOR = "FFFFFF"

_PLACEHOLDER_TYPE_MAP = {
    PP_PLACEHOLDER.TITLE: "TITLE",
    PP_PLACEHOLDER.CENTER_TITLE: "TITLE",
    PP_PLACEHOLDER.VERTICAL_TITLE: "TITLE",
    PP_PLACEHOLDER.SUBTITLE: "SUBTITLE",
    PP_PLACEHOLDER.BODY: "BODY",
    PP_PLACEHOLDER.OBJECT: "BODY",
    PP_PLACEHOLDER.VERTICAL_BODY: "BODY",
    PP_PLACEHOLDER.VERTICAL_OBJECT: "BODY",
    PP_PLACEHOLDER.PICTURE: "PICTURE",
    PP_PLACEHOLDER.TABLE: "TABLE",
    PP_PLACEHOLDER.CHART: "CHART",
}
_WARNING_ONLY_PLACEHOLDERS = {
    PP_PLACEHOLDER.MEDIA_CLIP,
    PP_PLACEHOLDER.ORG_CHART,
}


@dataclass(frozen=True)
class LearnResult:
    manifest_path: Path
    manifest: dict[str, object]
    layouts_found: int
    usable_layouts: int
    warnings: list[str]

    @property
    def source(self) -> str:
        return str(self.manifest["source"])


@dataclass(frozen=True)
class _Rect:
    left: float
    top: float
    width: float
    height: float

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    @property
    def area(self) -> float:
        return self.width * self.height


def read_template_manifest(
    template_path: str | Path,
    output_path: str | Path | None = None,
) -> LearnResult:
    """Read a PPTX template and persist its manifest JSON."""

    source_path = Path(template_path)
    if not source_path.exists():
        raise AgentSlidesError(
            FILE_NOT_FOUND, f"Template file not found: {source_path}"
        )

    manifest_path = (
        Path(output_path)
        if output_path is not None
        else _default_manifest_path(source_path)
    )
    presentation = _open_presentation(source_path)
    return _build_manifest(source_path, manifest_path, presentation)


def _default_manifest_path(template_path: Path) -> Path:
    return template_path.with_name(f"{template_path.stem}.manifest.json")


def _open_presentation(template_path: Path) -> Presentation:
    try:
        return Presentation(template_path)
    except FileNotFoundError as exc:
        raise AgentSlidesError(
            FILE_NOT_FOUND, f"Template file not found: {template_path}"
        ) from exc
    except (BadZipFile, KeyError, ValueError) as exc:
        if _looks_like_encrypted_office_file(template_path):
            raise AgentSlidesError(
                SCHEMA_ERROR, "password-protected files not supported"
            ) from exc
        raise AgentSlidesError(SCHEMA_ERROR, "not a valid PPTX file") from exc
    except Exception as exc:
        raise AgentSlidesError(SCHEMA_ERROR, "not a valid PPTX file") from exc


def _looks_like_encrypted_office_file(path: Path) -> bool:
    try:
        with path.open("rb") as handle:
            return handle.read(len(OLE_MAGIC)) == OLE_MAGIC
    except OSError:
        return False


def _build_manifest(
    template_path: Path, manifest_path: Path, presentation: Presentation
) -> LearnResult:
    layouts_found = 0
    usable_layouts = 0
    warnings: list[str] = []
    slug_counts: dict[str, int] = {}
    masters: list[dict[str, object]] = []
    theme: dict[str, object] | None = None
    scheme_colors: dict[str, str] = {}

    for master_index, slide_master in enumerate(presentation.slide_masters):
        layouts: list[dict[str, object]] = []
        for layout_index, slide_layout in enumerate(slide_master.slide_layouts):
            if theme is None:
                theme = _extract_theme(presentation)
                scheme_colors = _theme_scheme_colors(theme)
            layouts_found += 1
            layout_name = slide_layout.name or f"Layout {layout_index + 1}"
            (
                placeholders,
                slot_mapping,
                color_zones,
                editable_regions,
                layout_warnings,
            ) = _extract_layout(
                slide_layout,
                layout_name,
                scheme_colors=scheme_colors,
            )
            warnings.extend(layout_warnings)

            slug = _unique_slug(layout_name, slug_counts)
            usable = bool(placeholders) or _is_blank_layout_name(layout_name)
            if usable:
                usable_layouts += 1

            layouts.append(
                {
                    "index": layout_index,
                    "master_index": master_index,
                    "name": layout_name,
                    "slug": slug,
                    "usable": usable,
                    "placeholders": placeholders,
                    "slot_mapping": slot_mapping,
                    "color_zones": color_zones,
                    "editable_regions": editable_regions,
                }
            )

        masters.append(
            {
                "index": master_index,
                "name": f"Slide Master {master_index + 1}",
                "layouts": layouts,
            }
        )

    if layouts_found == 0:
        raise AgentSlidesError(SCHEMA_ERROR, "template has no slide layouts")
    if usable_layouts == 0:
        warnings.append("template has 0 usable layouts")
    if theme is None:
        theme = _extract_theme(presentation)

    manifest = {
        "source": _relative_source_path(template_path, manifest_path),
        "source_hash": _sha256(template_path),
        "slide_masters": masters,
        "theme": theme,
    }
    _write_manifest(manifest_path, manifest)

    return LearnResult(
        manifest_path=manifest_path,
        manifest=manifest,
        layouts_found=layouts_found,
        usable_layouts=usable_layouts,
        warnings=warnings,
    )


def _extract_layout(
    slide_layout: object,
    layout_name: str,
    *,
    scheme_colors: Mapping[str, str],
) -> tuple[
    list[dict[str, object]],
    dict[str, int],
    list[dict[str, object]],
    list[dict[str, object]],
    list[str],
]:
    placeholders: list[dict[str, object]] = []
    warnings: list[str] = []

    for placeholder in slide_layout.placeholders:
        placeholder_type = _normalize_placeholder_type(
            placeholder.placeholder_format.type
        )
        if placeholder_type is None:
            raw_type = placeholder.placeholder_format.type
            if raw_type in _WARNING_ONLY_PLACEHOLDERS:
                warnings.append(
                    f"layout '{layout_name}': skipped unsupported {raw_type.name.lower()} placeholder '{placeholder.name}'"
                )
            continue

        placeholder_record = _shape_record(
            idx=int(placeholder.placeholder_format.idx),
            type_tag=placeholder_type,
            name=placeholder.name,
            left=placeholder.left,
            top=placeholder.top,
            width=placeholder.width,
            height=placeholder.height,
            shape_id=getattr(placeholder, "shape_id", None),
            shape_kind="placeholder",
        )
        placeholders.append(placeholder_record)

    for shape in getattr(slide_layout, "shapes", []):
        if getattr(shape, "is_placeholder", False):
            continue
        extracted = _extract_non_placeholder_shape(shape)
        if extracted is not None:
            placeholders.append(extracted)

    placeholders.sort(key=lambda item: int(item["idx"]))
    slot_mapping = _build_slot_mapping(placeholders)
    color_zones = _extract_color_zones(
        slide_layout,
        placeholders,
        scheme_colors=scheme_colors,
    )
    return (
        placeholders,
        slot_mapping,
        color_zones,
        _extract_editable_regions(
            slide_layout,
            placeholders,
            color_zones=color_zones,
        ),
        warnings,
    )


def _extract_editable_regions(
    slide_layout: object,
    placeholders: list[dict[str, object]],
    *,
    color_zones: list[dict[str, object]],
) -> list[dict[str, object]]:
    placeholder_union = _placeholder_union_region(placeholders)
    if placeholder_union is not None:
        return [_region_record("content_box", placeholder_union, "placeholder_union")]

    inferred = _infer_editable_region(
        slide_layout,
        placeholders,
        color_zones=color_zones,
    )
    if inferred is None:
        return []
    return [
        _region_record(
            "content_area",
            inferred,
            "visual_inference_no_placeholders",
        )
    ]


def _placeholder_union_region(
    placeholders: list[dict[str, object]],
) -> _Rect | None:
    content_placeholders = [
        placeholder
        for placeholder in placeholders
        if placeholder.get("shape_kind") == "placeholder"
        and placeholder["type"] in {"BODY", "TABLE", "CHART"}
    ]
    return _union_rect(content_placeholders)


_MIN_EDITABLE_REGION_WIDTH_PT = 40.0
_MIN_EDITABLE_REGION_HEIGHT_PT = 40.0
_MIN_EDITABLE_REGION_AREA_PT = 4_000.0
_FIXED_LABEL_MAX_HEIGHT_PT = 24.0
_FIXED_LABEL_BOTTOM_BAND_PT = 72.0


def _infer_editable_region(
    slide_layout: object,
    placeholders: list[dict[str, object]],
    *,
    color_zones: list[dict[str, object]],
) -> _Rect | None:
    base_region = _select_visual_inference_region(
        color_zones,
        title_placeholder=_title_placeholder(placeholders),
    )
    if base_region is None:
        return None

    candidates = [base_region]
    for obstacle in _visual_obstacle_rects(slide_layout, placeholders, base_region):
        next_candidates: list[_Rect] = []
        for candidate in candidates:
            next_candidates.extend(_subtract_rect(candidate, obstacle))
        candidates = _filter_candidate_rects(next_candidates)
        if not candidates:
            return None

    return max(candidates, key=lambda candidate: candidate.area)


def _select_visual_inference_region(
    color_zones: list[dict[str, object]],
    *,
    title_placeholder: dict[str, object] | None,
) -> _Rect | None:
    title_zone = _zone_containing_title(color_zones, title_placeholder)
    if title_zone is not None:
        zone_rect = _zone_rect(title_zone)
        top = zone_rect.top
        if title_placeholder is not None:
            title_rect = _rect_from_placeholder(title_placeholder)
            if title_rect is not None:
                top = max(
                    top, min(SLIDE_HEIGHT_PT, title_rect.bottom + DEFAULT_GUTTER_PT)
                )
        return _rect_from_edges(zone_rect.left, top, zone_rect.right, SLIDE_HEIGHT_PT)

    zones = [
        _zone_rect(zone)
        for zone in color_zones
        if float(zone.get("width", 0)) > 0
        and (
            str(zone.get("region", "")).startswith("gap_")
            or zone.get("region") == "full_slide"
        )
    ]
    if not zones:
        zones = [
            _zone_rect(zone) for zone in color_zones if float(zone.get("width", 0)) > 0
        ]
    if not zones:
        return _Rect(left=0.0, top=0.0, width=SLIDE_WIDTH_PT, height=SLIDE_HEIGHT_PT)

    widest_zone = max(zones, key=lambda zone: (zone.width, zone.area))
    title_rect = _rect_from_placeholder(title_placeholder)
    top = widest_zone.top
    if title_rect is not None and _ranges_overlap(
        widest_zone.left, widest_zone.right, title_rect.left, title_rect.right
    ):
        top = max(top, min(SLIDE_HEIGHT_PT, title_rect.bottom + DEFAULT_GUTTER_PT))
    return _rect_from_edges(
        widest_zone.left,
        top,
        widest_zone.right,
        SLIDE_HEIGHT_PT,
    )


def _zone_containing_title(
    color_zones: list[dict[str, object]],
    title_placeholder: dict[str, object] | None,
) -> dict[str, object] | None:
    title_rect = _rect_from_placeholder(title_placeholder)
    if title_rect is None:
        return None
    midpoint = title_rect.left + (title_rect.width / 2)
    for zone in color_zones:
        left = float(zone.get("left", 0))
        right = left + float(zone.get("width", 0))
        if left <= midpoint <= right:
            return zone
    return None


def _visual_obstacle_rects(
    slide_layout: object,
    placeholders: list[dict[str, object]],
    base_region: _Rect,
) -> list[_Rect]:
    obstacles: list[_Rect] = []
    title_placeholder = _title_placeholder(placeholders)
    if title_placeholder is not None:
        title_rect = _rect_from_placeholder(title_placeholder)
        if title_rect is not None:
            obstacles.append(
                _expand_rect(
                    title_rect,
                    bottom=DEFAULT_GUTTER_PT,
                )
            )

    for shape in getattr(slide_layout, "shapes", []):
        if getattr(shape, "is_placeholder", False):
            continue
        if _is_panel_shape(shape):
            continue

        obstacle = _obstacle_rect_for_shape(shape)
        if obstacle is None:
            continue
        if not _rectangles_overlap(base_region, obstacle):
            continue
        obstacles.append(obstacle)

    obstacles.sort(key=lambda obstacle: (obstacle.top, obstacle.left))
    return obstacles


def _obstacle_rect_for_shape(shape: object) -> _Rect | None:
    rect = _rect_from_shape(shape)
    if rect is None:
        return None
    if getattr(shape, "has_text_frame", False) and _shape_text(shape).strip():
        if _is_fixed_label_shape(shape, rect):
            return _expand_rect(
                rect, top=DEFAULT_GUTTER_PT / 2, bottom=DEFAULT_GUTTER_PT / 2
            )
        return None
    if getattr(shape, "has_table", False) or getattr(shape, "has_chart", False):
        return None
    return rect


def _is_fixed_label_shape(shape: object, rect: _Rect) -> bool:
    text = _shape_text(shape).strip().casefold()
    name = str(getattr(shape, "name", "")).casefold()
    if "copyright" in text or "copyright" in name:
        return True
    if "footer" in name or "confidential" in text or "page " in text:
        return True
    if rect.height <= _FIXED_LABEL_MAX_HEIGHT_PT:
        return True
    return rect.top >= SLIDE_HEIGHT_PT - _FIXED_LABEL_BOTTOM_BAND_PT


def _is_panel_shape(shape: object) -> bool:
    fill = getattr(shape, "fill", None)
    if fill is None or getattr(fill, "type", None) != MSO_FILL.SOLID:
        return False
    width = _emu_to_points(getattr(shape, "width", 0))
    height = _emu_to_points(getattr(shape, "height", 0))
    return width > 0 and height > 0 and (width * height) >= PANEL_AREA_THRESHOLD_PT


def _filter_candidate_rects(candidates: list[_Rect]) -> list[_Rect]:
    return [
        candidate
        for candidate in candidates
        if candidate.width >= _MIN_EDITABLE_REGION_WIDTH_PT
        and candidate.height >= _MIN_EDITABLE_REGION_HEIGHT_PT
        and candidate.area >= _MIN_EDITABLE_REGION_AREA_PT
    ]


def _subtract_rect(candidate: _Rect, obstacle: _Rect) -> list[_Rect]:
    left = max(candidate.left, obstacle.left)
    top = max(candidate.top, obstacle.top)
    right = min(candidate.right, obstacle.right)
    bottom = min(candidate.bottom, obstacle.bottom)

    if left >= right or top >= bottom:
        return [candidate]

    remainder: list[_Rect] = []
    remainder.append(
        _rect_from_edges(candidate.left, candidate.top, candidate.right, top)
    )
    remainder.append(
        _rect_from_edges(candidate.left, bottom, candidate.right, candidate.bottom)
    )
    remainder.append(_rect_from_edges(candidate.left, top, left, bottom))
    remainder.append(_rect_from_edges(right, top, candidate.right, bottom))
    return [rect for rect in remainder if rect is not None]


def _union_rect(items: list[dict[str, object]]) -> _Rect | None:
    rects = [
        rect
        for rect in (_rect_from_placeholder(item) for item in items)
        if rect is not None
    ]
    if not rects:
        return None
    return _Rect(
        left=min(rect.left for rect in rects),
        top=min(rect.top for rect in rects),
        width=max(rect.right for rect in rects) - min(rect.left for rect in rects),
        height=max(rect.bottom for rect in rects) - min(rect.top for rect in rects),
    )


def _rect_from_placeholder(placeholder: dict[str, object] | None) -> _Rect | None:
    if placeholder is None:
        return None
    bounds = placeholder.get("bounds")
    if not isinstance(bounds, Mapping):
        return None
    left = float(bounds.get("x", 0))
    top = float(bounds.get("y", 0))
    width = float(bounds.get("w", bounds.get("width", 0)))
    height = float(bounds.get("h", bounds.get("height", 0)))
    return _rect_from_edges(left, top, left + width, top + height)


def _rect_from_shape(shape: object) -> _Rect | None:
    left = _emu_to_points(getattr(shape, "left", 0))
    top = _emu_to_points(getattr(shape, "top", 0))
    width = _emu_to_points(getattr(shape, "width", 0))
    height = _emu_to_points(getattr(shape, "height", 0))
    return _rect_from_edges(left, top, left + width, top + height)


def _zone_rect(zone: Mapping[str, object]) -> _Rect:
    left = float(zone.get("left", 0))
    width = float(zone.get("width", 0))
    return _Rect(
        left=left,
        top=0.0,
        width=width,
        height=SLIDE_HEIGHT_PT,
    )


def _expand_rect(
    rect: _Rect,
    *,
    left: float = 0.0,
    top: float = 0.0,
    right: float = 0.0,
    bottom: float = 0.0,
) -> _Rect:
    return _Rect(
        left=max(0.0, rect.left - left),
        top=max(0.0, rect.top - top),
        width=max(
            0.0,
            min(SLIDE_WIDTH_PT, rect.right + right) - max(0.0, rect.left - left),
        ),
        height=max(
            0.0,
            min(SLIDE_HEIGHT_PT, rect.bottom + bottom) - max(0.0, rect.top - top),
        ),
    )


def _rect_from_edges(
    left: float,
    top: float,
    right: float,
    bottom: float,
) -> _Rect | None:
    clamped_left = max(0.0, min(SLIDE_WIDTH_PT, left))
    clamped_top = max(0.0, min(SLIDE_HEIGHT_PT, top))
    clamped_right = max(0.0, min(SLIDE_WIDTH_PT, right))
    clamped_bottom = max(0.0, min(SLIDE_HEIGHT_PT, bottom))
    if clamped_right <= clamped_left or clamped_bottom <= clamped_top:
        return None
    return _Rect(
        left=clamped_left,
        top=clamped_top,
        width=clamped_right - clamped_left,
        height=clamped_bottom - clamped_top,
    )


def _rectangles_overlap(first: _Rect, second: _Rect) -> bool:
    return _ranges_overlap(
        first.left, first.right, second.left, second.right
    ) and _ranges_overlap(first.top, first.bottom, second.top, second.bottom)


def _ranges_overlap(
    first_start: float,
    first_end: float,
    second_start: float,
    second_end: float,
) -> bool:
    return max(first_start, second_start) < min(first_end, second_end)


def _region_record(name: str, rect: _Rect, source: str) -> dict[str, object]:
    return {
        "name": name,
        "left": round(rect.left, 3),
        "top": round(rect.top, 3),
        "width": round(rect.width, 3),
        "height": round(rect.height, 3),
        "source": source,
    }


def _extract_color_zones(
    slide_layout: object,
    placeholders: list[dict[str, object]],
    *,
    scheme_colors: Mapping[str, str],
) -> list[dict[str, object]]:
    background = _layout_background_color(slide_layout, scheme_colors=scheme_colors)
    panels = sorted(
        _extract_panel_candidates(slide_layout),
        key=lambda panel: (float(panel["left"]), float(panel["width"])),
    )

    if not panels:
        return [
            _build_zone(
                region="full_slide",
                left=0.0,
                width=SLIDE_WIDTH_PT,
                bg_color=background,
                title_placeholder=_title_placeholder(placeholders),
            )
        ]

    zones: list[dict[str, object]] = []
    cursor = 0.0
    gap_index = 0
    title_placeholder = _title_placeholder(placeholders)
    for panel_index, panel in enumerate(panels):
        panel_left = float(panel["left"])
        panel_width = float(panel["width"])
        panel_right = min(SLIDE_WIDTH_PT, panel_left + panel_width)
        if panel_left > cursor:
            zones.append(
                _build_zone(
                    region=f"gap_{gap_index}",
                    left=cursor,
                    width=panel_left - cursor,
                    bg_color=background,
                    title_placeholder=title_placeholder,
                )
            )
            gap_index += 1

        zones.append(
            _build_zone(
                region=f"panel_{panel_index}",
                left=panel_left,
                width=max(0.0, panel_right - panel_left),
                bg_color=str(panel["bg_color"]),
                title_placeholder=title_placeholder,
            )
        )
        cursor = max(cursor, panel_right)

    if cursor < SLIDE_WIDTH_PT:
        zones.append(
            _build_zone(
                region=f"gap_{gap_index}",
                left=cursor,
                width=SLIDE_WIDTH_PT - cursor,
                bg_color=background,
                title_placeholder=title_placeholder,
            )
        )

    return [zone for zone in zones if float(zone["width"]) > 0]


def _extract_panel_candidates(slide_layout: object) -> list[dict[str, object]]:
    panels: list[dict[str, object]] = []
    for shape in getattr(slide_layout, "shapes", []):
        if getattr(shape, "is_placeholder", False):
            continue
        fill = getattr(shape, "fill", None)
        if fill is None or getattr(fill, "type", None) != MSO_FILL.SOLID:
            continue
        bg_color = _shape_fill_color(shape)
        if bg_color is None:
            continue

        left = max(0.0, _emu_to_points(getattr(shape, "left", 0)))
        width = _emu_to_points(getattr(shape, "width", 0))
        height = _emu_to_points(getattr(shape, "height", 0))
        if width <= 0 or height <= 0:
            continue
        if width * height < PANEL_AREA_THRESHOLD_PT:
            continue

        panels.append(
            {
                "left": left,
                "width": max(0.0, min(SLIDE_WIDTH_PT - left, width)),
                "bg_color": bg_color,
            }
        )
    return panels


def _shape_fill_color(shape: object) -> str | None:
    fill = getattr(shape, "fill", None)
    if fill is None:
        return None
    color = getattr(fill, "fore_color", None)
    rgb = getattr(color, "rgb", None)
    if rgb is None:
        return None
    return _normalize_hex_color(str(rgb))


def _layout_background_color(
    slide_layout: object, *, scheme_colors: Mapping[str, str]
) -> str:
    for element in (
        getattr(slide_layout, "_element", None),
        getattr(getattr(slide_layout, "slide_master", None), "_element", None),
    ):
        bg_pr = _find_background_properties(element)
        if bg_pr is None:
            continue
        color = _solid_fill_color(bg_pr, scheme_colors=scheme_colors)
        if color is not None:
            return color
    return _normalize_hex_color(scheme_colors.get("lt1", DEFAULT_BG_COLOR))


def _find_background_properties(element: object | None) -> object | None:
    if element is None or not hasattr(element, "find"):
        return None
    background = element.find("./p:cSld/p:bg/p:bgPr", OOXML_NS)
    if background is not None:
        return background
    return element.find(".//p:bgPr", OOXML_NS)


def _solid_fill_color(
    parent: object | None, *, scheme_colors: Mapping[str, str]
) -> str | None:
    if parent is None or not hasattr(parent, "find"):
        return None
    solid_fill = parent.find("./a:solidFill", OOXML_NS)
    if solid_fill is not None:
        return _resolve_color_element(solid_fill, scheme_colors=scheme_colors)
    grad_fill = parent.find("./a:gradFill", OOXML_NS)
    if grad_fill is not None:
        first_stop = grad_fill.find("./a:gsLst/a:gs", OOXML_NS)
        if first_stop is not None:
            return _resolve_color_element(first_stop, scheme_colors=scheme_colors)
    return None


def _resolve_color_element(
    element: object, *, scheme_colors: Mapping[str, str]
) -> str | None:
    if element is None or not hasattr(element, "find"):
        return None
    srgb = element.find("./a:srgbClr", OOXML_NS)
    if srgb is not None:
        return _normalize_hex_color(srgb.attrib.get("val"))
    system = element.find("./a:sysClr", OOXML_NS)
    if system is not None:
        return _normalize_hex_color(
            system.attrib.get("lastClr", system.attrib.get("val"))
        )
    scheme = element.find("./a:schemeClr", OOXML_NS)
    if scheme is not None:
        scheme_name = scheme.attrib.get("val")
        if scheme_name:
            return _normalize_hex_color(
                scheme_colors.get(scheme_name, DEFAULT_BG_COLOR)
            )
    return None


def _title_placeholder(
    placeholders: list[dict[str, object]],
) -> dict[str, object] | None:
    titles = sorted(
        _placeholders_of_type(placeholders, "TITLE"), key=_sort_key_by_position
    )
    return titles[0] if titles else None


def _build_zone(
    *,
    region: str,
    left: float,
    width: float,
    bg_color: str,
    title_placeholder: dict[str, object] | None,
) -> dict[str, object]:
    zone = {
        "region": region,
        "left": round(left, 3),
        "width": round(width, 3),
        "bg_color": _normalize_hex_color(bg_color),
    }
    zone["text_color"] = _contrasting_text_color(str(zone["bg_color"]))
    _attach_editable_regions(zone, title_placeholder)
    return zone


def _attach_editable_regions(
    zone: dict[str, object], title_placeholder: dict[str, object] | None
) -> None:
    if title_placeholder is None:
        return
    bounds = title_placeholder.get("bounds")
    if not isinstance(bounds, dict):
        return
    zone_left = float(zone["left"])
    zone_right = zone_left + float(zone["width"])
    title_left = float(bounds.get("x", 0))
    title_width = float(bounds.get("w", bounds.get("width", 0)))
    title_midpoint = title_left + (title_width / 2)
    if not (zone_left <= title_midpoint <= zone_right):
        return

    title_top = max(0.0, float(bounds.get("y", 0)))
    title_height = max(0.0, float(bounds.get("h", bounds.get("height", 0))))
    title_bottom = min(SLIDE_HEIGHT_PT, title_top + title_height)

    if title_top > 0:
        zone["editable_above"] = {
            "left": zone_left,
            "top": 0.0,
            "width": float(zone["width"]),
            "height": round(title_top, 3),
        }
    if title_bottom < SLIDE_HEIGHT_PT:
        zone["editable_below"] = {
            "left": zone_left,
            "top": round(title_bottom, 3),
            "width": float(zone["width"]),
            "height": round(SLIDE_HEIGHT_PT - title_bottom, 3),
        }


def _contrasting_text_color(bg_color: str) -> str:
    red, green, blue = _hex_to_rgb(bg_color)
    luminance = ((0.2126 * red) + (0.7152 * green) + (0.0722 * blue)) / 255.0
    return DEFAULT_DARK_TEXT if luminance > 0.65 else DEFAULT_LIGHT_TEXT


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    normalized = _normalize_hex_color(value)
    return (
        int(normalized[0:2], 16),
        int(normalized[2:4], 16),
        int(normalized[4:6], 16),
    )


def _normalize_hex_color(value: str | None) -> str:
    if not value:
        return DEFAULT_BG_COLOR
    normalized = str(value).strip().lstrip("#").upper()
    if len(normalized) == 3:
        normalized = "".join(channel * 2 for channel in normalized)
    if len(normalized) != 6:
        return DEFAULT_BG_COLOR
    return normalized


def _theme_scheme_colors(theme: Mapping[str, object]) -> dict[str, str]:
    colors = theme.get("colors")
    if not isinstance(colors, Mapping):
        return {}
    dk1 = _normalize_hex_color(str(colors.get("text", DEFAULT_DARK_TEXT)))
    dk2 = _normalize_hex_color(str(colors.get("heading_text", DEFAULT_DARK_TEXT)))
    lt1 = _normalize_hex_color(str(colors.get("background", DEFAULT_BG_COLOR)))
    lt2 = _normalize_hex_color(str(colors.get("subtle_text", DEFAULT_LIGHT_TEXT)))
    return {
        "accent1": _normalize_hex_color(str(colors.get("primary", DEFAULT_BG_COLOR))),
        "accent2": _normalize_hex_color(str(colors.get("secondary", DEFAULT_BG_COLOR))),
        "accent3": _normalize_hex_color(str(colors.get("accent", DEFAULT_BG_COLOR))),
        "dk1": dk1,
        "dk2": dk2,
        "lt1": lt1,
        "lt2": lt2,
        # OOXML aliases: tx1/tx2 are the same as dk1/dk2, bg1/bg2 = lt1/lt2
        "tx1": dk1,
        "tx2": dk2,
        "bg1": lt1,
        "bg2": lt2,
    }


def _normalize_placeholder_type(raw_type: PP_PLACEHOLDER) -> str | None:
    return _PLACEHOLDER_TYPE_MAP.get(raw_type)


def _build_slot_mapping(placeholders: list[dict[str, object]]) -> dict[str, int]:
    slot_mapping: dict[str, int] = {}

    titles = _placeholders_of_type(placeholders, "TITLE")
    if titles:
        slot_mapping["heading"] = int(
            sorted(titles, key=_sort_key_by_position)[0]["idx"]
        )

    subtitles = _placeholders_of_type(placeholders, "SUBTITLE")
    if subtitles:
        slot_mapping["subheading"] = int(
            sorted(subtitles, key=_sort_key_by_position)[0]["idx"]
        )

    claimed_indexes = set(slot_mapping.values())
    suggestions = _collect_suggested_slots(placeholders)
    for slot_name in sorted(SPECIAL_SLOT_TYPES):
        suggested = suggestions.get(slot_name, [])
        if len(suggested) == 1:
            placeholder_idx = int(suggested[0]["idx"])
            slot_mapping[slot_name] = placeholder_idx
            claimed_indexes.add(placeholder_idx)

    unmapped_body_like = [
        placeholder
        for placeholder in placeholders
        if placeholder["type"] in BODY_SLOT_TYPES
        and int(placeholder["idx"]) not in claimed_indexes
        and _placeholder_height(placeholder) >= _MIN_BODY_HEIGHT_PT
    ]
    preferred_types = (
        {"BODY"}
        if any(placeholder["type"] == "BODY" for placeholder in unmapped_body_like)
        else BODY_SLOT_TYPES
    )
    bodies = sorted(
        [
            placeholder
            for placeholder in unmapped_body_like
            if placeholder["type"] in preferred_types
        ],
        key=_sort_key_by_position,
    )
    if len(bodies) == 1:
        slot_mapping["body"] = int(bodies[0]["idx"])
    elif len(bodies) > 1:
        if _same_row(bodies):
            for index, placeholder in enumerate(
                sorted(bodies, key=_sort_key_by_x), start=1
            ):
                slot_mapping[f"col{index}"] = int(placeholder["idx"])
        else:
            slot_mapping["body"] = int(bodies[0]["idx"])

    pictures = sorted(
        _placeholders_of_type(placeholders, "PICTURE"), key=_sort_key_by_position
    )
    if pictures:
        slot_mapping["image"] = int(pictures[0]["idx"])

    return slot_mapping


def _placeholders_of_type(
    placeholders: list[dict[str, object]], placeholder_type: str
) -> list[dict[str, object]]:
    return [
        placeholder
        for placeholder in placeholders
        if placeholder["type"] == placeholder_type
    ]


def _same_row(placeholders: list[dict[str, object]]) -> bool:
    top_values = [float(placeholder["bounds"]["y"]) for placeholder in placeholders]
    return max(top_values) - min(top_values) <= BODY_ROW_THRESHOLD_PT


def _sort_key_by_position(placeholder: dict[str, object]) -> tuple[float, float, int]:
    bounds = placeholder["bounds"]
    return float(bounds["y"]), float(bounds["x"]), int(placeholder["idx"])


def _sort_key_by_x(placeholder: dict[str, object]) -> tuple[float, int]:
    bounds = placeholder["bounds"]
    return float(bounds["x"]), int(placeholder["idx"])


def _extract_non_placeholder_shape(shape: Any) -> dict[str, object] | None:
    base = _shape_record(
        idx=_synthetic_shape_idx(getattr(shape, "shape_id", 0)),
        type_tag="UNKNOWN",
        name=getattr(shape, "name", "Shape"),
        left=getattr(shape, "left", 0),
        top=getattr(shape, "top", 0),
        width=getattr(shape, "width", 0),
        height=getattr(shape, "height", 0),
        shape_id=getattr(shape, "shape_id", None),
        shape_kind=_shape_kind(shape),
    )

    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        base["type"] = "GROUP"
        base["group"] = {"children": len(getattr(shape, "shapes", []))}
        base["suggested_slot"] = _suggest_slot_name(base)
        return base

    if getattr(shape, "has_table", False):
        table = shape.table
        base["type"] = "TABLE"
        base["table"] = {"rows": len(table.rows), "cols": len(table.columns)}
        base["suggested_slot"] = _suggest_slot_name(base)
        return base

    if getattr(shape, "has_chart", False):
        base["type"] = "CHART"
        base["chart"] = {"chart_type": int(shape.chart.chart_type)}
        base["suggested_slot"] = _suggest_slot_name(base)
        return base

    if getattr(shape, "has_text_frame", False) and _shape_text(shape).strip():
        base["type"] = "TEXT_BOX"
        base["text"] = _shape_text(shape)
        base["suggested_slot"] = _suggest_slot_name(base)
        return base

    return None


def _shape_record(
    *,
    idx: int,
    type_tag: str,
    name: str,
    left: int,
    top: int,
    width: int,
    height: int,
    shape_id: int | None,
    shape_kind: str,
) -> dict[str, object]:
    placeholder: dict[str, object] = {
        "idx": idx,
        "type": type_tag,
        "name": name,
        "bounds": {
            "x": _emu_to_points(left),
            "y": _emu_to_points(top),
            "w": _emu_to_points(width),
            "h": _emu_to_points(height),
        },
        "shape_kind": shape_kind,
    }
    if shape_id is not None:
        placeholder["shape_id"] = int(shape_id)
    placeholder["suggested_slot"] = _suggest_slot_name(placeholder)
    return placeholder


def _synthetic_shape_idx(shape_id: int) -> int:
    return NON_PLACEHOLDER_IDX_OFFSET + int(shape_id)


def _shape_kind(shape: Any) -> str:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type is None:
        return "shape"
    name = getattr(shape_type, "name", None)
    if isinstance(name, str) and name:
        return name.casefold()
    return str(shape_type).casefold()


def _shape_text(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text_frame.text
    return ""


def _collect_suggested_slots(
    placeholders: list[dict[str, object]],
) -> dict[str, list[dict[str, object]]]:
    suggestions: dict[str, list[dict[str, object]]] = {}
    for placeholder in placeholders:
        slot_name = placeholder.get("suggested_slot")
        if isinstance(slot_name, str) and slot_name:
            suggestions.setdefault(slot_name, []).append(placeholder)
    return suggestions


_MIN_BODY_HEIGHT_PT = 20.0


def _placeholder_height(placeholder: Mapping[str, object]) -> float:
    bounds = placeholder.get("bounds")
    if isinstance(bounds, dict):
        return float(bounds.get("h", bounds.get("height", 0)))
    return 0.0


def _suggest_slot_name(placeholder: Mapping[str, object]) -> str | None:
    placeholder_type = placeholder.get("type")
    if placeholder_type == "TITLE":
        return "heading"
    if placeholder_type == "SUBTITLE":
        return "subheading"
    if placeholder_type == "PICTURE":
        return "image"
    if placeholder_type not in BODY_SLOT_TYPES:
        return None

    # Skip tiny text boxes (< 20pt tall) — these are typically copyright or
    # footer lines, not real body content areas.
    bounds = placeholder.get("bounds")
    if isinstance(bounds, dict):
        height = float(bounds.get("h", bounds.get("height", 0)))
        if height < _MIN_BODY_HEIGHT_PT:
            return None

    role = infer_template_slot_role(str(placeholder.get("name", "")), placeholder)
    if role in {"quote", "attribution"}:
        return role
    return "body"


def _is_blank_layout_name(layout_name: str) -> bool:
    tokens = {token for token in SLUG_PATTERN.split(layout_name.casefold()) if token}
    return "blank" in tokens


def _unique_slug(layout_name: str, slug_counts: dict[str, int]) -> str:
    base_slug = _slugify(layout_name)
    count = slug_counts.get(base_slug, 0) + 1
    slug_counts[base_slug] = count
    return base_slug if count == 1 else f"{base_slug}_{count}"


def _slugify(value: str) -> str:
    slug = SLUG_PATTERN.sub("_", value.strip().lower()).strip("_")
    return slug or "layout"


def _extract_theme(presentation: Presentation) -> dict[str, object]:
    theme_root = _theme_root(presentation)
    colors = {
        "primary": _theme_color(theme_root, "accent1"),
        "secondary": _theme_color(theme_root, "accent2"),
        "accent": _theme_color(theme_root, "accent3"),
        "text": _theme_color(theme_root, "dk1"),
        "heading_text": _theme_color(theme_root, "dk2"),
        "subtle_text": _theme_color(theme_root, "lt2"),
        "background": _theme_color(theme_root, "lt1"),
    }
    fonts = {
        "heading": _theme_font(theme_root, "majorFont"),
        "body": _theme_font(theme_root, "minorFont"),
    }
    spacing = {
        "base_unit": DEFAULT_BASE_UNIT_PT,
        "margin": DEFAULT_MARGIN_PT,
        "gutter": DEFAULT_GUTTER_PT,
    }
    return {
        "colors": colors,
        "fonts": fonts,
        "spacing": spacing,
    }


def _theme_root(presentation: Presentation) -> ET.Element:
    for slide_master in presentation.slide_masters:
        for relationship in slide_master.part.rels.values():
            if relationship.reltype == THEME_RELATIONSHIP:
                return ET.fromstring(relationship.target_part.blob)
    raise AgentSlidesError(SCHEMA_ERROR, "template theme could not be read")


def _theme_color(theme_root: ET.Element, color_name: str) -> str:
    color_element = theme_root.find(f".//a:clrScheme/a:{color_name}", DRAWINGML_NS)
    if color_element is None:
        raise AgentSlidesError(
            SCHEMA_ERROR, f"template theme color '{color_name}' could not be read"
        )

    srgb = color_element.find("./a:srgbClr", DRAWINGML_NS)
    if srgb is not None:
        return f"#{srgb.attrib['val']}"

    system = color_element.find("./a:sysClr", DRAWINGML_NS)
    if system is not None:
        return f"#{system.attrib.get('lastClr', system.attrib['val'])}"

    raise AgentSlidesError(
        SCHEMA_ERROR, f"template theme color '{color_name}' could not be read"
    )


def _theme_font(theme_root: ET.Element, font_family: str) -> str:
    font = theme_root.find(f".//a:fontScheme/a:{font_family}/a:latin", DRAWINGML_NS)
    if font is None or "typeface" not in font.attrib:
        raise AgentSlidesError(
            SCHEMA_ERROR, f"template theme font '{font_family}' could not be read"
        )
    return font.attrib["typeface"]


def _relative_source_path(template_path: Path, manifest_path: Path) -> str:
    relative_path = os.path.relpath(
        template_path.resolve(), manifest_path.parent.resolve()
    )
    return Path(relative_path).as_posix()


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(64 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _write_manifest(path: Path, manifest: dict[str, object]) -> None:
    payload = f"{json.dumps(manifest, indent=2)}\n"
    try:
        path.write_text(payload, encoding="utf-8")
    except OSError as exc:
        raise AgentSlidesError(
            SCHEMA_ERROR,
            f"Failed to write manifest {path}: {exc.strerror or str(exc)}",
        ) from exc


def _emu_to_points(value: int) -> float:
    return round(float(value) / EMU_PER_POINT, 3)
